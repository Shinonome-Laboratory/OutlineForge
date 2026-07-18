"""FastAPI application for outline — content structure analyzer (Function B).

Serves on port 8001 with API prefix /api/outline/.
Shares corpus.db with the corpus module (Function A).
"""

import asyncio
import json
import logging
import os
import sqlite3
import statistics
import threading
import traceback
import uuid
from contextlib import asynccontextmanager
from dataclasses import dataclass, field
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Dict

import httpx
from fastapi import FastAPI, File, HTTPException, Request, UploadFile
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse, StreamingResponse

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

PROJECT_ROOT = Path(__file__).resolve().parent
# OUTLINE_DB_PATH 覆盖共享生产库路径（测试隔离用，与 01/03 的 env 兜底模式一致）
DB_PATH = Path(os.environ.get(
    "OUTLINE_DB_PATH", r"d:\Project\All for Style\00-data\corpus.db"
))
HTML_PATH = PROJECT_ROOT / "outline.html"
DB_VIEWER_PATH = PROJECT_ROOT / "db-viewer-02outline.html"

# ---------------------------------------------------------------------------
# Config defaults
# ---------------------------------------------------------------------------

_OB_PROMPT_ROUND1 = (
    "You are an expert curriculum designer and text structure analyst. "
    "Segment the numbered lecture transcript below into pedagogical topics.\n"
    "\n"
    "Your segmentation should behave like **condensed milk poured on a "
    "flat surface** — spreading naturally, maintaining consistent "
    "thickness, and staying highly cohesive without breaking apart.\n"
    "\n"
    "WORKFLOW:\n"
    "1. Count the total number of paragraphs in the transcript. This is "
    "your `total_paragraphs`.\n"
    "2. Scan the text. Identify the main pedagogical topics — each a "
    "conceptually complete unit. Write a concise title and a 1-sentence "
    "concept summary for each topic (this helps verify even density).\n"
    "3. For each topic (except the last), determine the paragraph index "
    "where it ENDS. The next topic begins at end+1.\n"
    "\n"
    "PRINCIPLES:\n"
    "• Consistent Density: Cognitive load of each topic must be roughly "
    "equal. Don't cram heavy concepts together; don't isolate trivial "
    "transitions as standalone topics.\n"
    "• Conceptual Completeness: Only cut after a concept is fully "
    "resolved (introduction → explanation → examples → wrap-up).\n"
    "• Natural Boundaries: Cut at semantic shifts or breathing pauses. "
    "If a paragraph bridges two ideas, assign it to the heavier side.\n"
    "\n"
    "CHAIN CALCULATION (Zero overlap, zero gap):\n"
    "• Topic 1 starts at paragraph 1.\n"
    "• Each subsequent topic starts at (previous topic's end + 1).\n"
    "• The final topic MUST end at total_paragraphs.\n"
    "• breaks array: the end paragraph index for each topic except the "
    "last. N topics → exactly N−1 breaks.\n"
    "\n"
    "OUTPUT — strict JSON only. No markdown fences, no extra text:\n"
    "{\n"
    '  "course_name": "Concise lecture title",\n'
    '  "total_paragraphs": <integer>,\n'
    '  "topics": ["Topic 1 Title", "Topic 2 Title", "Topic 3 Title"],\n'
    '  "breaks": [end_para_of_topic1, end_para_of_topic2]\n'
    "}\n"
    "\n"
    "FATAL RED LINES:\n"
    "• breaks[i] < breaks[i+1] (strictly increasing).\n"
    "• breaks[-1] < total_paragraphs (last break must leave room for final topic).\n"
    "• len(breaks) == len(topics) − 1.\n"
    "• All values are integers.\n"
    "\n"
    "Full transcript below:"
)

_OB_PROMPT_ROUND2 = (
    "You are a subject knowledge structure analyst. Below is the teaching "
    'content for the topic "{topic_name}" from a lecture transcript.\n'
    "\n"
    "Analyze the knowledge structure within this text and generate a knowledge "
    "subtree for this topic. Let the depth be determined by the actual content "
    "structure — shallow for simple content, deep for complex content.\n"
    "\n"
    "Output strictly in the following JSON format. Do not output any other text:\n"
    "\n"
    "{\n"
    '  "topic": "{topic_name}",\n'
    '  "subtree": {\n'
    '    "children": [\n'
    '      { "name": "sub-topic name", "children": [...] }\n'
    "    ]\n"
    "  }\n"
    "}\n"
    "\n"
    "If the topic content is too simple to subdivide further, children can be "
    "an empty array.\n"
    "\n"
    "Topic text below:"
)

OB_CONFIG_DEFAULTS = [
    ("ob_llm_model", "qwen2.5:14b-instruct"),
    ("ob_llm_temperature", "0.0"),
    # Ollama context-window size (tokens). Ollama defaults to ~2048 and SILENTLY
    # truncates longer prompts; this is user-tunable from the settings panel.
    ("ob_llm_num_ctx", "8192"),
    ("ob_prompt_round1", _OB_PROMPT_ROUND1),
    ("ob_prompt_round2", _OB_PROMPT_ROUND2),
]

# ---------------------------------------------------------------------------
# CKG concept extraction prompt (Step 1) — see ADR 0002 / CONTEXT.md「CK 维度」
# Stored in config under key ``ck_prompt_concepts`` (front-end editable).
# ---------------------------------------------------------------------------

_CK_PROMPT_CONCEPTS = (
    "You are a curriculum analyst reconstructing how an instructor DECOMPOSES a lecture's title into teaching concepts.\n"
    "\n"
    "Read the transcript and break the lecture's TITLE down into the teaching concepts the instructor actually carves out, as a DECOMPOSITION TREE: each concept is a piece the instructor splits a higher concept (or the lecture title) into. This captures HOW THIS INSTRUCTOR organizes the content — grounded entirely in THIS transcript, not the subject's objective structure.\n"
    "\n"
    "A TEACHING CONCEPT is a nameable, definable knowledge point the instructor actually teaches — something that could be a bullet on the syllabus, or that the instructor effectively says \"now let's learn X\" about.\n"
    "\n"
    "INCLUDE: named methods, models, principles, terms the instructor explains or builds on (e.g. \"Gradient Descent\", \"Learning Rate\", \"Cost Function\").\n"
    "EXCLUDE:\n"
    "- Tools/sub-steps used only to explain something else, not taught as their own topic (e.g. \"partial derivative\" mentioned only to derive gradient descent).\n"
    "- Generic words, examples, anecdotes, course logistics, greetings.\n"
    "- Anything NOT covered in this transcript. Do NOT add knowledge from outside the transcript — extract only what is actually taught here.\n"
    "\n"
    "GRANULARITY: syllabus-level — each concept is a nameable, definable knowledge point as described above. Extract AS MANY as the lecture genuinely teaches at this granularity; let the COUNT follow the actual content (a rich lecture yields many, a light one few). Do NOT pad, drop, split, or merge concepts just to hit any target number. Keep the granularity (what counts as a concept) consistent — never the count.\n"
    "\n"
    "For each concept output:\n"
    "- NAME: a short name.\n"
    "- DEFINITION: a ONE-SENTENCE definition grounded in how the transcript uses it.\n"
    "- FIRST_PARA: the integer index of the [段落 N] marker where the concept is first substantively taught.\n"
    "- PARENT: the NAME of the higher concept this one is carved out of (the concept the instructor was decomposing when this came up). If it is a top-level piece carved directly from the lecture title, set PARENT to the lecture TITLE string. PARENT must be EXACTLY another concept's NAME or the TITLE — never invent a new label here.\n"
    "\n"
    "Also output TITLE: the lecture's title / overall topic, used as the root of the decomposition tree.\n"
    "\n"
    "OUTPUT — strict JSON only. No markdown fences, no extra text:\n"
    "{\n"
    '  "title": "...",\n'
    '  "concepts": [ {"name": "...", "definition": "...", "first_para": <int>, "parent": "..."} ]\n'
    "}\n"
    "\n"
    "Full transcript below:"
)

# ---------------------------------------------------------------------------
# CKG prerequisite-edge extraction prompt (Step 2) — see ADR 0002 / CONTEXT.md
# Stored in config under key ``ck_prompt_edges`` (front-end editable).
# ---------------------------------------------------------------------------

_CK_PROMPT_EDGES = (
    "You are a learning-dependency analyst. Below is a list of teaching concepts (with definitions) from a single lecture.\n"
    "\n"
    "Identify PREREQUISITE relations. An edge A -> B means \"a learner must understand A before they can understand B.\" This is learning dependency / ordering — NOT containment or categorization.\n"
    "\n"
    "CRITICAL — what is and isn't a prerequisite edge:\n"
    "  ✓  \"Derivative\" -> \"Gradient Descent\"   (must grasp derivatives first)\n"
    "  ✗  \"Classification\" / \"Regression\" are both kinds of \"Supervised Learning\" → category relation, NOT prerequisite. No edge.\n"
    "  ✗  \"Classification\" vs \"Regression\" → siblings, either can be learned first. No edge.\n"
    "\n"
    "RULES:\n"
    "- Use ONLY the concepts in the list; use their exact names; invent nothing.\n"
    "- Edge only when the dependency is genuinely real. A few high-confidence edges beat many speculative ones — not every concept needs an edge.\n"
    "- Direction: prerequisite (earlier) -> dependent (later).\n"
    "- No cycles: if A -> B (directly or transitively), do not also add B -> A.\n"
    "\n"
    "OUTPUT — strict JSON only. No markdown fences, no extra text:\n"
    "{\n"
    '  "edges": [ {"from": "...", "to": "..."} ]\n'
    "}\n"
    "\n"
    "Concept list:"
)

# ---------------------------------------------------------------------------
# CKG grounded cross-relation prompt (Step 2, 讲述关联) — see ADR 0003.
# Stored in config under key ``ck_prompt_relations`` (front-end editable). The
# relation TYPE is FREE-LABELLED by the LLM (emergent, not a fixed taxonomy);
# builds-on / explains-via / motivates are only example anchors.
# ---------------------------------------------------------------------------

_CK_PROMPT_RELATIONS = (
    "You are analyzing how an instructor CONNECTS concepts ACROSS different branches of a lecture, beyond the parent→child decomposition.\n"
    "\n"
    "Below is the list of teaching concepts (with definitions) the instructor decomposed the lecture into, followed by the full transcript.\n"
    "\n"
    "Identify GROUNDED CROSS-RELATIONS: directed relations the instructor EXPLICITLY establishes between two concepts while teaching — connections the instructor actually states or demonstrates in THIS transcript, that are NOT already captured by the decomposition (parent→child) structure.\n"
    "\n"
    "For each relation output:\n"
    "- FROM / TO: the exact NAMES of two concepts from the list (the relation goes FROM → TO).\n"
    "- TYPE: a short label for HOW the instructor connects them — describe the relation in YOUR OWN WORDS (e.g. \"builds-on\", \"explains-via\", \"motivates\", \"contrasts-with\"). Use whatever label best fits; you are NOT limited to these examples.\n"
    "- EVIDENCE: a brief quote or the [段落 N] index grounding this relation in the transcript.\n"
    "\n"
    "RULES:\n"
    "- Use ONLY concept names from the list; use their exact names; invent no new concepts.\n"
    "- Only relations the instructor genuinely establishes in the transcript — a few well-grounded relations beat many speculative ones; not every concept needs one.\n"
    "- These are CROSS-branch narrative links, NOT the decomposition hierarchy (do not just restate parent→child).\n"
    "\n"
    "OUTPUT — strict JSON only. No markdown fences, no extra text:\n"
    "{\n"
    '  "relations": [ {"from": "...", "to": "...", "type": "...", "evidence": "..."} ]\n'
    "}\n"
    "\n"
    "Concept list:"
)

# ck_prompt_edges (prerequisite-edge prompt) is RETIRED for CKG extraction —
# the graph backbone is now the title→concept DECOMPOSITION tree (ADR 0003).
# _CK_PROMPT_EDGES / parse_edges_response are kept only for lesson generation.
CK_CONFIG_DEFAULTS = [
    ("ck_prompt_concepts", _CK_PROMPT_CONCEPTS),
    ("ck_prompt_relations", _CK_PROMPT_RELATIONS),
]

# ---------------------------------------------------------------------------
# Lesson generation (教案生成) prompts — see CONTEXT.md.
# topic-mode concept enumeration (no transcript: LLM uses its own knowledge);
# and the style-aware outline generation prompt (3rd LLM call).
# ---------------------------------------------------------------------------

_LESSON_GEN_TOPIC_PROMPT = (
    "You are a curriculum designer. Given a single teaching TOPIC, produce a "
    "DECOMPOSITION of the topic into teaching concepts, plus grounded relations.\n"
    "\n"
    "Use your own subject knowledge — there is no transcript. Produce roughly "
    "10-25 syllabus-level concepts (named methods, models, principles, terms), "
    "not micro-terms or anecdotes.\n"
    "\n"
    "DECOMPOSITION: organize concepts as a tree under the TOPIC. For each concept "
    "give a PARENT — the name of the higher concept it is carved out of, or the "
    "TOPIC string for a top-level concept. PARENT must be exactly another "
    "concept's NAME or the TOPIC.\n"
    "\n"
    "RELATIONS: directed cross-links between concepts beyond the decomposition "
    "(e.g. one concept builds on / explains-via / motivates another). Give each a "
    "free-form TYPE label describing the connection; you are NOT limited to a "
    "fixed set. Use exact concept names; a few well-founded relations beat many.\n"
    "\n"
    "OUTPUT — strict JSON only. No markdown fences, no extra text:\n"
    "{\n"
    '  "title": "<the topic>",\n'
    '  "concepts": [ {"name": "...", "definition": "...", "parent": "..."} ],\n'
    '  "relations": [ {"from": "...", "to": "...", "type": "..."} ]\n'
    "}\n"
    "\n"
    "TOPIC:"
)

_LESSON_GEN_OUTLINE_PROMPT = (
    "You are organizing teaching content in a specific teacher's style.\n"
    "\n"
    "TEACHER'S CONTENT-ORGANIZATION STYLE:\n"
    "{generation_rules}\n"
    "Target shape: depth≈{depth_mean}, branching≈{branch_mean}; "
    "sequencing: {sequencing}.\n"
    "\n"
    "CONCEPTS (with one-line definitions):\n"
    "{concept_lines}\n"
    "DECOMPOSITION (parent -> child = teacher carves child out of parent):\n"
    "{edge_lines}\n"
    "GROUNDED RELATIONS (A -[type]-> B = teacher connects A to B):\n"
    "{relation_lines}\n"
    "\n"
    "Organize these concepts into a teaching outline IN THIS TEACHER'S STYLE:\n"
    "- Nest to roughly the target depth/branching, following the decomposition.\n"
    "- Order following the sequencing tendency above.\n"
    "- Use the grounded relations to connect concepts across branches where it helps.\n"
    "\n"
    "OUTPUT strict JSON only, no fences:\n"
    "{{\n"
    '  "outline": [ {{"name":"...", "children":[ {{"name":"...","children":[]}} ]}} ],\n'
    '  "sequence": ["concept names in teaching order"]\n'
    "}}\n"
)

# Long-material protection for lesson generation: the material is split into
# chunks and concepts are extracted chunk-by-chunk, so no global truncation is
# needed. A single oversized chunk is still capped to fit the context window.
#   - "page"    chunk mode: one chunk per PDF page / PPTX slide (see _PAGE_DELIM)
#   - "segment" chunk mode: pages merged then re-split into ~N-char segments
_LESSON_GEN_MAX_CHUNK_CHARS = 12000
_LESSON_GEN_SEGMENT_CHARS = 6000

# Big-lecture → small-units split (e.g. a 2h PDF → Andrew-Ng-sized ~7–18min
# units). When the merged concept count exceeds the threshold, the UI offers to
# split into K units (~_CONCEPTS_PER_UNIT concepts each), then generates a
# styled lesson plan per unit.
_UNIT_SPLIT_THRESHOLD = 18
_CONCEPTS_PER_UNIT = 8

# ---------------------------------------------------------------------------
# Schema
# ---------------------------------------------------------------------------

COURSE_TOPICS_DDL = """
CREATE TABLE IF NOT EXISTS course_topics (
    id                INTEGER PRIMARY KEY AUTOINCREMENT,
    video_id          INTEGER NOT NULL,
    start_para_index  INTEGER NOT NULL,
    end_para_index    INTEGER NOT NULL,
    start_time        REAL    NOT NULL,
    end_time          REAL    NOT NULL,
    topic_name        TEXT    NOT NULL,
    subtree_json      TEXT,
    FOREIGN KEY (video_id) REFERENCES videos(id)
);
"""

# Course Knowledge Graph (CKG) — one row per video.
# graph_json (schema_version 2) holds {"schema_version", "title", "concepts"
# (each with parent/first_para), "decomposition_edges", "relations"}. The
# decomposition backbone is the title→concept tree (issue 01); grounded
# cross-relations (讲述关联) fill "relations" in issue 02.
COURSE_CKG_DDL = """
CREATE TABLE IF NOT EXISTS course_ckg (
    video_id    INTEGER PRIMARY KEY,
    graph_json  TEXT,
    model       TEXT,
    created_at  TEXT,
    FOREIGN KEY (video_id) REFERENCES videos(id)
);
"""

# Saved lesson plans (教案). payload_json holds the full generated result
# ({title, units:[plan,...]}) so it can be reloaded into the outline + graphs.
LESSON_PLANS_DDL = """
CREATE TABLE IF NOT EXISTS lesson_plans (
    id           INTEGER PRIMARY KEY AUTOINCREMENT,
    name         TEXT NOT NULL,
    created_at   TEXT,
    payload_json TEXT
);
"""

# ---------------------------------------------------------------------------
# Database helpers
# ---------------------------------------------------------------------------


def get_db(db_path: str) -> sqlite3.Connection:
    """Return a read-write connection with FK and WAL pragmas set."""
    conn = sqlite3.connect(db_path, timeout=10.0)
    conn.row_factory = sqlite3.Row
    conn.execute("PRAGMA foreign_keys = ON")
    conn.execute("PRAGMA journal_mode = WAL")
    return conn


def init_db(db_path: str) -> None:
    """Create course_topics table (if not exists) and seed ob_ config defaults."""
    conn = sqlite3.connect(db_path)
    conn.execute("PRAGMA foreign_keys = ON")
    conn.execute("PRAGMA journal_mode = WAL")
    conn.executescript(COURSE_TOPICS_DDL)
    conn.executescript(COURSE_CKG_DDL)
    conn.executescript(LESSON_PLANS_DDL)

    # Idempotent migration: CREATE TABLE IF NOT EXISTS never adds columns to an
    # already-existing course_ckg, so add the topology / delivery-direction
    # columns one by one only when missing (issue 03).
    _CKG_TOPO_COLUMNS = [
        ("depth", "INTEGER"),            # = 拆解深度 (decomposition-tree height)
        ("branch_factor", "REAL"),       # = 拆解宽度 (mean children of non-leaf)
        ("convergence_count", "INTEGER"),  # in-degree>1 over decomposition+relations
        ("relation_density", "REAL"),    # = 讲述关联边数 / 概念数 (issue 03)
        ("density", "REAL"),             # RETIRED (issue 03): no longer populated
        ("avg_path_length", "REAL"),     # RETIRED (issue 03): no longer populated
        ("clustering", "REAL"),          # RETIRED (issue 03): no longer populated
        ("bottomup_ratio", "REAL"),      # 辅助信号 (delivery direction)
    ]
    existing_cols = {
        row[1] for row in conn.execute("PRAGMA table_info(course_ckg)")
    }
    for col_name, col_type in _CKG_TOPO_COLUMNS:
        if col_name not in existing_cols:
            conn.execute(
                f"ALTER TABLE course_ckg ADD COLUMN {col_name} {col_type}"
            )

    conn.executemany(
        "INSERT OR IGNORE INTO config (key, value) VALUES (?, ?)",
        OB_CONFIG_DEFAULTS,
    )
    conn.executemany(
        "INSERT OR IGNORE INTO config (key, value) VALUES (?, ?)",
        CK_CONFIG_DEFAULTS,
    )

    # Migrate old round-1 prompts to the latest version
    _OLD_MARKERS = [
        '"start_para"',                           # v1: range-based prompt
        'classroom instruction content analyst',  # v2: breakpoint prompt
        '"last_para"',                            # v3: curriculum designer prompt
    ]
    row = conn.execute(
        "SELECT value FROM config WHERE key = 'ob_prompt_round1'"
    ).fetchone()
    if row and any(m in (row[0] or "") for m in _OLD_MARKERS):
        conn.execute(
            "UPDATE config SET value = ? WHERE key = 'ob_prompt_round1'",
            (_OB_PROMPT_ROUND1,),
        )
        conn.execute("COMMIT")

    # Migrate any pre-decomposition ck_prompt_concepts (concept-list only, no
    # PARENT field) to the latest DECOMPOSITION-tree prompt (ADR 0003). The new
    # prompt is detected by the presence of the "parent" output field.
    ck_row = conn.execute(
        "SELECT value FROM config WHERE key = 'ck_prompt_concepts'"
    ).fetchone()
    if ck_row and '"parent"' not in (ck_row[0] or ""):
        conn.execute(
            "UPDATE config SET value = ? WHERE key = 'ck_prompt_concepts'",
            (_CK_PROMPT_CONCEPTS,),
        )
        conn.execute("COMMIT")

    conn.commit()
    conn.close()


# ---------------------------------------------------------------------------
# SSE Logging infrastructure — Slice 02
# ---------------------------------------------------------------------------


@dataclass
class LogEvent:
    """A structured log event for SSE streaming.

    Attributes:
        category: SSE event type — one of info, debug, success, error, progress.
        message: Human-readable log message.
        ts: ISO-8601 UTC timestamp, auto-generated.
        traceback: Optional traceback string (populated for error events).
        progress_pct: Optional progress percentage (populated for progress events).
    """

    category: str
    message: str
    ts: str = field(default_factory=lambda: datetime.now(timezone.utc).isoformat())
    traceback: str | None = None
    progress_pct: float | None = None


_log_queue: asyncio.Queue | None = None
"""Module-level queue shared between log handler and SSE endpoint.

Retained as the *default* queue so the existing SSE tests (which set
``_log_queue`` directly, push, then drain it) keep passing. ``_broadcast``
always delivers to this queue in addition to every per-connection subscriber.
"""

_log_subscribers: set[asyncio.Queue] = set()
"""Per-connection SSE queues. Each ``stream_logs`` connection registers its own
queue here so events fan out to *all* live connections (broadcast), instead of
being consumed by whichever single connection happens to win the race on a
shared queue."""

_event_loop: asyncio.AbstractEventLoop | None = None
"""Event loop captured at lifespan startup for thread-safe queue access."""


def get_log_queue() -> asyncio.Queue:
    """Return the module-level log queue, creating it lazily if needed."""
    global _log_queue
    if _log_queue is None:
        _log_queue = asyncio.Queue()
    return _log_queue


def _get_loop() -> asyncio.AbstractEventLoop | None:
    """Return the captured event loop, the current running loop, or None."""
    global _event_loop
    if _event_loop is not None:
        return _event_loop
    try:
        return asyncio.get_running_loop()
    except RuntimeError:
        return None


class SSELogHandler(logging.Handler):
    """Custom logging handler that converts LogRecords into LogEvents
    and pushes them onto the asyncio Queue for SSE streaming.

    Thread-safety: emit() may be called from any thread.  We capture the
    event loop at construction time and use call_soon_threadsafe to push.
    """

    def __init__(self, loop: asyncio.AbstractEventLoop):
        super().__init__()
        self._loop = loop

    def emit(self, record: logging.LogRecord) -> None:
        """Convert a LogRecord to a LogEvent and push it to the queue."""
        try:
            category = _level_to_category(record.levelno)
            message = record.getMessage()
            tb = None
            if record.exc_info and record.exc_info[1] is not None:
                tb = traceback.format_exc()

            event = LogEvent(category=category, message=message, traceback=tb)
            _broadcast(event)
        except Exception:
            self.handleError(record)


def _level_to_category(levelno: int) -> str:
    """Map a Python log level number to an SSE event category string."""
    if levelno >= logging.ERROR:
        return "error"
    elif levelno >= logging.INFO:
        return "info"
    else:
        return "debug"


def _event_to_sse(event: LogEvent) -> str:
    """Convert a LogEvent to an SSE-formatted string.

    Format:
        event: <category>
        data: <json>
        <blank line>
    """
    data: dict = {
        "ts": event.ts,
        "msg": event.message,
        "category": event.category,
    }
    if event.traceback is not None:
        data["traceback"] = event.traceback
    if event.progress_pct is not None:
        data["progress"] = event.progress_pct

    return (
        f"event: {event.category}\n"
        f"data: {json.dumps(data, ensure_ascii=False)}\n\n"
    )


async def sse_event_generator(queue: asyncio.Queue):
    """Async generator that yields SSE-formatted log events from the queue.

    Runs indefinitely — the SSE endpoint will close when the client disconnects.
    """
    while True:
        event = await queue.get()
        yield _event_to_sse(event)


def push_log_event(
    category: str,
    message: str,
    traceback: str | None = None,
    progress_pct: float | None = None,
) -> None:
    """Push a structured event directly to the SSE log queue.

    ⚠ 签名冲突警告：03-actions 也有同名函数 ``push_log_event``，但签名不同
    （02 传 ``message: str``，03 传 dict payload）。若未来抽公共库，两者不可
    直接合并——必须先统一签名或改名。

    This bypasses the Python logging module entirely, allowing callers to
    emit custom categories (``success``, ``progress``) that don't map to
    standard log levels.

    Args:
        category: SSE category — info, debug, success, error, or progress.
        message: Human-readable log message.
        traceback: Optional traceback string for error events.
        progress_pct: Optional progress percentage (0.0–100.0).
    """
    event = LogEvent(
        category=category,
        message=message,
        traceback=traceback,
        progress_pct=progress_pct,
    )
    _broadcast(event)


def _broadcast(event: LogEvent) -> None:
    """Fan-out a LogEvent to every live SSE consumer.

    Delivers to the *default* queue (``get_log_queue()``) plus each
    per-connection queue registered in ``_log_subscribers``. Without this,
    multiple ``EventSource`` connections would compete for events on one shared
    ``asyncio.Queue`` (each event goes to exactly one consumer), so the second
    connection would silently miss events — the root cause of the previous
    "button stuck / fake progress" behaviour.

    The default queue is kept so the existing SSE tests (which push to and drain
    ``_log_queue`` directly, without ever opening ``stream_logs``) still pass.

    Thread-safe: ``_run_ckg_extraction`` runs in a worker thread, so we hop to
    the captured event loop via ``call_soon_threadsafe`` when available. Each
    queue is delivered independently so one full queue never blocks the rest.
    """
    loop = _get_loop()
    targets = [get_log_queue()]
    targets.extend(_log_subscribers)
    for q in targets:
        try:
            if loop is not None:
                loop.call_soon_threadsafe(q.put_nowait, event)
            else:
                q.put_nowait(event)
        except asyncio.QueueFull:
            import sys
            print(f"[SSE_DROP] QueueFull: [{event.category}] {event.message[:120]}", file=sys.stderr, flush=True)
        except Exception as exc:
            import sys
            print(f"[SSE_DROP] {type(exc).__name__}: [{event.category}] {event.message[:120]}", file=sys.stderr, flush=True)


# ---------------------------------------------------------------------------
# Lifespan
# ---------------------------------------------------------------------------


@asynccontextmanager
async def lifespan(application: FastAPI):
    """Startup: initialise database and configure SSE log handler."""
    init_db(str(DB_PATH))

    # Configure SSE log handler on the "outline" logger
    loop = asyncio.get_running_loop()
    global _event_loop
    _event_loop = loop
    handler = SSELogHandler(loop)
    handler.setLevel(logging.DEBUG)
    logger = logging.getLogger("outline")
    logger.addHandler(handler)
    logger.setLevel(logging.DEBUG)
    # Prevent duplicate messages in the root logger
    logger.propagate = False

    yield

    # Cleanup
    logger.removeHandler(handler)
    _event_loop = None


# ---------------------------------------------------------------------------
# FastAPI app
# ---------------------------------------------------------------------------

app = FastAPI(lifespan=lifespan)


# ---------------------------------------------------------------------------
# Routes
# ---------------------------------------------------------------------------


@app.get("/", response_class=HTMLResponse)
async def root():
    """Serve the outline single-page application shell."""
    return HTMLResponse(content=HTML_PATH.read_text(encoding="utf-8"))


@app.get("/db-viewer", response_class=HTMLResponse)
async def db_viewer():
    """Serve the database viewer for the outline module."""
    return HTMLResponse(content=DB_VIEWER_PATH.read_text(encoding="utf-8"))


# ---------------------------------------------------------------------------
# DB viewer API — shared with db-viewer-02outline.html
# ---------------------------------------------------------------------------


@app.get("/api/outline/db-path")
async def get_db_path():
    """Return the absolute filesystem path of the SQLite database."""
    return {"absolute_path": str(DB_PATH.resolve())}


@app.get("/api/outline/db")
async def serve_db_file():
    """Serve the raw SQLite database file for sql.js loading."""
    if not DB_PATH.exists():
        raise HTTPException(status_code=404, detail="Database file not found")
    return FileResponse(
        path=str(DB_PATH),
        media_type="application/vnd.sqlite3",
        filename="corpus.db",
    )


@app.post("/api/outline/db-table/_sql")
async def run_sql(request: Request):
    """Run any SQL statement. SELECT returns rows; INSERT/UPDATE/DELETE
    commits and returns rowcount; DDL commits and returns success."""
    body = await request.json()
    sql = (body.get("sql") or "").strip()
    if not sql:
        raise HTTPException(status_code=400, detail="Empty SQL")

    first_word = sql.split(maxsplit=1)[0].upper() if sql else ""
    is_read = first_word in ("SELECT", "WITH", "EXPLAIN", "PRAGMA", "DESCRIBE")

    conn = get_db(str(DB_PATH))
    conn.row_factory = None
    try:
        if is_read:
            cur = conn.execute(sql)
            cols = [c[0] for c in cur.description] if cur.description else []
            rows = cur.fetchall()
            conn.close()
            return {"type": "select", "columns": cols, "rows": rows, "total": len(rows)}
        else:
            cur = conn.execute(sql)
            conn.execute("COMMIT")
            rowcount = cur.rowcount if cur.rowcount >= 0 else 0
            conn.close()
            return {"type": "write", "affected": rowcount}
    except Exception as e:
        conn.close()
        raise HTTPException(status_code=400, detail=str(e))


# ---------------------------------------------------------------------------
# Config panel API — Slice 11
# ---------------------------------------------------------------------------


@app.get("/api/outline/config-defaults")
async def get_config_defaults():
    """配置默认值单一真相源（db-viewer 用）。

    内嵌在 db-viewer 里的副本曾经过期漂移（含已废弃的 ob_max_topics），01 的
    db-viewer 修过同类问题——统一改为前端 fetch 此端点。
    """
    return [[k, v] for k, v in (*OB_CONFIG_DEFAULTS, *CK_CONFIG_DEFAULTS)]


@app.get("/api/outline/config")
async def get_config():
    """Return all ob_ prefix config keys as a {key: value} JSON object."""
    conn = get_db(str(DB_PATH))
    rows = conn.execute(
        "SELECT key, value FROM config WHERE key LIKE 'ob_%'"
    ).fetchall()
    conn.close()
    return {row["key"]: row["value"] for row in rows}


@app.put("/api/outline/config")
async def put_config(payload: Dict[str, str]):
    """UPSERT config keys. Only keys with the 'ob_' prefix are allowed."""
    if not payload:
        raise HTTPException(status_code=400, detail="Request body must not be empty")

    # Validate: all keys must start with ob_
    for key in payload:
        if not key.startswith("ob_"):
            raise HTTPException(
                status_code=422,
                detail=f"Config key '{key}' does not have the required 'ob_' prefix",
            )

    conn = get_db(str(DB_PATH))
    try:
        for key, value in payload.items():
            conn.execute(
                "INSERT INTO config (key, value) VALUES (?, ?) "
                "ON CONFLICT(key) DO UPDATE SET value = excluded.value",
                (key, value),
            )
        conn.execute("COMMIT")

        # Return the updated values for these keys
        placeholders = ",".join("?" for _ in payload)
        rows = conn.execute(
            f"SELECT key, value FROM config WHERE key IN ({placeholders})",
            list(payload.keys()),
        ).fetchall()
        conn.close()
        return {row["key"]: row["value"] for row in rows}
    except Exception as e:
        conn.close()
        raise HTTPException(status_code=500, detail=str(e))


# ---------------------------------------------------------------------------
# CKG prompt editor API — issue 05
# The two CK prompt keys live in the same ``config`` table but are NOT exposed
# through the ob_-only config endpoints (those stay restricted to ob_ keys).
# A dedicated read/write pair keeps the CK prompts front-end editable without
# touching ob_ config behaviour.
# ---------------------------------------------------------------------------

_CK_PROMPT_KEYS = ("ck_prompt_concepts", "ck_prompt_relations")


@app.get("/api/outline/ckg/prompts")
async def get_ckg_prompts():
    """Return the editable CK prompt config values as a {key: value} JSON object."""
    conn = get_db(str(DB_PATH))
    rows = conn.execute(
        "SELECT key, value FROM config WHERE key IN (%s)"
        % ",".join("?" * len(_CK_PROMPT_KEYS)),
        _CK_PROMPT_KEYS,
    ).fetchall()
    conn.close()
    result = {row["key"]: row["value"] for row in rows}
    # Fall back to seeded defaults if a key is somehow missing.
    result.setdefault("ck_prompt_concepts", _CK_PROMPT_CONCEPTS)
    result.setdefault("ck_prompt_relations", _CK_PROMPT_RELATIONS)
    return result


@app.put("/api/outline/ckg/prompts")
async def put_ckg_prompts(payload: Dict[str, str]):
    """UPSERT CK prompt keys. Only ck_prompt_concepts is editable."""
    if not payload:
        raise HTTPException(status_code=400, detail="Request body must not be empty")

    for key in payload:
        if key not in _CK_PROMPT_KEYS:
            raise HTTPException(
                status_code=422,
                detail=f"Config key '{key}' is not an editable CK prompt key",
            )

    conn = get_db(str(DB_PATH))
    try:
        for key, value in payload.items():
            conn.execute(
                "INSERT INTO config (key, value) VALUES (?, ?) "
                "ON CONFLICT(key) DO UPDATE SET value = excluded.value",
                (key, value),
            )
        conn.execute("COMMIT")

        placeholders = ",".join("?" for _ in payload)
        rows = conn.execute(
            f"SELECT key, value FROM config WHERE key IN ({placeholders})",
            list(payload.keys()),
        ).fetchall()
        conn.close()
        return {row["key"]: row["value"] for row in rows}
    except Exception as e:
        conn.close()
        raise HTTPException(status_code=500, detail=str(e))


def _fetch_ollama_tags() -> list[dict]:
    """Call Ollama's GET /api/tags and return the raw list of model dicts.

    Extracted as a module-level function so tests can mock it without
    interfering with TestClient's own httpx usage.
    """
    with httpx.Client(timeout=5.0) as client:
        resp = client.get("http://localhost:11434/api/tags")
        resp.raise_for_status()
        return resp.json().get("models", [])


@app.get("/api/outline/ollama/models")
async def get_ollama_models():
    """Proxy to Ollama's /api/tags endpoint. Returns a flat list of model name
    strings. Returns empty list + 503 if Ollama is unreachable."""
    try:
        raw_models = _fetch_ollama_tags()
        models = [m["name"] for m in raw_models]
        return JSONResponse(content=models)
    except httpx.ConnectError:
        return JSONResponse(status_code=503, content=[])
    except Exception:
        return JSONResponse(status_code=503, content=[])


# ---------------------------------------------------------------------------
# Video list + paragraph browsing API — Slice 04
# ---------------------------------------------------------------------------


@app.get("/api/outline/videos")
async def get_videos():
    """Return all videos with status='done', including paragraph count.

    每个视频带归属字段 course_id / course_name / teacher_id / teacher_name
    （未分类视频全为 NULL）。此处 courses 指教师开设的一门课（teacher_course，
    01 建模），与 02 自己的 course_topics（单节课内容）无关，勿混淆。
    """
    conn = get_db(str(DB_PATH))
    rows = conn.execute(
        "SELECT v.id, v.name, v.duration, v.course_id, "
        "       c.name AS course_name, c.teacher_id, t.name AS teacher_name, "
        "       COALESCE(p.cnt, 0) AS paragraph_count, "
        "       COALESCE(tc.cnt, 0) AS topic_count, "
        "       COALESCE(tc.done, 0) AS subtree_count "
        "FROM videos v "
        "LEFT JOIN courses  c ON c.id = v.course_id "
        "LEFT JOIN teachers t ON t.id = c.teacher_id "
        "LEFT JOIN ("
        "    SELECT video_id, COUNT(*) AS cnt "
        "    FROM corpus_paragraphs "
        "    GROUP BY video_id"
        ") p ON v.id = p.video_id "
        "LEFT JOIN ("
        "    SELECT video_id, COUNT(*) AS cnt, COUNT(subtree_json) AS done "
        "    FROM course_topics "
        "    GROUP BY video_id"
        ") tc ON v.id = tc.video_id "
        "WHERE v.status = 'done' "
        "ORDER BY v.id"
    ).fetchall()
    conn.close()
    return [dict(row) for row in rows]


# ---------------------------------------------------------------------------
# Teachers — 教师档案（只读，P2）
# ---------------------------------------------------------------------------
# teachers / courses 由 01-corpus 独家建模与写入，02 只读消费（HANDOFF01-teachers
# §2）；编辑入口在 01（http://localhost:8000）。此处 courses = 教师开设的一门课
# （teacher_course），与 02 自己的 course_topics / course_ckg 里的 "course"
# （= 单个视频的一节课）没有任何外键关系——后者只能经 videos 中转 JOIN 到
# courses，绝不可直连（HANDOFF01-teachers §6-1）。

_OB_TEACHER_SELECT = """
    SELECT t.id, t.name, t.title, t.affiliation, t.email, t.note, t.created_at,
           (SELECT COUNT(*) FROM courses c
             WHERE c.teacher_id = t.id) AS course_count,
           (SELECT COUNT(*) FROM videos v
              JOIN courses c ON c.id = v.course_id
             WHERE c.teacher_id = t.id) AS video_count,
           (SELECT COUNT(*) FROM corpus_paragraphs cp
              JOIN videos v  ON v.id = cp.video_id
              JOIN courses c ON c.id = v.course_id
             WHERE c.teacher_id = t.id) AS para_count,
           (SELECT COUNT(*) FROM course_ckg k
              JOIN videos v  ON v.id = k.video_id
              JOIN courses c ON c.id = v.course_id
             WHERE c.teacher_id = t.id) AS extracted_count,
           (SELECT COALESCE(SUM(v.duration), 0) FROM videos v
              JOIN courses c ON c.id = v.course_id
             WHERE c.teacher_id = t.id) AS total_duration
      FROM teachers t
"""


@app.get("/api/outline/teachers")
async def get_teachers():
    """Return all teachers with aggregate stats (read-only).

    02 特有列 ``extracted_count`` = 该教师已抽取 CKG 的课数；其余聚合列与 01
    的 ``_TEACHER_SELECT`` 同构。
    """
    conn = get_db(str(DB_PATH))
    try:
        rows = conn.execute(_OB_TEACHER_SELECT + " ORDER BY t.name").fetchall()
    finally:
        conn.close()
    return [dict(r) for r in rows]


@app.get("/api/outline/teachers/{teacher_id}/ck-profile")
async def get_teacher_ck_profile(teacher_id: int):
    """某教师的 CK 分析档案：基本信息 + 每门课的 CK 统计 + CKG 指标均值。

    每门课的均值用标量子查询逐列聚合（01 的 ``_COURSE_SELECT`` 模式）——不能
    把 course_topics 与 course_ckg 同时 LEFT JOIN 再 AVG，行叉积会让话题多的
    视频在均值里被重复加权。RETIRED 列（density / avg_path_length /
    clustering）不消费。
    """
    conn = get_db(str(DB_PATH))
    try:
        teacher = conn.execute(
            _OB_TEACHER_SELECT + " WHERE t.id = ?", (teacher_id,)
        ).fetchone()
        if teacher is None:
            raise HTTPException(status_code=404, detail="Teacher not found")

        _ckg_scalar = (
            "(SELECT {expr} FROM course_ckg k "
            "  JOIN videos v ON v.id = k.video_id "
            " WHERE v.course_id = c.id)"
        )
        course_sql = (
            "SELECT c.id, c.name, c.semester, c.note, c.created_at, "
            "  (SELECT COUNT(*) FROM videos v WHERE v.course_id = c.id) "
            "    AS video_count, "
            "  (SELECT COUNT(*) FROM course_topics ct "
            "     JOIN videos v ON v.id = ct.video_id "
            "    WHERE v.course_id = c.id) AS topic_count, "
            + _ckg_scalar.format(expr="COUNT(*)") + " AS extracted_count, "
            + _ckg_scalar.format(expr="AVG(k.depth)") + " AS avg_depth, "
            + _ckg_scalar.format(expr="AVG(k.branch_factor)")
            + " AS avg_branch_factor, "
            + _ckg_scalar.format(expr="AVG(k.convergence_count)")
            + " AS avg_convergence_count, "
            + _ckg_scalar.format(expr="AVG(k.relation_density)")
            + " AS avg_relation_density, "
            + _ckg_scalar.format(expr="AVG(k.bottomup_ratio)")
            + " AS avg_bottomup_ratio "
            "FROM courses c WHERE c.teacher_id = ? ORDER BY c.name"
        )
        courses = conn.execute(course_sql, (teacher_id,)).fetchall()

        means = conn.execute(
            "SELECT AVG(k.depth)             AS depth, "
            "       AVG(k.branch_factor)     AS branch_factor, "
            "       AVG(k.convergence_count) AS convergence_count, "
            "       AVG(k.relation_density)  AS relation_density, "
            "       AVG(k.bottomup_ratio)    AS bottomup_ratio, "
            "       COUNT(*)                 AS n_videos "
            "FROM course_ckg k "
            "JOIN videos  v ON v.id = k.video_id "
            "JOIN courses c ON c.id = v.course_id "
            "WHERE c.teacher_id = ?",
            (teacher_id,),
        ).fetchone()
    finally:
        conn.close()

    return {
        "teacher": dict(teacher),
        "courses": [dict(r) for r in courses],
        "ckg_means": dict(means) if means else None,
    }


@app.get("/api/outline/video/{video_id}/paragraphs")
async def get_video_paragraphs(video_id: int):
    """Return all paragraphs for a video, with text truncated to 100 chars.

    Returns 404 if the video does not exist or has no paragraphs.
    """
    conn = get_db(str(DB_PATH))

    # Check video exists
    video = conn.execute(
        "SELECT id FROM videos WHERE id = ?", (video_id,)
    ).fetchone()
    if video is None:
        conn.close()
        raise HTTPException(status_code=404, detail="Video not found")

    rows = conn.execute(
        "SELECT paragraph_index, start_time, end_time, "
        "       SUBSTR(text, 1, 100) AS text "
        "FROM corpus_paragraphs "
        "WHERE video_id = ? "
        "ORDER BY paragraph_index",
        (video_id,),
    ).fetchall()
    conn.close()

    if not rows:
        raise HTTPException(status_code=404, detail="No paragraphs found for this video")

    return [dict(row) for row in rows]


@app.get("/api/outline/video/{video_id}/file")
async def serve_video_file(video_id: int):
    """Serve the video file for a given video_id.

    Reads the file path from the ``videos`` table and streams the file.
    Returns 404 if the video or its file does not exist.
    """
    conn = get_db(str(DB_PATH))
    row = conn.execute(
        "SELECT path FROM videos WHERE id = ?", (video_id,)
    ).fetchone()
    conn.close()

    if row is None:
        raise HTTPException(status_code=404, detail="Video not found")

    file_path = Path(row["path"])
    if not file_path.exists():
        raise HTTPException(status_code=404, detail="Video file not found on disk")

    # Guess media type from extension
    ext = file_path.suffix.lower()
    media_map = {
        ".mp4": "video/mp4",
        ".mkv": "video/x-matroska",
        ".webm": "video/webm",
        ".avi": "video/x-msvideo",
        ".mov": "video/quicktime",
    }
    return FileResponse(
        path=str(file_path),
        media_type=media_map.get(ext, "video/mp4"),
        filename=file_path.name,
    )


@app.get("/api/outline/topics/{video_id}")
async def get_topics(video_id: int):
    """Return topic segmentation results for a video from course_topics table.

    Each topic includes start/end times, paragraph range, duration, and
    topic name.  Returns an empty list if no topics exist for the video or
    the video itself does not exist.
    """
    conn = get_db(str(DB_PATH))
    try:
        rows = conn.execute(
            "SELECT id, video_id, start_para_index, end_para_index, "
            "start_time, end_time, topic_name, subtree_json "
            "FROM course_topics "
            "WHERE video_id = ? "
            "ORDER BY start_para_index",
            (video_id,),
        ).fetchall()

        topics = []
        for row in rows:
            topics.append({
                "id": row["id"],
                "video_id": row["video_id"],
                "start_para_index": row["start_para_index"],
                "end_para_index": row["end_para_index"],
                "start_time": row["start_time"],
                "end_time": row["end_time"],
                "duration": round(row["end_time"] - row["start_time"], 2),
                "topic_name": row["topic_name"],
                "subtree_json": row["subtree_json"],
            })
        return topics
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# PUT /api/outline/topics/{topic_id} — Slice 10 思维导图编辑回写
# ---------------------------------------------------------------------------


@app.put("/api/outline/topics/{topic_id}")
async def update_topic(topic_id: int, payload: dict):
    """Update a single course_topics row.

    Accepts a JSON body with optional ``topic_name`` and ``subtree_json``
    fields.  At least one field must be present.  Returns the updated row
    as JSON.

    Returns 404 if the topic_id does not exist.
    """
    if not payload:
        raise HTTPException(status_code=400, detail="Request body must not be empty")

    # Validate field names
    allowed = {"topic_name", "subtree_json"}
    for key in payload:
        if key not in allowed:
            raise HTTPException(
                status_code=422,
                detail=f"Unexpected field: '{key}'. Allowed fields: topic_name, subtree_json",
            )

    if not any(k in payload for k in ("topic_name", "subtree_json")):
        raise HTTPException(status_code=400, detail="At least one of topic_name or subtree_json is required")

    conn = get_db(str(DB_PATH))
    try:
        # Check topic exists
        existing = conn.execute(
            "SELECT id FROM course_topics WHERE id = ?", (topic_id,)
        ).fetchone()
        if existing is None:
            raise HTTPException(status_code=404, detail="Topic not found")

        # Build dynamic UPDATE
        set_clauses: list[str] = []
        values: list = []
        if "topic_name" in payload:
            set_clauses.append("topic_name = ?")
            values.append(payload["topic_name"])
        if "subtree_json" in payload:
            set_clauses.append("subtree_json = ?")
            values.append(payload["subtree_json"])

        values.append(topic_id)
        sql = f"UPDATE course_topics SET {', '.join(set_clauses)} WHERE id = ?"
        conn.execute(sql, values)
        conn.execute("COMMIT")

        # Return updated row
        row = conn.execute(
            "SELECT id, video_id, start_para_index, end_para_index, "
            "start_time, end_time, topic_name, subtree_json "
            "FROM course_topics WHERE id = ?",
            (topic_id,),
        ).fetchone()

        return dict(row)
    except HTTPException:
        raise
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# SSE streaming & error trigger — Slice 02
# ---------------------------------------------------------------------------


@app.get("/api/stream/logs/outline")
async def stream_logs():
    """SSE endpoint that streams log events from the outline processing pipeline.

    Returns ``text/event-stream`` with events categorised as info, debug,
    success, error, or progress.  Each event carries a JSON payload.

    The stream runs until the client disconnects.

    Each connection gets its *own* queue, registered in ``_log_subscribers`` so
    ``_broadcast`` fans every event out to all live connections. On disconnect
    the queue is removed in the generator's ``finally`` block.
    """
    q: asyncio.Queue = asyncio.Queue()
    _log_subscribers.add(q)

    async def _per_connection_generator():
        try:
            while True:
                yield _event_to_sse(await q.get())
        finally:
            _log_subscribers.discard(q)

    return StreamingResponse(
        _per_connection_generator(),
        media_type="text/event-stream",
        headers={
            "Cache-Control": "no-cache",
            "Connection": "keep-alive",
            "X-Accel-Buffering": "no",
        },
    )


@app.post("/api/outline/trigger-error")
async def trigger_error():
    """Deliberately trigger a ZeroDivisionError to test the SSE error display.

    The exception is caught, logged via the ``outline`` logger (which pushes
    an error event with full traceback to the SSE stream), and a 200 response
    is returned to confirm the trigger succeeded.
    """
    logger = logging.getLogger("outline")
    try:
        1 / 0
    except ZeroDivisionError:
        logger.error("除零错误测试 — deliberately triggered", exc_info=True)
    return {"status": "error_triggered"}


# ---------------------------------------------------------------------------
# Round 1 analysis — shared state (Slice 06)
# ---------------------------------------------------------------------------

_topic_segmentation_results: Dict[int, dict] = {}
"""Module-level dict storing parsed round-1 results keyed by video_id.
Written by the worker thread; consumed by Slice 07 for DB persistence."""

_analysis_cancel_flags: Dict[int, bool] = {}
"""Module-level dict tracking cancel requests keyed by video_id.
True = cancel requested; False = running (flag exists but not cancelled);
absent = no analysis running.  The worker thread checks this between steps."""


# ---------------------------------------------------------------------------
# JSON parsing helper (Slice 06)
# ---------------------------------------------------------------------------


def parse_round1_response(raw: str, total_para_count: int = 0) -> dict | None:
    """Parse the LLM's raw output into the round-1 result structure.

    The LLM uses ``last_para`` per topic (breakpoints).  This function converts
    those into non-overlapping ``start_para`` / ``end_para`` ranges.

    Handles common LLM output quirks:
    - Markdown code fences (```json ... ```)
    - Leading/trailing non-JSON text
    - Whitespace / BOM

    Returns a dict with keys ``course_name`` and ``topics`` on success,
    or ``None`` if the response cannot be parsed.
    """
    import re

    text = raw.strip()
    if not text:
        return None

    # 1) Try to extract JSON from markdown code fences
    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    # 2) Try to find the first '{' ... last '}' pair
    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace : last_brace + 1]

    # 3) Parse JSON
    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        return None

    # 4) Validate structure
    if not isinstance(data, dict):
        return None
    if "course_name" not in data:
        return None
    if "topics" not in data:
        return None
    if not isinstance(data["topics"], list):
        return None
    if len(data["topics"]) == 0:
        return None
    # Accept both new format (list of strings) and old format (list of dicts)
    if not all(isinstance(t, (str, dict)) for t in data["topics"]):
        return None

    # 5) Convert topic names + breaks → non-overlapping {start_para, end_para}
    if all(isinstance(t, str) for t in data["topics"]):
        # New format: { "topics": ["A","B","C"], "breaks": [5, 12] }
        topic_names = data["topics"]
        raw_breaks = data.get("breaks", [])

        # Use LLM-declared paragraph count if available, else fall back to caller
        declared_total = data.get("total_paragraphs")
        effective_total = total_para_count
        if isinstance(declared_total, (int, float)) and int(declared_total) > 0:
            effective_total = int(declared_total)

        if not isinstance(raw_breaks, list):
            raw_breaks = []

        # Validate breaks: integers, strictly increasing, within bounds
        breaks = []
        for b in raw_breaks:
            try:
                val = int(b)
            except (ValueError, TypeError):
                continue
            if val <= 0:
                continue
            if effective_total > 0 and val >= effective_total:
                val = effective_total  # cap at last paragraph
            if breaks and val <= breaks[-1]:
                continue  # skip non-increasing duplicates
            breaks.append(val)

        # Fix wrong number of breaks
        expected = len(topic_names) - 1
        if len(breaks) > expected:
            breaks = breaks[:expected]
        elif len(breaks) < expected and effective_total > 0:
            remaining = topic_names[len(breaks):]
            last_para = breaks[-1] if breaks else 0
            total_remaining = effective_total - last_para
            if total_remaining > 0:
                step = max(2, total_remaining // (len(remaining) + 1))
                for _ in range(len(remaining)):
                    last_para += step
                    if last_para >= effective_total:
                        last_para = effective_total
                    breaks.append(last_para)

        import sys
        print(f"[PARSE1] topics={len(topic_names)} breaks={breaks} total_para={effective_total}", file=sys.stderr, flush=True)

        result_topics = []
        last_seen = 0
        for i, name in enumerate(topic_names):
            start = last_seen + 1
            if i < len(breaks):
                end = int(breaks[i])
            else:
                end = effective_total if effective_total > 0 else total_para_count if total_para_count > 0 else 999999
            if end < start:
                # LLM assigned too few paragraphs → skip this topic
                print(f"[PARSE1] skip topic '{name}': start={start} > end={end}", file=sys.stderr, flush=True)
                continue
            result_topics.append({"name": name, "start_para": start, "end_para": end})
            last_seen = end
        data["topics"] = result_topics
        data.pop("breaks", None)
    elif any("last_para" in t for t in data.get("topics", [])):
        # Old format: per-topic "last_para"
        topics = data["topics"]
        last_seen_para = 0
        for i, t in enumerate(topics):
            raw_last = t.get("last_para")
            if raw_last is not None:
                end = int(raw_last)
            else:
                end = total_para_count if total_para_count > 0 else 999999
            start = last_seen_para + 1
            if end < start:
                end = start
            t["start_para"] = start
            t["end_para"] = end
            t.pop("last_para", None)
            last_seen_para = end
    else:
        # Legacy: start_para/end_para
        for t in data.get("topics", []):
            t["start_para"] = int(t.get("start_para", 1))
            t["end_para"] = int(t.get("end_para", t["start_para"]))
        data["topics"].sort(key=lambda t: t["start_para"])

    return data


# ---------------------------------------------------------------------------
# Round 2 JSON parsing (Slice 07)
# ---------------------------------------------------------------------------


def parse_round2_response(raw: str, topic_name: str = "") -> dict | None:
    """Parse the LLM's round-2 output into the subtree structure.

    Handles common LLM output quirks:
    - Markdown code fences (```json ... ```)
    - Leading/trailing non-JSON text
    - Whitespace / BOM
    - Nameless / empty leaf nodes (stripped)

    Returns a dict with keys ``topic`` and ``subtree`` on success,
    or ``None`` if the response cannot be parsed.
    """
    import re

    text = raw.strip()
    if not text:
        return None

    # 1) Try to extract JSON from markdown code fences
    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    # 2) Try to find the first '{' ... last '}' pair
    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace : last_brace + 1]

    # 3) Parse JSON — with auto-repair for truncated responses
    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        # Auto-repair: try adding missing closing brackets
        repaired = None
        for attempt in range(3):
            # Count open vs close braces/brackets
            opens = text.count("{") - text.count("}")
            closes = text.count("[") - text.count("]")
            if opens <= 0 and closes <= 0:
                break  # balanced — can't fix
            suffix = "}" * max(0, opens) + "]" * max(0, closes)
            try:
                data = json.loads(text + suffix)
                repaired = True
                break
            except (json.JSONDecodeError, ValueError):
                # Try removing trailing comma before adding brackets
                stripped = text.rstrip()
                if stripped.endswith(","):
                    text = stripped[:-1]
                    continue
                text += suffix
                try:
                    data = json.loads(text)
                    repaired = True
                    break
                except (json.JSONDecodeError, ValueError):
                    break
        if repaired:
            import sys
            print(f"[PARSE2] auto-repaired truncated JSON", file=sys.stderr, flush=True)
        else:
            import sys
            print(f"[PARSE2] unrepairable JSON len={len(text)}", file=sys.stderr, flush=True)
            return None
    except ValueError as e:
        import sys
        print(f"[PARSE2] ValueError: {e}", file=sys.stderr, flush=True)
        return None

    # 4) Validate structure
    if not isinstance(data, dict):
        return None
    if "subtree" not in data:
        return None

    return data


# ---------------------------------------------------------------------------
# Paragraph concatenation (Slice 06)
# ---------------------------------------------------------------------------


def concat_paragraphs(paragraphs: list[dict]) -> tuple[str, dict]:
    """Concatenate a list of paragraph dicts into a single full-text string
    with paragraph markers, suitable for feeding to the LLM.

    Each paragraph dict must have ``paragraph_index`` (int) and ``text`` (str).

    Returns:
        (full_text, para_map) — para_map is a dict mapping paragraph_index
        to the paragraph dict (for deriving time anchors later).
    """
    lines: list[str] = []
    para_map: dict = {}

    for p in paragraphs:
        idx = p["paragraph_index"]
        text = p.get("text", "")
        lines.append(f"[段落 {idx}] {text}")
        para_map[idx] = p

    return "\n\n".join(lines), para_map


# ---------------------------------------------------------------------------
# DB fetch helper (extracted for mockability in tests — Slice 06)
# ---------------------------------------------------------------------------


def _fetch_paragraphs_full_text(video_id: int) -> list[dict]:
    """Fetch ALL paragraphs for a video with their full text from corpus.db.

    Returns a list of dicts with keys: paragraph_index, start_time, end_time, text.
    Each paragraph.text is the FULL text (not truncated to 100 chars).
    """
    conn = get_db(str(DB_PATH))
    rows = conn.execute(
        "SELECT paragraph_index, start_time, end_time, text "
        "FROM corpus_paragraphs "
        "WHERE video_id = ? "
        "ORDER BY paragraph_index",
        (video_id,),
    ).fetchall()
    conn.close()
    return [dict(row) for row in rows]


# ---------------------------------------------------------------------------
# Ollama generate helper (extracted for mockability in tests — Slice 06)
# ---------------------------------------------------------------------------


def _call_ollama_generate(model: str, prompt: str, stream: bool = False, temperature: float = 0.0,
                          fmt: str | None = None, num_predict: int = 2048,
                          num_ctx: int | None = None) -> dict:
    """Call Ollama's POST /api/generate with the given model and prompt.

    Args:
        model: Ollama model name.
        prompt: The full prompt text.
        stream: If True, returns a streaming response.
        temperature: Sampling temperature (0.0 = deterministic).
        fmt: When truthy, sent as the top-level ``format`` field. Pass
            ``"json"`` to force Ollama to emit syntactically valid JSON,
            preventing the corrupted-array failures qwen produces on free-form
            generation.
        num_predict: Max tokens to predict (``options.num_predict``). Long
            definitions + many concepts need more than the 2048 default, so the
            CKG worker raises this to 4096.
        num_ctx: Context-window size (``options.num_ctx``). Ollama defaults to
            a small window (~2048) and SILENTLY truncates longer prompts, so
            callers feeding long material (uploaded PDFs) pass a larger value to
            actually fit the text. ``None`` leaves Ollama's default.

    Returns:
        The parsed JSON response body from Ollama.
    """
    import sys, time as _time
    t0 = _time.time()
    try:
        print(f"[OLLAMA] REQ model={model} prompt_len={len(prompt)} temp={temperature} fmt={fmt} num_predict={num_predict} num_ctx={num_ctx}", file=sys.stderr, flush=True)
        with httpx.Client(timeout=120.0) as client:
            _options = {
                "num_predict": num_predict,
                "temperature": temperature,
            }
            if num_ctx is not None:
                _options["num_ctx"] = num_ctx
            body = {
                "model": model,
                "prompt": prompt,
                "stream": stream,
                "options": _options,
            }
            if fmt:
                body["format"] = fmt
            resp = client.post(
                "http://localhost:11434/api/generate",
                json=body,
            )
            elapsed = _time.time() - t0
            print(f"[OLLAMA] RES status={resp.status_code} elapsed={elapsed:.1f}s", file=sys.stderr, flush=True)
            resp.raise_for_status()
            data = resp.json()
            # Keep stderr summary short
            response_text = data.get("response", "")
            done = data.get("done", False)
            ctx_len = data.get("context", [])
            print(f"[OLLAMA] OK done={done} resp_len={len(response_text)} ctx_tokens={len(ctx_len)} elapsed={elapsed:.1f}s", file=sys.stderr, flush=True)
            return data
    except httpx.ConnectError as e:
        print(f"[OLLAMA] CONNECT_ERROR: {e}", file=sys.stderr, flush=True)
        raise
    except httpx.ReadTimeout as e:
        print(f"[OLLAMA] READ_TIMEOUT after {_time.time()-t0:.1f}s: {e}", file=sys.stderr, flush=True)
        raise
    except httpx.HTTPStatusError as e:
        body = ""
        try: body = e.response.text[:500]
        except: pass
        print(f"[OLLAMA] HTTP_ERROR {e.response.status_code}: {body}", file=sys.stderr, flush=True)
        raise
    except Exception as e:
        print(f"[OLLAMA] ERROR {type(e).__name__}: {e}", file=sys.stderr, flush=True)
        raise


# ---------------------------------------------------------------------------
# Worker thread for round-1 analysis (Slice 06)
# ---------------------------------------------------------------------------


def _run_round1_analysis(video_id: int) -> None:
    """Background worker that performs round-1 LLM topic segmentation.

    Steps:
    1. Fetch all paragraphs for the video (full text).
    2. Concatenate paragraphs into one full transcript with markers.
    3. Load the round-1 prompt from config and append the transcript.
    4. Call Ollama /api/generate with the assembled prompt.
    5. Parse the JSON response.
    6. Store the parsed result in _topic_segmentation_results[video_id].
    7. Push SSE events (info, progress, success, error) throughout.
    """
    logger = logging.getLogger("outline")

    try:
        # --- Check cancel flag before each major step ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event("info", f"分析已中止 (video_id={video_id})")
            logger.info(f"Round-1 analysis cancelled for video_id={video_id}")
            return

        # --- Step 1: Fetch paragraphs ---
        push_log_event("info", "正在加载段落全文…", progress_pct=5.0)
        paragraphs = _fetch_paragraphs_full_text(video_id)

        if not paragraphs:
            push_log_event("error", f"视频 {video_id} 没有段落数据")
            return

        para_count = len(paragraphs)
        push_log_event(
            "info",
            f"已加载 {para_count} 个段落，正在拼接全文…",
            progress_pct=10.0,
        )

        # --- Check cancel ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event("info", f"分析已中止 (video_id={video_id})")
            logger.info(f"Round-1 analysis cancelled for video_id={video_id}")
            return

        # --- Step 2: Concatenate ---
        full_text, _para_map = concat_paragraphs(paragraphs)
        total_chars = len(full_text)
        push_log_event(
            "info",
            f"全文拼接完成，共 {total_chars} 字符，准备调用 LLM…",
            progress_pct=15.0,
        )

        # --- Check cancel ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event("info", f"分析已中止 (video_id={video_id})")
            logger.info(f"Round-1 analysis cancelled for video_id={video_id}")
            return

        # --- Step 3: Load config and build prompt ---
        conn = get_db(str(DB_PATH))
        config_rows = conn.execute(
            "SELECT key, value FROM config WHERE key IN (?, ?, ?, ?)",
            ("ob_prompt_round1", "ob_llm_model", "ob_llm_temperature",
             "ob_llm_num_ctx"),
        ).fetchall()
        conn.close()

        config = {row["key"]: row["value"] for row in config_rows}
        prompt_template = config.get("ob_prompt_round1", _OB_PROMPT_ROUND1)
        model = config.get("ob_llm_model", "qwen2.5:14b-instruct")
        temperature = float(config.get("ob_llm_temperature", 0.0))
        num_ctx = _parse_num_ctx(config.get("ob_llm_num_ctx"))

        full_prompt = prompt_template + "\n\n" + full_text

        # --- Step 4: Call Ollama ---
        push_log_event("info", "调用 LLM 进行话题切分…", progress_pct=20.0)
        logger.info(
            f"Calling Ollama model={model} prompt_length={len(full_prompt)}"
        )

        # --- Check cancel ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event("info", f"分析已中止 (video_id={video_id})")
            logger.info(f"Round-1 analysis cancelled for video_id={video_id}")
            return

        try:
            ollama_response = _call_ollama_generate(model, full_prompt, stream=False, temperature=temperature, num_ctx=num_ctx)
        except Exception as e:
            push_log_event(
                "error",
                f"Ollama 调用失败: {e}",
                traceback=traceback.format_exc(),
            )
            logger.error(f"Ollama call failed: {e}", exc_info=True)
            return

        # --- Check cancel ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event("info", f"分析已中止 (video_id={video_id})")
            logger.info(f"Round-1 analysis cancelled for video_id={video_id}")
            return

        # --- Step 5: Parse JSON ---
        push_log_event("info", "正在解析 LLM 返回结果…", progress_pct=80.0)
        raw_output = ollama_response.get("response", "")

        parsed = parse_round1_response(raw_output, para_count)
        if parsed is None:
            # JSON parse failed — push error with first 500 chars of raw output
            preview = raw_output[:500]
            push_log_event(
                "error",
                f"LLM 返回的 JSON 解析失败。原始输出 (前500字):\n{preview}",
            )
            logger.error(
                f"JSON parse failed for video_id={video_id}. "
                f"Raw output (first 500 chars): {preview}"
            )
            return

        # --- Step 6: Store result ---
        _topic_segmentation_results[video_id] = parsed

        topic_count = len(parsed.get("topics", []))
        course_name = parsed.get("course_name", "未知课程")
        push_log_event(
            "success",
            f"话题切分完成: 课程「{course_name}」, 共 {topic_count} 个话题",
            progress_pct=100.0,
        )
        logger.info(
            f"Round-1 analysis complete for video_id={video_id}: "
            f"course_name={course_name}, topics={topic_count}"
        )

    except Exception:
        push_log_event(
            "error",
            f"分析过程发生异常 (video_id={video_id})",
            traceback=traceback.format_exc(),
        )
        logger.error(
            f"Round-1 analysis failed for video_id={video_id}",
            exc_info=True,
        )
    # Note: cancel-flag cleanup is handled by _run_full_analysis wrapper


# ---------------------------------------------------------------------------
# Round 2 analysis — subtree generation + DB persistence (Slice 07)
# ---------------------------------------------------------------------------


def _run_round2_analysis(video_id: int) -> None:
    """Process each topic from round-1 through round-2 LLM for subtree generation.

    For each topic discovered in round 1:
    1. Extract paragraphs for that topic's range (start_para to end_para).
    2. Concatenate them.
    3. Load ``ob_prompt_round2`` from config, replace ``{topic_name}``.
    4. Call Ollama ``/api/generate``.
    5. Parse subtree JSON.
    6. INSERT into ``course_topics`` table immediately.
    7. Push SSE ``success`` event ("话题 N/M: {topic_name} 完成").

    If a single topic fails: skip it, push ``error`` event, continue with remaining.
    Supports cancel between topics via ``_analysis_cancel_flags``.

    Preconditions:
        ``_topic_segmentation_results[video_id]`` must exist (round 1 succeeded).
    """
    logger = logging.getLogger("outline")

    import sys
    print(f"[PROBE] Round2 entry: video_id={video_id}", file=sys.stderr, flush=True)

    round1_result = _topic_segmentation_results.get(video_id)
    if not round1_result:
        push_log_event("error", f"没有第一轮分析结果，无法执行第二轮 (video_id={video_id})")
        print(f"[PROBE] Round2 abort: no round1 result for video_id={video_id}", file=sys.stderr, flush=True)
        return

    course_name = round1_result.get("course_name", "未知课程")
    topics = round1_result.get("topics", [])
    total = len(topics)
    print(f"[PROBE] Round2: {total} topics to process", file=sys.stderr, flush=True)

    if total == 0:
        push_log_event("success", "没有话题需要生成子树")
        return

    # --- Load config ---
    conn = get_db(str(DB_PATH))
    config_rows = conn.execute(
        "SELECT key, value FROM config WHERE key IN (?, ?, ?, ?)",
        ("ob_prompt_round2", "ob_llm_model", "ob_llm_temperature",
         "ob_llm_num_ctx"),
    ).fetchall()
    conn.close()
    config = {row["key"]: row["value"] for row in config_rows}
    prompt_template = config.get("ob_prompt_round2", _OB_PROMPT_ROUND2)
    model = config.get("ob_llm_model", "qwen2.5:14b-instruct")
    temperature = float(config.get("ob_llm_temperature", 0.0))
    num_ctx = _parse_num_ctx(config.get("ob_llm_num_ctx"))

    # --- Fetch all paragraphs once (for time derivation) ---
    all_paragraphs = _fetch_paragraphs_full_text(video_id)

    completed = 0
    failed = 0

    for i, topic in enumerate(topics, 1):
        # --- Check cancel between topics ---
        if _analysis_cancel_flags.get(video_id):
            push_log_event(
                "info",
                f"第二轮分析已中止 (video_id={video_id})，"
                f"已完成 {completed}/{total} 个话题",
            )
            logger.info(
                f"Round-2 analysis cancelled for video_id={video_id} "
                f"after {completed}/{total} topics"
            )
            return

        topic_name = topic.get("name", f"话题{i}")
        start_para = topic.get("start_para", 0)
        end_para = topic.get("end_para", 0)

        import sys
        print(f"[PROBE] Round2 topic {i}/{total}: {topic_name} start", file=sys.stderr, flush=True)

        push_log_event(
            "progress",
            f"话题 {i}/{total}: {topic_name} 正在生成子树…",
            progress_pct=(i / total) * 100.0,
        )

        try:
            # --- Filter paragraphs for this topic's range ---
            # Defensive: skip if paragraph range is out of bounds
            max_para = max(p["paragraph_index"] for p in all_paragraphs) if all_paragraphs else 0
            if start_para > max_para:
                push_log_event(
                    "error",
                    f"话题 '{topic_name}' 段落范围越界 "
                    f"(start_para={start_para}, max_para={max_para})，已跳过",
                )
                failed += 1
                continue

            topic_paragraphs = [
                p for p in all_paragraphs
                if start_para <= p["paragraph_index"] <= end_para
            ]

            if not topic_paragraphs:
                push_log_event(
                    "error",
                    f"话题 '{topic_name}' 没有找到对应段落 "
                    f"(start_para={start_para}, end_para={end_para})",
                )
                failed += 1
                continue

            # --- Concatenate ---
            topic_text, _ = concat_paragraphs(topic_paragraphs)

            # --- Build prompt with topic name ---
            prompt = prompt_template.replace("{topic_name}", topic_name)
            full_prompt = prompt + "\n\n" + topic_text

            # --- Check cancel before Ollama call ---
            if _analysis_cancel_flags.get(video_id):
                push_log_event(
                    "info",
                    f"第二轮分析已中止 (video_id={video_id})，"
                    f"已完成 {completed}/{total} 个话题",
                )
                return

            # --- Call Ollama ---
            logger.info(
                f"Round-2 calling Ollama model={model} topic={topic_name!r} "
                f"prompt_length={len(full_prompt)}"
            )
            print(f"[PROBE] Ollama call start: topic={topic_name!r} len={len(full_prompt)}", file=sys.stderr, flush=True)
            ollama_response = _call_ollama_generate(model, full_prompt, stream=False, temperature=temperature, num_ctx=num_ctx)
            print(f"[PROBE] Ollama call done: topic={topic_name!r} resp_len={len(ollama_response.get('response',''))}", file=sys.stderr, flush=True)
            raw_output = ollama_response.get("response", "")

            # --- Parse subtree JSON ---
            parsed = parse_round2_response(raw_output, topic_name)
            if parsed is None:
                preview = raw_output[:500]
                push_log_event(
                    "error",
                    f"话题 '{topic_name}' 子树 JSON 解析失败。"
                    f"原始输出 (前500字):\n{preview}",
                )
                logger.error(
                    f"Round-2 JSON parse failed for topic={topic_name!r}, "
                    f"video_id={video_id}. "
                    f"Raw output (first 500 chars): {preview}"
                )
                failed += 1
                continue

            # --- Derive time fields from first/last paragraph ---
            start_time = topic_paragraphs[0]["start_time"]
            end_time = topic_paragraphs[-1]["end_time"]

            # --- Build subtree_json for storage ---
            subtree_json = json.dumps(
                {
                    "course_name": course_name,
                    "subtree": parsed.get("subtree", {"children": []}),
                },
                ensure_ascii=False,
            )

            # --- INSERT into course_topics ---
            print(f"[PROBE] DB insert: topic={topic_name!r}", file=sys.stderr, flush=True)
            conn = get_db(str(DB_PATH))
            try:
                conn.execute(
                    "INSERT INTO course_topics "
                    "(video_id, start_para_index, end_para_index, "
                    "start_time, end_time, topic_name, subtree_json) "
                    "VALUES (?, ?, ?, ?, ?, ?, ?)",
                    (
                        video_id,
                        start_para,
                        end_para,
                        start_time,
                        end_time,
                        topic_name,
                        subtree_json,
                    ),
                )
                conn.execute("COMMIT")
            finally:
                conn.close()

            completed += 1
            push_log_event(
                "success",
                f"话题 {i}/{total}: {topic_name} 完成",
            )
            logger.info(
                f"Round-2 topic completed: {topic_name!r} "
                f"({completed}/{total}), video_id={video_id}"
            )

        except Exception as e:
            push_log_event(
                "error",
                f"话题 '{topic_name}' 处理失败: {e}",
                traceback=traceback.format_exc(),
            )
            logger.error(
                f"Round-2 topic failed: {topic_name!r}, video_id={video_id}: {e}",
                exc_info=True,
            )
            failed += 1
            continue

    # --- Final summary ---
    import sys
    print(f"[PROBE] Round2 summary: completed={completed} failed={failed} total={total}", file=sys.stderr, flush=True)
    if failed == 0:
        push_log_event(
            "success",
            f"全部 {completed} 个话题完成，已写入 course_topics",
        )
    else:
        push_log_event(
            "success",
            f"全部 {total} 个话题处理完成: {completed} 成功, {failed} 失败",
        )
    print(f"[PROBE] Round2 summary sent to SSE", file=sys.stderr, flush=True)
    logger.info(
        f"Round-2 analysis complete for video_id={video_id}: "
        f"completed={completed}, failed={failed}"
    )
    print(f"[PROBE] Round2 exit: video_id={video_id}", file=sys.stderr, flush=True)


# ---------------------------------------------------------------------------
# Combined analysis pipeline (Slice 07)
# ---------------------------------------------------------------------------


def _run_full_analysis(video_id: int) -> None:
    """Background worker that performs the complete two-round LLM analysis.

    Round 1 (Slice 06): topic segmentation → stored in
    ``_topic_segmentation_results``.
    Round 2 (Slice 07): per-topic subtree generation → INSERT into
    ``course_topics``.

    Runs both rounds sequentially in the same background thread.
    The cancel flag (``_analysis_cancel_flags``) is checked between steps
    and between topics; cleanup happens in ``finally``.
    """
    import sys
    print(f"[PROBE] _run_full_analysis start: video_id={video_id}", file=sys.stderr, flush=True)
    try:
        # --- Round 1: Topic Segmentation ---
        print(f"[PROBE] Round 1 starting: video_id={video_id}", file=sys.stderr, flush=True)
        _run_round1_analysis(video_id)
        print(f"[PROBE] Round 1 done: video_id={video_id} result={'OK' if video_id in _topic_segmentation_results else 'FAIL'}", file=sys.stderr, flush=True)

        # If round1 was cancelled or failed, don't proceed to round2
        if _analysis_cancel_flags.get(video_id):
            print(f"[PROBE] Round 2 skipped: cancelled", file=sys.stderr, flush=True)
            return
        if video_id not in _topic_segmentation_results:
            print(f"[PROBE] Round 2 skipped: no topic result", file=sys.stderr, flush=True)
            return

        # --- Round 2: Subtree Generation ---
        print(f"[PROBE] Round 2 starting: video_id={video_id}", file=sys.stderr, flush=True)
        _run_round2_analysis(video_id)
        print(f"[PROBE] Round 2 done: video_id={video_id}", file=sys.stderr, flush=True)
    finally:
        # Clean up cancel flag so the user can re-trigger
        _analysis_cancel_flags.pop(video_id, None)
        print(f"[PROBE] _run_full_analysis exit: video_id={video_id}", file=sys.stderr, flush=True)


# ---------------------------------------------------------------------------
# POST /api/outline/analyze/{video_id} (Slice 06 / Slice 07)
# ---------------------------------------------------------------------------


@app.post("/api/outline/analyze/{video_id}")
async def analyze_video(video_id: int):
    """Start the full two-round LLM analysis for a video.

    Returns 202 Accepted immediately.  A background thread performs:
    1. Round 1: topic segmentation → stored in memory.
    2. Round 2: per-topic subtree generation → INSERT into ``course_topics``.

    SSE events are pushed throughout the process.
    Call POST /api/outline/stop/{video_id} to cancel.
    """
    # --- Validate video exists and has paragraphs ---
    conn = get_db(str(DB_PATH))
    video = conn.execute(
        "SELECT id FROM videos WHERE id = ?", (video_id,)
    ).fetchone()
    if video is None:
        conn.close()
        raise HTTPException(status_code=404, detail="Video not found")

    para_exists = conn.execute(
        "SELECT 1 FROM corpus_paragraphs WHERE video_id = ? LIMIT 1",
        (video_id,),
    ).fetchone()
    conn.close()
    if para_exists is None:
        raise HTTPException(status_code=404, detail="No paragraphs found for this video")

    # --- Check if analysis is already running ---
    if video_id in _analysis_cancel_flags:
        raise HTTPException(
            status_code=409,
            detail="Analysis is already running for this video",
        )

    # --- Mark as running and spawn background thread ---
    _analysis_cancel_flags[video_id] = False  # False = running, not cancelled

    import threading

    thread = threading.Thread(
        target=_run_full_analysis,
        args=(video_id,),
        daemon=True,
    )
    thread.start()

    return JSONResponse(
        status_code=202,
        content={"status": "started", "video_id": video_id},
    )


# ---------------------------------------------------------------------------
# POST /api/outline/stop/{video_id} (Slice 06)
# ---------------------------------------------------------------------------


@app.post("/api/outline/stop/{video_id}")
async def stop_analysis(video_id: int):
    """Request cancellation of a running round-1 analysis.

    Sets the cancel flag for the given video_id.  The worker thread checks
    this flag between major steps and stops if set.  Idempotent — calling
    stop when no analysis is running is a no-op.
    """
    _analysis_cancel_flags[video_id] = True
    return {"status": "cancelling", "video_id": video_id}


@app.post("/api/outline/regenerate/{video_id}")
async def regenerate_analysis(video_id: int):
    """Clear existing course_topics for the video, then re-run full analysis.

    Deletes all ``course_topics`` rows for the given ``video_id``, then
    spawns the same two-round LLM pipeline as ``POST /analyze/{video_id}``.
    Returns 409 if an analysis is already running.
    """
    # Validate video exists
    conn = get_db(str(DB_PATH))
    video = conn.execute(
        "SELECT id FROM videos WHERE id = ?", (video_id,)
    ).fetchone()
    if video is None:
        conn.close()
        raise HTTPException(status_code=404, detail="Video not found")

    para_exists = conn.execute(
        "SELECT 1 FROM corpus_paragraphs WHERE video_id = ? LIMIT 1",
        (video_id,),
    ).fetchone()
    conn.close()
    if para_exists is None:
        raise HTTPException(status_code=404, detail="No paragraphs found for this video")

    # Check if analysis is already running
    if video_id in _analysis_cancel_flags:
        raise HTTPException(
            status_code=409,
            detail="Analysis is already running for this video",
        )

    # Clear old topics
    conn = get_db(str(DB_PATH))
    conn.execute(
        "DELETE FROM course_topics WHERE video_id = ?", (video_id,)
    )
    conn.execute("COMMIT")
    deleted = conn.total_changes
    conn.close()

    push_log_event(
        "info",
        f"已清空 {deleted} 条旧话题记录，开始重新分析…",
    )

    # Clean up any stale round-1 result in memory
    _topic_segmentation_results.pop(video_id, None)

    # Mark as running and spawn thread
    _analysis_cancel_flags[video_id] = False

    import threading
    thread = threading.Thread(
        target=_run_full_analysis,
        args=(video_id,),
        daemon=True,
    )
    thread.start()

    return JSONResponse(
        status_code=202,
        content={"status": "started", "video_id": video_id, "cleared": deleted},
    )


# ---------------------------------------------------------------------------
# CK Dimension — Knowledge Graph Topology (CKG Analysis)
# ---------------------------------------------------------------------------


def _normalize_tree_children(children: list) -> list:
    """Recursively normalize subtree children, filtering blank/empty nodes.

    Mirrors the JS convertChildren() logic to produce a clean tree for
    topology computation.
    """
    result = []
    for child in children:
        name = (child.get("name") or child.get("topic") or "").strip()
        kids = child.get("children") or []
        if not name and not kids:
            continue  # drop truly empty leaf nodes (LLM artefacts)
        node = {
            "name": name or "···",
            "children": _normalize_tree_children(kids) if kids else [],
        }
        result.append(node)
    return result


def _build_tree_for_topology(rows: list[dict]) -> dict | None:
    """Build the full knowledge tree from course_topics DB rows.

    Mirrors the JS buildFullTree() logic so topology is computed on the
    same tree that is visualised in the mind map.

    Args:
        rows: dicts with keys ``topic_name`` and ``subtree_json``.

    Returns:
        A tree dict ``{name, children}`` rooted at the course name,
        or ``None`` if rows is empty.
    """
    if not rows:
        return None

    course_name = "课程"
    try:
        first = json.loads(rows[0].get("subtree_json") or "{}")
        if first.get("course_name"):
            course_name = first["course_name"]
    except Exception:
        pass

    root: dict = {"name": course_name, "children": []}
    for row in rows:
        topic_node: dict = {"name": row["topic_name"], "children": []}
        try:
            subtree_data = json.loads(row.get("subtree_json") or "{}")
            raw_children = (subtree_data.get("subtree") or {}).get("children") or []
            topic_node["children"] = _normalize_tree_children(raw_children)
        except Exception:
            pass
        root["children"].append(topic_node)
    return root


def compute_ckg_topology(root: dict) -> dict:
    """Compute CKG topological parameters from a knowledge tree.

    Traverses the tree via BFS, builds an undirected adjacency list, then
    computes each metric.  All operations are O(N) or O(N²) — the latter
    (average path length) is acceptable for typical course trees (< 500 nodes).

    Parameters computed (matching the paper's CK dimension):
    - depth: max depth from root (root = level 0).
    - branch_factor: mean children count per non-leaf node.
    - avg_path_length: mean shortest path between all ordered node pairs.
    - node_count: total number of nodes (including root).
    - leaf_count: nodes with no children.
    """
    from collections import deque

    if not root:
        return {
            "depth": 0, "branch_factor": 0.0,
            "avg_path_length": 0.0,
            "node_count": 0, "leaf_count": 0,
        }

    # BFS: assign integer IDs, build undirected adjacency list
    adj: dict[int, list[int]] = {}
    depths: dict[int, int] = {}
    child_counts: dict[int, int] = {}

    node_id_counter = 0
    queue: deque = deque([(root, -1, 0)])

    while queue:
        node, parent_id, depth = queue.popleft()
        nid = node_id_counter
        node_id_counter += 1

        depths[nid] = depth
        adj[nid] = []

        if parent_id >= 0:
            adj[nid].append(parent_id)
            adj[parent_id].append(nid)
            child_counts[parent_id] = child_counts.get(parent_id, 0) + 1

        for child in (node.get("children") or []):
            queue.append((child, nid, depth + 1))

    N = node_id_counter
    if N == 0:
        return {
            "depth": 0, "branch_factor": 0.0,
            "avg_path_length": 0.0,
            "node_count": 0, "leaf_count": 0,
        }

    max_depth = max(depths.values())
    leaf_count = sum(1 for nid in range(N) if nid not in child_counts)

    non_leaf = [nid for nid in range(N) if nid in child_counts]
    branch_factor = (
        sum(child_counts[nid] for nid in non_leaf) / len(non_leaf)
        if non_leaf else 0.0
    )

    # Average path length: BFS from every node, O(N²)
    if N <= 1:
        avg_path_length = 0.0
    else:
        total_dist = 0
        for start in range(N):
            visited: dict[int, int] = {start: 0}
            bfs_q: deque = deque([start])
            while bfs_q:
                curr = bfs_q.popleft()
                for nb in adj[curr]:
                    if nb not in visited:
                        visited[nb] = visited[curr] + 1
                        total_dist += visited[nb]
                        bfs_q.append(nb)
        avg_path_length = total_dist / (N * (N - 1))

    return {
        "depth": max_depth,
        "branch_factor": round(branch_factor, 2),
        "avg_path_length": round(avg_path_length, 2),
        "node_count": N,
        "leaf_count": leaf_count,
    }


def _interpret_ck_style(topology: dict) -> dict:
    """Derive a style label and descriptive text from CKG topology params.

    Thresholds are set based on the paper's illustrative examples:
    depth ≥ 5 → "deep"; branch_factor ≥ 3.0 → "divergent".
    """
    depth = topology.get("depth", 0)
    branch = topology.get("branch_factor", 0.0)
    avg_path = topology.get("avg_path_length", 0.0)
    n_nodes = topology.get("node_count", 0)

    deep = depth >= 5
    branchy = branch >= 3.0

    if deep and not branchy:
        style_key = "deep"
        style_label = "深度递进型"
        description = (
            f"该教师的内容组织以纵向深入为主（深度={depth:.2f}），横向发散较少（分支因子={branch:.2f}）。"
            f"知识体系层层深入，概念链条清晰，平均路径长度为 {avg_path:.2f}，适合学生逐步构建深层理解。"
        )
    elif not deep and branchy:
        style_key = "divergent"
        style_label = "横向发散型"
        description = (
            f"该教师的内容组织以横向展开为主（分支因子={branch:.2f}），知识深度适中（深度={depth:.2f}）。"
            f"同一层级覆盖多个子话题，共 {n_nodes:.2f} 个知识节点，信息密度高，适合广覆盖式的知识介绍。"
        )
    elif deep and branchy:
        style_key = "balanced"
        style_label = "立体均衡型"
        description = (
            f"该教师的内容组织在深度（深度={depth:.2f}）和广度（分支因子={branch:.2f}）上均较突出。"
            f"知识体系共 {n_nodes:.2f} 个节点，既有纵向深入的概念链，也有丰富的横向分支，结构立体而全面。"
        )
    else:
        style_key = "compact"
        style_label = "简洁聚焦型"
        description = (
            f"该教师的内容组织相对简洁（深度={depth:.2f}，分支因子={branch:.2f}），"
            f"共 {n_nodes:.2f} 个知识节点。知识结构紧凑聚焦，重点突出，适合核心概念的精讲与强化。"
        )

    return {"style_key": style_key, "style_label": style_label, "description": description}


@app.get("/api/outline/video/{video_id}/ck-profile")
async def get_ck_profile(video_id: int):
    """Return CK dimension topology profile for a video's knowledge tree.

    Fetches all completed topics for the video, rebuilds the full knowledge
    tree (mirroring the mind map visualisation), computes CKG topological
    parameters, and returns them together with an interpretive style label.

    Returns 404 if no topics with subtree data exist for the video.
    """
    conn = get_db(str(DB_PATH))
    try:
        rows = conn.execute(
            "SELECT topic_name, subtree_json FROM course_topics "
            "WHERE video_id = ? AND subtree_json IS NOT NULL "
            "ORDER BY start_para_index",
            (video_id,),
        ).fetchall()
    finally:
        conn.close()

    if not rows:
        raise HTTPException(
            status_code=404,
            detail="No completed topics found — run analysis first",
        )

    rows_list = [dict(row) for row in rows]
    tree = _build_tree_for_topology(rows_list)
    if not tree:
        raise HTTPException(status_code=500, detail="Failed to build knowledge tree")

    topology = compute_ckg_topology(tree)
    style = _interpret_ck_style(topology)

    return {**topology, **style}


# ---------------------------------------------------------------------------
# CKG concept extraction — Step 1 (issue 01)
# ---------------------------------------------------------------------------

_ckg_cancel_flags: Dict[int, bool] = {}
"""Tracks cancel requests for CKG extraction, keyed by video_id.
True = cancel requested; False = running; absent = no extraction running."""

def parse_concepts_response(raw: str) -> dict | None:
    """Parse the LLM's raw Step-1 output into a concepts structure.

    Robust against common LLM output quirks (markdown fences, leading/trailing
    prose, whitespace/BOM), mirroring ``parse_round1_response``.

    Returns ``{"title": <str>, "concepts": [{"name", "definition",
    "first_para", "parent"}, ...]}`` with names de-duplicated (case/space-
    insensitive, first occurrence wins), or ``None`` if the response cannot be
    parsed into a valid concepts list. ``parent`` is the name of the higher
    concept (or the lecture title) this concept is decomposed from; it defaults
    to ``""`` when absent. ``title`` defaults to ``""`` when absent.
    """
    import re

    text = (raw or "").strip()
    if not text:
        return None

    # 1) Strip markdown code fences if present
    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    # 2) Narrow to first '{' ... last '}'
    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace : last_brace + 1]

    # 3) Parse JSON; on failure, salvage flat objects from the broken text.
    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        # The LLM occasionally corrupts the array mid-stream (dropping a '{' /
        # '"name":'), which poisons json.loads for the whole payload even though
        # the leading objects are perfectly fine. Concept objects are flat (no
        # nested braces), so we can re-capture each {...} and keep the ones that
        # carry a non-empty name.
        salvaged = []
        for m in re.findall(r"\{[^{}]*\}", text):
            try:
                obj = json.loads(m)
            except (json.JSONDecodeError, ValueError):
                continue
            if isinstance(obj, dict) and str(obj.get("name", "")).strip():
                salvaged.append(obj)
        if not salvaged:
            return None
        data = {"concepts": salvaged}

    # 4) Validate structure
    if not isinstance(data, dict):
        return None
    if not isinstance(data.get("concepts"), list):
        return None

    # 5) Normalize + de-duplicate by normalized name
    seen: set[str] = set()
    concepts: list[dict] = []
    for c in data["concepts"]:
        if not isinstance(c, dict):
            continue
        name = str(c.get("name", "")).strip()
        if not name:
            continue
        key = name.lower().strip()
        if key in seen:
            continue
        seen.add(key)

        definition = str(c.get("definition", "")).strip()
        raw_fp = c.get("first_para")
        try:
            first_para = int(raw_fp)
        except (TypeError, ValueError):
            first_para = 0

        parent = str(c.get("parent", "")).strip()

        concepts.append(
            {
                "name": name,
                "definition": definition,
                "first_para": first_para,
                "parent": parent,
            }
        )

    if not concepts:
        return None

    title = ""
    if isinstance(data, dict):
        title = str(data.get("title", "")).strip()

    return {"title": title, "concepts": concepts}


def parse_edges_response(raw: str) -> dict | None:
    """Parse the LLM's raw Step-2 output into a prerequisite-edge structure.

    Robust against markdown fences / surrounding prose, mirroring
    ``parse_concepts_response``.

    Returns ``{"edges": [{"from", "to"}, ...]}`` keeping only edges with both
    non-empty ``from`` and ``to`` strings (malformed entries are skipped). An
    empty edge list is valid output. Returns ``None`` only when the response
    cannot be parsed into a JSON object carrying an ``edges`` list.
    """
    import re

    text = (raw or "").strip()
    if not text:
        return None

    # 1) Strip markdown code fences if present
    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    # 2) Narrow to first '{' ... last '}'
    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace : last_brace + 1]

    # 3) Parse JSON; on failure, salvage flat edge objects from the broken text.
    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        # Edge objects are flat ({"from","to"}); re-capture each {...} and keep
        # the well-formed ones so a corrupt tail doesn't discard valid edges.
        salvaged = []
        for m in re.findall(r"\{[^{}]*\}", text):
            try:
                obj = json.loads(m)
            except (json.JSONDecodeError, ValueError):
                continue
            if (isinstance(obj, dict)
                    and str(obj.get("from", "")).strip()
                    and str(obj.get("to", "")).strip()):
                salvaged.append(obj)
        if not salvaged:
            return None
        data = {"edges": salvaged}

    # 4) Validate structure
    if not isinstance(data, dict):
        return None
    if not isinstance(data.get("edges"), list):
        return None

    # 5) Keep only well-formed edges
    edges: list[dict] = []
    for e in data["edges"]:
        if not isinstance(e, dict):
            continue
        src = str(e.get("from", "")).strip()
        dst = str(e.get("to", "")).strip()
        if not src or not dst:
            continue
        edges.append({"from": src, "to": dst})

    return {"edges": edges}


def validate_edges(edges: list[dict], concepts: list[dict]) -> list[dict]:
    """Drop edges referencing unknown concepts; normalize endpoint names.

    Matching is by normalized name (lower-cased, stripped). Endpoints are
    rewritten to the canonical concept name so the graph and edges agree.
    Self-loops (from == to after normalization) are dropped.
    """
    canonical: dict[str, str] = {}
    for c in concepts:
        name = str(c.get("name", "")).strip()
        if name:
            canonical.setdefault(name.lower(), name)

    valid: list[dict] = []
    for e in edges:
        src_key = str(e.get("from", "")).strip().lower()
        dst_key = str(e.get("to", "")).strip().lower()
        if src_key not in canonical or dst_key not in canonical:
            continue
        if src_key == dst_key:
            continue
        valid.append({"from": canonical[src_key], "to": canonical[dst_key]})
    return valid


def parse_relations_response(raw: str) -> dict | None:
    """Parse the LLM's grounded cross-relation (讲述关联) output.

    Robust against markdown fences / surrounding prose / corrupt tails, mirroring
    ``parse_edges_response``. Returns ``{"relations": [{"from", "to", "type",
    "evidence"}, ...]}`` keeping only entries with non-empty ``from`` and
    ``to``. ``type`` is the LLM's FREE-LABELLED relation kind, stored verbatim
    (default ``""``); ``evidence`` is optional grounding text (default ``""``).
    An empty relation list is valid. Returns ``None`` only when the response
    cannot be parsed into a JSON object carrying a ``relations`` list.
    """
    import re

    text = (raw or "").strip()
    if not text:
        return None

    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace : last_brace + 1]

    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        # Relation objects are flat; re-capture each {...} and keep well-formed
        # ones so a corrupt tail doesn't discard valid relations.
        salvaged = []
        for m in re.findall(r"\{[^{}]*\}", text):
            try:
                obj = json.loads(m)
            except (json.JSONDecodeError, ValueError):
                continue
            if (isinstance(obj, dict)
                    and str(obj.get("from", "")).strip()
                    and str(obj.get("to", "")).strip()):
                salvaged.append(obj)
        if not salvaged:
            return None
        data = {"relations": salvaged}

    if not isinstance(data, dict):
        return None
    if not isinstance(data.get("relations"), list):
        return None

    relations: list[dict] = []
    for r in data["relations"]:
        if not isinstance(r, dict):
            continue
        src = str(r.get("from", "")).strip()
        dst = str(r.get("to", "")).strip()
        if not src or not dst:
            continue
        relations.append(
            {
                "from": src,
                "to": dst,
                "type": str(r.get("type", "")).strip(),
                "evidence": str(r.get("evidence", "")).strip(),
            }
        )

    return {"relations": relations}


def validate_relations(relations: list[dict], concepts: list[dict]) -> list[dict]:
    """Drop cross-relations referencing unknown concepts; normalize endpoints.

    Like ``validate_edges`` but preserves the free-labelled ``type`` and
    ``evidence`` fields. Endpoints are rewritten to canonical concept names;
    self-loops (from == to after normalization) are dropped. Relations are NOT
    de-cycled — narrative links may legitimately form cycles.
    """
    canonical: dict[str, str] = {}
    for c in concepts:
        name = str(c.get("name", "")).strip()
        if name:
            canonical.setdefault(name.lower(), name)

    valid: list[dict] = []
    for r in relations:
        src_key = str(r.get("from", "")).strip().lower()
        dst_key = str(r.get("to", "")).strip().lower()
        if src_key not in canonical or dst_key not in canonical:
            continue
        if src_key == dst_key:
            continue
        valid.append(
            {
                "from": canonical[src_key],
                "to": canonical[dst_key],
                "type": str(r.get("type", "")).strip(),
                "evidence": str(r.get("evidence", "")).strip(),
            }
        )
    return valid


def compute_convergence_count(
    concepts: list[dict],
    decomposition_edges: list[dict],
    relations: list[dict],
) -> int:
    """Number of concepts reached by more than one incoming edge (issue 03).

    In-degree is counted over BOTH the decomposition backbone AND the grounded
    relations. In a pure decomposition tree every node has in-degree ≤ 1, so a
    convergence (in-degree > 1) can only arise when a relation points into a
    concept that already has a parent (or another relation) — i.e. the teacher
    arrives at that concept from multiple directions.
    """
    names = {str(c.get("name", "")).strip().lower() for c in concepts}
    names.discard("")
    indeg: dict[str, int] = {}
    for e in list(decomposition_edges) + list(relations):
        src = str(e.get("from", "")).strip().lower()
        dst = str(e.get("to", "")).strip().lower()
        if dst in names and src in names and src != dst:
            indeg[dst] = indeg.get(dst, 0) + 1
    return sum(1 for v in indeg.values() if v > 1)


def compute_relation_density(concepts: list[dict], relations: list[dict]) -> float:
    """讲述关联密度 = grounded-relation count / concept count (issue 03).

    0.0 when there are no concepts. Measures how much the teacher horizontally
    connects concepts on top of the decomposition.
    """
    n = len([c for c in concepts if str(c.get("name", "")).strip()])
    if n == 0:
        return 0.0
    return len(relations) / n


def compute_dag_topology(concepts: list[dict], edges: list[dict]) -> dict:
    """Compute directed-graph topology parameters for a real CKG.

    Brand-new for the prerequisite DAG (issue 03) — do NOT confuse with the
    undirected-tree ``compute_ckg_topology`` scaffold (ADR 0002).

    Inputs are ``graph_json``'s concepts (each a dict with ``name``) and edges
    (each ``{"from", "to"}``). Edges referencing unknown nodes are ignored so
    counts stay consistent with the node set.

    Returns a dict with:
    - ``depth``: edge count of the longest directed path from a source
      (in-degree 0) to a sink (out-degree 0). 0 for empty / edge-less graphs.
    - ``branch_factor``: mean out-degree of non-sink nodes (out-degree > 0).
    - ``convergence_count``: number of nodes with in-degree > 1 (the DAG
      signature absent in trees).
    - ``density``: edges / (N*(N-1)); 0 when N < 2.
    - ``avg_path_length``: mean shortest path over all reachable ordered pairs
      on the UNDIRECTED projection (unreachable pairs excluded).
    - ``clustering``: average clustering coefficient on the undirected
      projection (a second-order observable — may be 0 on tree-like data).
    """
    from collections import defaultdict, deque

    names = [c.get("name") for c in concepts if c.get("name") is not None]
    node_set = set(names)
    N = len(node_set)

    # Directed adjacency + degree bookkeeping (known nodes only).
    out_adj: dict = defaultdict(list)
    out_deg: dict = defaultdict(int)
    in_deg: dict = defaultdict(int)
    valid_edges = []
    for e in edges:
        src, dst = e.get("from"), e.get("to")
        if src in node_set and dst in node_set:
            out_adj[src].append(dst)
            out_deg[src] += 1
            in_deg[dst] += 1
            valid_edges.append((src, dst))

    E = len(valid_edges)

    # depth = longest directed path (edge count). DAG → DP via topo order;
    # but break_cycles already guarantees acyclicity, so memoised DFS is safe.
    longest_cache: dict = {}

    def longest_from(node: str) -> int:
        if node in longest_cache:
            return longest_cache[node]
        best = 0
        for nxt in out_adj[node]:
            best = max(best, 1 + longest_from(nxt))
        longest_cache[node] = best
        return best

    depth = max((longest_from(n) for n in node_set), default=0)

    # branch_factor = mean out-degree over non-sink nodes (out-degree > 0).
    non_sinks = [n for n in node_set if out_deg[n] > 0]
    branch_factor = (
        sum(out_deg[n] for n in non_sinks) / len(non_sinks)
        if non_sinks else 0.0
    )

    # convergence_count = nodes with in-degree > 1.
    convergence_count = sum(1 for n in node_set if in_deg[n] > 1)

    # density = E / (N*(N-1)).
    density = E / (N * (N - 1)) if N >= 2 else 0.0

    # Undirected projection for path length + clustering.
    und: dict = defaultdict(set)
    for src, dst in valid_edges:
        und[src].add(dst)
        und[dst].add(src)

    # avg_path_length: BFS from each node over the undirected projection,
    # averaging shortest paths over reachable ordered pairs only.
    total_dist = 0
    reachable_pairs = 0
    for start in node_set:
        dist = {start: 0}
        q = deque([start])
        while q:
            cur = q.popleft()
            for nb in und[cur]:
                if nb not in dist:
                    dist[nb] = dist[cur] + 1
                    total_dist += dist[nb]
                    reachable_pairs += 1
                    q.append(nb)
    avg_path_length = (total_dist / reachable_pairs) if reachable_pairs else 0.0

    # clustering: average local clustering coefficient on undirected projection.
    if N == 0:
        clustering = 0.0
    else:
        coeffs = []
        for n in node_set:
            neigh = und[n]
            k = len(neigh)
            if k < 2:
                coeffs.append(0.0)
                continue
            links = 0
            neigh_list = list(neigh)
            for i in range(len(neigh_list)):
                for j in range(i + 1, len(neigh_list)):
                    if neigh_list[j] in und[neigh_list[i]]:
                        links += 1
            coeffs.append((2 * links) / (k * (k - 1)))
        clustering = sum(coeffs) / N

    return {
        "depth": depth,
        "branch_factor": branch_factor,
        "convergence_count": convergence_count,
        "density": density,
        "avg_path_length": avg_path_length,
        "clustering": clustering,
    }


def compute_bottomup_ratio(concepts: list[dict], edges: list[dict]):
    """Compute the delivery-direction "bottom-up ratio" (讲授走向).

    For each prerequisite edge A->B (A is the prerequisite), compare the two
    concepts' ``first_para``:
    - the edge is "prerequisite-taught-first" when first_para(A) < first_para(B).

    ``bottomup_ratio`` = (# prerequisite-taught-first edges) /
    (# edges whose endpoints have DIFFERENT first_para). Edges where the two
    first_para values are equal are excluded from BOTH numerator and
    denominator. Returns ``None`` when there are no valid (unequal) edges, so
    the front-end can render a dash.
    """
    fp = {
        c.get("name"): c.get("first_para")
        for c in concepts
        if c.get("name") is not None
    }

    valid = 0
    bottomup = 0
    for e in edges:
        a, b = e.get("from"), e.get("to")
        fa, fb = fp.get(a), fp.get(b)
        if fa is None or fb is None:
            continue
        if fa == fb:
            continue  # equal first_para excluded from numerator and denominator
        valid += 1
        if fa < fb:
            bottomup += 1

    if valid == 0:
        return None
    return bottomup / valid


def break_cycles(edges: list[dict]) -> list[dict]:
    """Return a subset of ``edges`` that forms a DAG.

    Greedy incremental construction: process edges in order, adding each only
    if doing so does NOT create a cycle (i.e. there is no existing directed
    path from the edge's ``to`` back to its ``from``). Any edge that would
    close a cycle is dropped. This guarantees the result is acyclic regardless
    of what the LLM produced — the no-cycle invariant comes from code.
    """
    from collections import defaultdict, deque

    adj: dict[str, set] = defaultdict(set)

    def reachable(start: str, target: str) -> bool:
        """True if ``target`` is reachable from ``start`` along current edges."""
        if start == target:
            return True
        seen = {start}
        queue = deque([start])
        while queue:
            node = queue.popleft()
            for nxt in adj[node]:
                if nxt == target:
                    return True
                if nxt not in seen:
                    seen.add(nxt)
                    queue.append(nxt)
        return False

    kept: list[dict] = []
    for e in edges:
        src, dst = e["from"], e["to"]
        # Adding src->dst closes a cycle iff src is already reachable from dst.
        if reachable(dst, src):
            continue
        adj[src].add(dst)
        kept.append(e)
    return kept


def parse_lesson_outline_response(raw: str) -> dict | None:
    """Parse the style-aware lesson-generation LLM output.

    Robust against markdown fences / surrounding prose (mirrors
    ``parse_concepts_response``): strip fences, narrow to first '{' … last '}',
    then ``json.loads``. Returns ``{"outline": [...], "sequence": [...]}`` with
    a recursively-sanitised outline (each node ``{"name", "children"}``) and a
    list of string ``sequence`` names, or ``None`` if it cannot be parsed into
    an object carrying either an outline list or a sequence list.
    """
    import re

    text = (raw or "").strip()
    if not text:
        return None

    fence_match = re.search(r"```(?:json)?\s*([\s\S]*?)```", text)
    if fence_match:
        text = fence_match.group(1).strip()

    first_brace = text.find("{")
    last_brace = text.rfind("}")
    if first_brace != -1 and last_brace != -1 and last_brace > first_brace:
        text = text[first_brace: last_brace + 1]

    try:
        data = json.loads(text)
    except (json.JSONDecodeError, ValueError):
        return None
    if not isinstance(data, dict):
        return None

    def _clean_nodes(nodes) -> list:
        cleaned = []
        if not isinstance(nodes, list):
            return cleaned
        for n in nodes:
            if not isinstance(n, dict):
                continue
            name = str(n.get("name", "")).strip()
            if not name:
                continue
            cleaned.append({
                "name": name,
                "children": _clean_nodes(n.get("children")),
            })
        return cleaned

    outline = _clean_nodes(data.get("outline"))

    sequence = []
    if isinstance(data.get("sequence"), list):
        for s in data["sequence"]:
            name = str(s).strip()
            if name:
                sequence.append(name)

    # Need at least one of the two to be a usable result.
    if not outline and not sequence:
        return None

    return {"outline": outline, "sequence": sequence}


def derive_decomposition_edges(concepts: list[dict]) -> list[dict]:
    """Build the decomposition backbone (parent → child edges) from concepts.

    Each concept may carry a ``parent`` naming the higher concept it is carved
    out of. An edge ``{"from": parent, "to": child}`` is emitted only when the
    parent resolves (case/space-insensitive) to ANOTHER concept in the list —
    top-level concepts whose parent is the lecture title, empty, or unknown
    simply have no incoming edge (they are roots of the decomposition forest).
    Endpoints are rewritten to canonical concept names; self-loops and
    duplicates are dropped. The result is guaranteed acyclic via break_cycles.
    """
    canonical: dict[str, str] = {}
    for c in concepts:
        name = str(c.get("name", "")).strip()
        if name:
            canonical.setdefault(name.lower().strip(), name)

    edges: list[dict] = []
    seen: set[tuple[str, str]] = set()
    for c in concepts:
        child = str(c.get("name", "")).strip()
        parent_raw = str(c.get("parent", "")).strip()
        if not child or not parent_raw:
            continue
        parent = canonical.get(parent_raw.lower().strip())
        if not parent:
            continue  # parent is the title / unknown → child is a root
        if parent.lower().strip() == child.lower().strip():
            continue  # self-loop
        key = (parent, child)
        if key in seen:
            continue
        seen.add(key)
        edges.append({"from": parent, "to": child})

    return break_cycles(edges)


def _run_ckg_extraction(video_id: int) -> None:
    """Background worker performing CKG decomposition extraction (issue 01).

    Steps: fetch paragraphs → concat full text → load prompt + model from
    config → call Ollama once with the full transcript → parse the
    title→concept DECOMPOSITION tree → derive decomposition edges (parent→child)
    → persist to ``course_ckg``.  SSE events are pushed throughout.  Parse
    failures push an ``error`` event carrying the first 500 chars of the raw
    LLM output and do NOT write to the DB or crash. Grounded cross-relations
    (讲述关联) are added by issue 02; ``relations`` is persisted empty here.
    """
    logger = logging.getLogger("outline")

    try:
        if _ckg_cancel_flags.get(video_id):
            push_log_event("info", f"概念抽取已中止 (video_id={video_id})")
            return

        # --- Step 1: fetch paragraphs ---
        push_log_event("info", "正在加载段落全文…", progress_pct=5.0)
        paragraphs = _fetch_paragraphs_full_text(video_id)
        if not paragraphs:
            push_log_event("error", f"视频 {video_id} 没有段落数据")
            return

        para_count = len(paragraphs)
        push_log_event(
            "info",
            f"已加载 {para_count} 个段落，正在拼接全文…",
            progress_pct=15.0,
        )

        if _ckg_cancel_flags.get(video_id):
            push_log_event("info", f"概念抽取已中止 (video_id={video_id})")
            return

        # --- Step 2: concatenate ---
        full_text, _para_map = concat_paragraphs(paragraphs)

        # --- Step 3: load config ---
        conn = get_db(str(DB_PATH))
        config_rows = conn.execute(
            "SELECT key, value FROM config WHERE key IN (?, ?, ?, ?, ?)",
            ("ck_prompt_concepts", "ck_prompt_relations",
             "ob_llm_model", "ob_llm_temperature", "ob_llm_num_ctx"),
        ).fetchall()
        conn.close()
        config = {row["key"]: row["value"] for row in config_rows}
        prompt_template = config.get("ck_prompt_concepts", _CK_PROMPT_CONCEPTS)
        model = config.get("ob_llm_model", "qwen2.5:14b-instruct")
        temperature = float(config.get("ob_llm_temperature", 0.0))
        num_ctx = _parse_num_ctx(config.get("ob_llm_num_ctx"))

        full_prompt = prompt_template + "\n\n" + full_text

        # --- Step 4: call Ollama ---
        push_log_event("info", "调用 LLM 抽取概念…", progress_pct=25.0)
        logger.info(
            f"CKG extraction calling Ollama model={model} "
            f"prompt_length={len(full_prompt)} video_id={video_id}"
        )

        if _ckg_cancel_flags.get(video_id):
            push_log_event("info", f"概念抽取已中止 (video_id={video_id})")
            return

        try:
            ollama_response = _call_ollama_generate(
                model, full_prompt, stream=False, temperature=temperature,
                fmt="json", num_predict=4096, num_ctx=num_ctx,
            )
        except Exception as e:
            push_log_event(
                "error",
                f"Ollama 调用失败: {e}",
                traceback=traceback.format_exc(),
            )
            logger.error(f"CKG Ollama call failed: {e}", exc_info=True)
            return

        if _ckg_cancel_flags.get(video_id):
            push_log_event("info", f"概念抽取已中止 (video_id={video_id})")
            return

        # --- Step 5: parse ---
        push_log_event("info", "正在解析概念清单…", progress_pct=80.0)
        raw_output = ollama_response.get("response", "")
        parsed = parse_concepts_response(raw_output)
        if parsed is None:
            preview = raw_output[:500]
            push_log_event(
                "error",
                f"概念 JSON 解析失败。原始输出 (前500字):\n{preview}",
            )
            logger.error(
                f"CKG concept parse failed for video_id={video_id}. "
                f"Raw output (first 500 chars): {preview}"
            )
            return

        concepts = parsed["concepts"]
        title = parsed.get("title", "")

        # --- Step 6: derive the decomposition backbone (parent → child) ---
        push_log_event("info", "正在构建拆解树…", progress_pct=80.0)
        decomposition_edges = derive_decomposition_edges(concepts)

        # --- Step 6b: extract grounded cross-relations (讲述关联) ---
        relations: list[dict] = []
        if _ckg_cancel_flags.get(video_id):
            push_log_event("info", f"概念抽取已中止 (video_id={video_id})")
            return
        push_log_event("info", "正在抽取讲述关联…", progress_pct=88.0)
        rel_prompt_template = config.get(
            "ck_prompt_relations", _CK_PROMPT_RELATIONS
        )
        concept_lines = "\n".join(
            f"- {c['name']}: {c.get('definition', '')}" for c in concepts
        )
        rel_prompt = (
            rel_prompt_template + "\n" + concept_lines
            + "\n\nFull transcript:\n" + full_text
        )
        try:
            rel_response = _call_ollama_generate(
                model, rel_prompt, stream=False, temperature=temperature,
                fmt="json", num_predict=4096, num_ctx=num_ctx,
            )
            rel_raw = rel_response.get("response", "")
            rel_parsed = parse_relations_response(rel_raw)
            if rel_parsed is None:
                push_log_event(
                    "info",
                    "讲述关联解析失败，本课图将只含拆解骨架（概念照常保存）。"
                    f"原始输出 (前500字):\n{rel_raw[:500]}",
                )
                logger.warning(
                    f"CKG relation parse failed for video_id={video_id}; "
                    f"persisting decomposition only."
                )
            else:
                relations = validate_relations(rel_parsed["relations"], concepts)
        except Exception as e:
            # A relation-step failure must not lose the concepts / backbone.
            push_log_event(
                "info",
                f"讲述关联抽取失败 ({e})，本课图将只含拆解骨架（概念照常保存）。",
            )
            logger.warning(
                f"CKG relation extraction failed for video_id={video_id}: {e}"
            )

        # --- Step 7: compute the decomposition-style parameters (issue 03) ---
        # Core: 拆解深度 (depth) / 拆解宽度 (branch_factor) on the decomposition
        # backbone; 关联密度 (relation_density); 汇聚数 (convergence over
        # decomposition + relations). density/avg_path_length/clustering are
        # RETIRED (collinear with depth/breadth on a tree). bottomup_ratio is an
        # AUXILIARY delivery-direction signal, not a core parameter.
        topology = compute_dag_topology(concepts, decomposition_edges)
        decompose_depth = topology["depth"]
        decompose_breadth = topology["branch_factor"]
        convergence_count = compute_convergence_count(
            concepts, decomposition_edges, relations
        )
        relation_density = compute_relation_density(concepts, relations)
        bottomup_ratio = compute_bottomup_ratio(concepts, decomposition_edges)

        # --- Step 8: persist to course_ckg ---
        graph_json = json.dumps(
            {
                "schema_version": 2,
                "title": title,
                "concepts": concepts,
                "decomposition_edges": decomposition_edges,
                "relations": relations,
            },
            ensure_ascii=False,
        )
        created_at = datetime.now(timezone.utc).isoformat()
        conn = get_db(str(DB_PATH))
        try:
            # density / avg_path_length / clustering are RETIRED (issue 03):
            # persisted as NULL so the columns stay but carry no signal.
            conn.execute(
                "INSERT INTO course_ckg "
                "(video_id, graph_json, model, created_at, "
                "depth, branch_factor, convergence_count, relation_density, "
                "density, avg_path_length, clustering, bottomup_ratio) "
                "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?) "
                "ON CONFLICT(video_id) DO UPDATE SET "
                "graph_json = excluded.graph_json, "
                "model = excluded.model, "
                "created_at = excluded.created_at, "
                "depth = excluded.depth, "
                "branch_factor = excluded.branch_factor, "
                "convergence_count = excluded.convergence_count, "
                "relation_density = excluded.relation_density, "
                "density = excluded.density, "
                "avg_path_length = excluded.avg_path_length, "
                "clustering = excluded.clustering, "
                "bottomup_ratio = excluded.bottomup_ratio",
                (
                    video_id, graph_json, model, created_at,
                    decompose_depth, decompose_breadth,
                    convergence_count, relation_density,
                    None, None, None,
                    bottomup_ratio,
                ),
            )
            conn.execute("COMMIT")
        finally:
            conn.close()

        push_log_event(
            "success",
            f"CKG 抽取完成: 共 {len(concepts)} 个概念、"
            f"{len(decomposition_edges)} 条拆解边、"
            f"{len(relations)} 条讲述关联，已写入 course_ckg",
            progress_pct=100.0,
        )
        logger.info(
            f"CKG extraction complete for video_id={video_id}: "
            f"concepts={len(concepts)} decomposition_edges={len(decomposition_edges)} "
            f"relations={len(relations)}"
        )

    except Exception:
        push_log_event(
            "error",
            f"概念抽取过程发生异常 (video_id={video_id})",
            traceback=traceback.format_exc(),
        )
        logger.error(
            f"CKG extraction failed for video_id={video_id}", exc_info=True
        )
    finally:
        _ckg_cancel_flags.pop(video_id, None)


@app.post("/api/outline/ckg/{video_id}")
async def extract_ckg(video_id: int):
    """Start Step-1 CKG concept extraction for a video in a background thread.

    Returns 202 Accepted immediately; progress is pushed over the shared SSE
    log stream (``/api/stream/logs/outline``).  Returns 404 if the video does
    not exist or has no paragraphs, 409 if extraction is already running.
    """
    conn = get_db(str(DB_PATH))
    video = conn.execute(
        "SELECT id FROM videos WHERE id = ?", (video_id,)
    ).fetchone()
    if video is None:
        conn.close()
        raise HTTPException(status_code=404, detail="Video not found")

    para_exists = conn.execute(
        "SELECT 1 FROM corpus_paragraphs WHERE video_id = ? LIMIT 1",
        (video_id,),
    ).fetchone()
    conn.close()
    if para_exists is None:
        raise HTTPException(status_code=404, detail="No paragraphs found for this video")

    if video_id in _ckg_cancel_flags:
        raise HTTPException(
            status_code=409,
            detail="CKG extraction is already running for this video",
        )

    _ckg_cancel_flags[video_id] = False

    import threading

    thread = threading.Thread(
        target=_run_ckg_extraction,
        args=(video_id,),
        daemon=True,
    )
    thread.start()

    return JSONResponse(
        status_code=202,
        content={"status": "started", "video_id": video_id},
    )


@app.get("/api/outline/ckg")
async def get_corpus_ckg():
    """Return the CKG topology params for *every* video that has been extracted.

    One entry per ``course_ckg`` row (i.e. only videos that have actually been
    extracted appear). Each entry carries the video name (joined from the
    ``videos`` table) plus the topology / delivery-direction params, so the
    front-end corpus scatter view can plot one point per lecture with
    ``x = depth``, ``y = bottomup_ratio`` (and swap axes later without a new
    round-trip).

    Usage: single-teacher sanity check — "do Andrew Ng's lectures cluster?".
    When a second teacher is added this same payload upgrades to a
    within/between-teacher distance comparison (not implemented here — see
    issue 04 / PRD).

    ``bottomup_ratio`` may be ``null`` for a lecture with no valid prerequisite
    edges; it is returned as ``null`` (NOT coerced to 0) so the scatter view can
    skip or specially mark it rather than letting it masquerade as a real 0.
    """
    conn = get_db(str(DB_PATH))
    try:
        # 归属字段（course_id / teacher_id）是追加列：tc = teacher_course
        # （01 的 courses 表，≠ 02 的 course_topics），未分类视频时为 NULL。
        rows = conn.execute(
            "SELECT c.video_id, v.name AS name, c.model, c.created_at, "
            "c.depth, c.branch_factor, c.convergence_count, c.density, "
            "c.avg_path_length, c.clustering, c.bottomup_ratio, c.graph_json, "
            "v.course_id, tc.teacher_id "
            "FROM course_ckg c "
            "LEFT JOIN videos v ON v.id = c.video_id "
            "LEFT JOIN courses tc ON tc.id = v.course_id "
            "ORDER BY c.video_id"
        ).fetchall()
    finally:
        conn.close()

    def _counts(graph_json: str | None) -> tuple[int, int]:
        """Return (concept_count, edge_count) parsed from a graph_json blob,
        tolerating null / malformed rows (counts fall back to 0)."""
        try:
            graph = json.loads(graph_json) if graph_json else {}
        except (json.JSONDecodeError, ValueError, TypeError):
            return 0, 0
        if not isinstance(graph, dict):
            return 0, 0
        concepts = graph.get("concepts")
        # schema_version 2: decomposition_edges + relations; legacy: edges.
        decomp = graph.get("decomposition_edges")
        if decomp is None:
            decomp = graph.get("edges")
        relations = graph.get("relations")
        c = len(concepts) if isinstance(concepts, list) else 0
        e = len(decomp) if isinstance(decomp, list) else 0
        e += len(relations) if isinstance(relations, list) else 0
        return c, e

    result = []
    for row in rows:
        concept_count, edge_count = _counts(row["graph_json"])
        result.append(
            {
                "video_id": row["video_id"],
                "name": row["name"],
                "model": row["model"],
                "created_at": row["created_at"],
                "depth": row["depth"],
                "branch_factor": row["branch_factor"],
                "convergence_count": row["convergence_count"],
                "density": row["density"],
                "avg_path_length": row["avg_path_length"],
                "clustering": row["clustering"],
                "bottomup_ratio": row["bottomup_ratio"],
                "concept_count": concept_count,
                "edge_count": edge_count,
                "course_id": row["course_id"],
                "teacher_id": row["teacher_id"],
            }
        )
    return result


@app.delete("/api/outline/ckg/{video_id}")
async def delete_ckg(video_id: int):
    """Delete the stored CKG row for a video.  Idempotent: deleting a
    non-existent row still returns 200."""
    conn = get_db(str(DB_PATH))
    try:
        conn.execute("DELETE FROM course_ckg WHERE video_id = ?", (video_id,))
        conn.execute("COMMIT")
    finally:
        conn.close()
    return {"status": "deleted", "video_id": video_id}


@app.put("/api/outline/ckg/{video_id}")
async def save_ckg(video_id: int, payload: Dict[str, Any]):
    """Manually (re)store a CKG result for a video.

    Body: ``{"concepts": [...], "edges": [...]}``.  The server re-runs
    ``break_cycles`` on the edges (guaranteeing a DAG regardless of client
    input) and re-computes the topology / delivery-direction params via
    ``compute_dag_topology`` + ``compute_bottomup_ratio`` — the client cannot
    smuggle in stale params.  UPSERTs into ``course_ckg`` with the same write
    as the worker, then returns 200 + the recomputed params.

    Empty ``concepts`` → 400 (nothing meaningful to store).
    """
    concepts = payload.get("concepts") or []
    if not concepts:
        raise HTTPException(status_code=400, detail="concepts must be non-empty")

    raw_edges = payload.get("edges") or []
    edges = break_cycles(raw_edges)

    topology = compute_dag_topology(concepts, edges)
    bottomup_ratio = compute_bottomup_ratio(concepts, edges)

    graph_json = json.dumps(
        {"concepts": concepts, "edges": edges}, ensure_ascii=False
    )
    created_at = datetime.now(timezone.utc).isoformat()

    conn = get_db(str(DB_PATH))
    try:
        # Preserve an existing model label if this row was first written by the
        # worker; otherwise mark it as a manual save.
        existing = conn.execute(
            "SELECT model FROM course_ckg WHERE video_id = ?", (video_id,)
        ).fetchone()
        model = existing["model"] if existing and existing["model"] else "manual"

        conn.execute(
            "INSERT INTO course_ckg "
            "(video_id, graph_json, model, created_at, "
            "depth, branch_factor, convergence_count, density, "
            "avg_path_length, clustering, bottomup_ratio) "
            "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?) "
            "ON CONFLICT(video_id) DO UPDATE SET "
            "graph_json = excluded.graph_json, "
            "model = excluded.model, "
            "created_at = excluded.created_at, "
            "depth = excluded.depth, "
            "branch_factor = excluded.branch_factor, "
            "convergence_count = excluded.convergence_count, "
            "density = excluded.density, "
            "avg_path_length = excluded.avg_path_length, "
            "clustering = excluded.clustering, "
            "bottomup_ratio = excluded.bottomup_ratio",
            (
                video_id, graph_json, model, created_at,
                topology["depth"], topology["branch_factor"],
                topology["convergence_count"], topology["density"],
                topology["avg_path_length"], topology["clustering"],
                bottomup_ratio,
            ),
        )
        conn.execute("COMMIT")
    finally:
        conn.close()

    return {
        "video_id": video_id,
        "model": model,
        "created_at": created_at,
        "concepts": concepts,
        "edges": edges,
        "depth": topology["depth"],
        "branch_factor": topology["branch_factor"],
        "convergence_count": topology["convergence_count"],
        "density": topology["density"],
        "avg_path_length": topology["avg_path_length"],
        "clustering": topology["clustering"],
        "bottomup_ratio": bottomup_ratio,
    }


def _build_profile_descriptor(means: dict) -> dict:
    """Turn the corpus-wide mean topology into a style descriptor.

    Reuses ``_interpret_ck_style`` for the base label/description, then folds
    in a delivery-direction read (``bottomup_ratio`` → 自底向上/自顶向下) and
    emits ``generation_rules`` — a plain-language organisation instruction that
    later lesson-generation can splice into an LLM prompt.
    """
    depth = means.get("depth") or 0
    branch = means.get("branch_factor") or 0.0
    bottomup = means.get("bottomup_ratio")

    base = _interpret_ck_style(
        {
            "depth": depth,
            "branch_factor": branch,
            "avg_path_length": means.get("avg_path_length") or 0.0,
            "node_count": means.get("node_count") or 0,
        }
    )

    # Delivery direction suffix from bottomup_ratio (null → unknown, skip).
    if bottomup is None:
        direction_key = "unknown"
        direction_label = ""
        direction_phrase = ""
        gen_direction = "按概念拆解的层级顺序展开"
    elif bottomup >= 0.6:
        direction_key = "bottomup"
        direction_label = "·自底向上型"
        direction_phrase = (
            f"内容多自底向上推进（自底向上比例={bottomup:.2f}）：先铺垫下位/细节概念，再逐层上收到上位主题。"
        )
        gen_direction = "先引入下位/细节概念，再自底向上收口到上位概念"
    elif bottomup <= 0.4:
        direction_key = "topdown"
        direction_label = "·自顶向下型"
        direction_phrase = (
            f"内容多自顶向下展开（自底向上比例={bottomup:.2f}）：先给出整体框架，再逐步细化到具体概念。"
        )
        gen_direction = "先给出整体框架/高层概念，再自顶向下细化到具体概念"
    else:
        direction_key = "mixed"
        direction_label = "·双向混合型"
        direction_phrase = (
            f"内容自底向上与自顶向下交替（自底向上比例={bottomup:.2f}），框架与细节穿插推进。"
        )
        gen_direction = "框架与细节交替推进，兼顾自顶向下铺陈与自底向上推导"

    style_label = base["style_label"] + direction_label
    description = base["description"]
    if direction_phrase:
        description = description + " " + direction_phrase

    # Translate the numbers into an organisation instruction.
    if depth >= 5:
        depth_rule = f"组织成深而{'少' if branch < 3.0 else '多'}分支的先修链（典型深度约 {depth:.0f} 层）"
    else:
        depth_rule = f"组织成浅而紧凑的结构（典型深度约 {depth:.0f} 层）"
    if branch >= 3.0:
        branch_rule = f"每个主题下展开多个子话题（平均分支因子约 {branch:.1f}）"
    else:
        branch_rule = f"每个主题下分支较少、聚焦主线（平均分支因子约 {branch:.1f}）"
    generation_rules = "；".join([depth_rule, gen_direction, branch_rule]) + "。"

    return {
        "style_key": base.get("style_key", "compact"),
        "direction_key": direction_key,
        "style_label": style_label,
        "description": description,
        "generation_rules": generation_rules,
    }


def _aggregate_ckg_profile(teacher_id: int | None = None) -> dict:
    """Aggregate ``course_ckg`` into one teacher style card.

    ``teacher_id=None``（缺省）＝全库聚合（含未分类视频），与历史行为一致；
    传入教师 id 时经 videos 中转 JOIN 到 courses 过滤（INNER JOIN 天然排除
    course_id IS NULL 的未分类视频——按教师聚合时正是想要的语义）。

    Shared by ``GET /api/outline/ckg/profile`` and ``POST .../lesson-gen``:
    both need the same mean/sd + style descriptor. See the endpoint docstring
    below for the field semantics.
    """
    conn = get_db(str(DB_PATH))
    try:
        if teacher_id is None:
            rows = conn.execute(
                "SELECT graph_json, depth, branch_factor, convergence_count, "
                "relation_density, bottomup_ratio "
                "FROM course_ckg"
            ).fetchall()
        else:
            rows = conn.execute(
                "SELECT k.graph_json, k.depth, k.branch_factor, "
                "k.convergence_count, k.relation_density, k.bottomup_ratio "
                "FROM course_ckg k "
                "JOIN videos  v ON v.id = k.video_id "
                "JOIN courses c ON c.id = v.course_id "
                "WHERE c.teacher_id = ?",
                (teacher_id,),
            ).fetchall()
    finally:
        conn.close()

    lecture_count = len(rows)

    # Collect per-param value lists (skipping None so null doesn't pollute the
    # mean/sd). node_count is derived from each row's graph_json.
    # Core decomposition-style params + auxiliary bottomup (issue 03).
    # density / avg_path_length / clustering are RETIRED.
    param_keys = [
        "depth", "branch_factor", "relation_density",
        "convergence_count", "bottomup_ratio",
    ]
    collected: Dict[str, list] = {k: [] for k in param_keys}
    collected["node_count"] = []

    for row in rows:
        for k in param_keys:
            v = row[k]
            if v is not None:
                collected[k].append(v)
        try:
            graph = json.loads(row["graph_json"] or "{}")
            concepts = graph.get("concepts") if isinstance(graph, dict) else None
            collected["node_count"].append(
                len(concepts) if isinstance(concepts, list) else 0
            )
        except (json.JSONDecodeError, ValueError, TypeError):
            collected["node_count"].append(0)

    def _mean_sd(values: list) -> dict:
        if not values:
            return {"mean": None, "sd": None}
        mean = sum(values) / len(values)
        sd = statistics.pstdev(values) if len(values) > 1 else 0.0
        return {"mean": mean, "sd": sd}

    params = {k: _mean_sd(v) for k, v in collected.items()}

    if lecture_count == 0:
        return {
            "lecture_count": 0,
            "params": params,
            "style_label": "",
            "description": "",
            "generation_rules": "",
        }

    means = {k: params[k]["mean"] for k in params}
    descriptor = _build_profile_descriptor(means)

    return {
        "lecture_count": lecture_count,
        "params": params,
        "means": means,
        **descriptor,
    }


@app.get("/api/outline/ckg/profile")
async def get_ckg_profile(teacher_id: int | None = None):
    """Aggregate ``course_ckg`` into one teacher style card.

    ``?teacher_id=`` 可选：缺省＝全库聚合（向后兼容的历史行为，含未分类视频）；
    给定教师 id 时只聚合该教师名下课程的视频。当前生产库 113 个视频全属同一位
    教师，两种作用域数值一致；录入第二位教师后 within/between 对比自动成立。

    For each topology param we report ``mean`` and (population) ``sd`` across
    lectures, plus ``node_count`` (from ``len(concepts)`` per row).
    ``bottomup_ratio`` rows that are ``null`` are skipped *for that param only*
    (not coerced to 0).

    Empty corpus → ``{"lecture_count": 0, ...}`` (front-end shows empty state).
    """
    profile = _aggregate_ckg_profile(teacher_id)
    # ``means`` is an internal aggregate used by lesson-gen; not part of the
    # public profile contract — strip it so the response shape is unchanged.
    profile.pop("means", None)
    if teacher_id is not None:
        profile["teacher_id"] = teacher_id
    return profile


def _parse_num_ctx(value, default: int = 8192) -> int:
    """Parse the configured num_ctx, clamped to a sane range [1024, 131072]."""
    try:
        n = int(float(value))
    except (TypeError, ValueError):
        return default
    return max(1024, min(131072, n))


def _load_ob_llm_config() -> tuple[str, float, int]:
    """Read the Ollama model + temperature + num_ctx from config."""
    conn = get_db(str(DB_PATH))
    try:
        rows = conn.execute(
            "SELECT key, value FROM config WHERE key IN (?, ?, ?)",
            ("ob_llm_model", "ob_llm_temperature", "ob_llm_num_ctx"),
        ).fetchall()
    finally:
        conn.close()
    cfg = {row["key"]: row["value"] for row in rows}
    model = cfg.get("ob_llm_model", "qwen2.5:14b-instruct")
    try:
        temperature = float(cfg.get("ob_llm_temperature", 0.0))
    except (TypeError, ValueError):
        temperature = 0.0
    num_ctx = _parse_num_ctx(cfg.get("ob_llm_num_ctx"))
    return model, temperature, num_ctx


# Job registry for asynchronous lesson generation.
# job_id -> {"state": "running"} | {"state": "done", "result": {...}}
#          | {"state": "error", "detail": "..."}
_lesson_gen_jobs: Dict[str, dict] = {}
"""Tracks background lesson-gen jobs, keyed by hex job_id. The worker runs in
a daemon thread (so the blocking ``_call_ollama_generate`` calls don't stall
the event loop) and pushes SSE progress; the result endpoint reads this dict."""


def _split_by_chars(text: str, budget: int) -> list[str]:
    """Greedily pack lines into chunks no longer than ``budget`` chars.

    Splits on line boundaries so segments break at natural points; a single line
    longer than the budget becomes its own (over-budget) chunk rather than being
    cut mid-line. Empty chunks are dropped.
    """
    chunks: list[str] = []
    cur = ""
    for line in text.split("\n"):
        piece = (line + "\n")
        if cur and len(cur) + len(piece) > budget:
            chunks.append(cur.strip())
            cur = ""
        cur += piece
    if cur.strip():
        chunks.append(cur.strip())
    return [c for c in chunks if c]


def _chunk_material(text: str, chunk_mode: str) -> list[str]:
    """Split lesson-gen material into chunks per the chosen strategy.

    - ``"segment"``: merge page boundaries, then re-split into ~N-char segments
      (good for many tiny pages — fewer LLM calls).
    - ``"page"`` (default): one chunk per page/slide (preserves page structure).
    Both fall back to a single chunk when the text has no internal boundaries.
    """
    if chunk_mode == "segment":
        full = text.replace(_PAGE_DELIM, "\n\n")
        segs = _split_by_chars(full, _LESSON_GEN_SEGMENT_CHARS)
        return segs or [full.strip() or text]
    pages = [p.strip() for p in text.split(_PAGE_DELIM) if p.strip()]
    return pages or [text]


def _suggest_unit_count(n_concepts: int) -> int:
    """Suggested number of units for ``n_concepts`` (~_CONCEPTS_PER_UNIT each)."""
    if n_concepts <= 0:
        return 1
    return max(1, round(n_concepts / _CONCEPTS_PER_UNIT))


def _partition_concepts_in_order(concepts: list[dict], k: int) -> list[list[dict]]:
    """Split concepts into ``k`` contiguous, near-equal units by teaching order.

    Concepts are sorted by ``first_para`` (document/delivery order) and cut into
    k consecutive slices — like segmenting a long video into time-ordered units.
    The first ``n % k`` units get one extra concept so sizes differ by at most 1.
    """
    k = max(1, min(int(k or 1), len(concepts) or 1))
    ordered = sorted(concepts, key=lambda c: (c.get("first_para") if c.get("first_para") is not None else 0))
    n = len(ordered)
    base, extra = divmod(n, k)
    units: list[list[dict]] = []
    start = 0
    for i in range(k):
        size = base + (1 if i < extra else 0)
        units.append(ordered[start:start + size])
        start += size
    return [u for u in units if u]


def _lessongen_extract_concepts(input_text, mode, chunk_mode, model,
                                temperature, num_ctx):
    """Extract the concept set from material (page-by-page) or a topic.

    Returns ``(concepts, title, topic_relations)``. ``topic_relations`` is the
    relations parsed from the topic-enumeration call (material mode returns
    ``None`` — relations are extracted later per concept set). Pushes per-chunk
    SSE progress over the 35→55 band.
    """
    logger = logging.getLogger("outline")
    title = ""
    if mode == "material":
        chunks = _chunk_material(input_text, chunk_mode)
        n_chunks = len(chunks)
        unit = "段" if chunk_mode == "segment" else "页"
        merged: dict[str, dict] = {}
        order = 0
        for i, page in enumerate(chunks):
            pct = 35.0 + (i / max(1, n_chunks)) * 20.0
            push_log_event(
                "progress",
                f"教案生成 · 抽取概念拆解（第 {i + 1}/{n_chunks} {unit}）",
                progress_pct=pct,
            )
            chunk = page[:_LESSON_GEN_MAX_CHUNK_CHARS]
            try:
                c_resp = _call_ollama_generate(
                    model, _CK_PROMPT_CONCEPTS + "\n\n" + chunk,
                    stream=False, temperature=temperature,
                    fmt="json", num_predict=4096, num_ctx=num_ctx,
                )
            except Exception as e:
                logger.warning(f"lesson-gen chunk {i + 1}/{n_chunks} concept call failed: {e}")
                continue
            c_parsed = parse_concepts_response(c_resp.get("response", ""))
            if c_parsed is None:
                continue
            if not title:
                title = c_parsed.get("title", "")
            for c in c_parsed["concepts"]:
                key = str(c.get("name", "")).strip().lower()
                if not key or key in merged:
                    continue
                c["first_para"] = order
                order += 1
                merged[key] = c
        return list(merged.values()), title, None

    # topic mode: one enumeration call yields concepts + relations together.
    push_log_event("progress", "教案生成 · 枚举主题知识点", progress_pct=45.0)
    t_resp = _call_ollama_generate(
        model, _LESSON_GEN_TOPIC_PROMPT + " " + input_text,
        stream=False, temperature=temperature, fmt="json",
        num_predict=4096, num_ctx=num_ctx,
    )
    raw = t_resp.get("response", "")
    c_parsed = parse_concepts_response(raw)
    if c_parsed is None:
        raise RuntimeError("主题概念枚举失败：LLM 输出无法解析。")
    r_parsed = parse_relations_response(raw)
    topic_relations = r_parsed["relations"] if r_parsed else []
    return c_parsed["concepts"], c_parsed.get("title", ""), topic_relations


def _lessongen_build_plan(concepts, relations, means, generation_rules,
                          style_label, model, temperature, num_ctx,
                          unit_title="", rel_progress=None, gen_progress=None):
    """Build one styled lesson plan from a concept set.

    Extracts grounded relations (if ``relations`` is None), derives the
    decomposition backbone, then runs the style-aware generation. Returns the
    result dict stored under a job's ``result``. ``rel_progress`` / ``gen_progress``
    are optional ``(pct, message)`` tuples pushed as SSE progress before the
    relations / generation steps respectively.
    """
    if rel_progress:
        push_log_event("progress", rel_progress[1], progress_pct=rel_progress[0])
    if relations is None:
        concept_lines = "\n".join(
            f"- {c['name']}: {c.get('definition', '')}" for c in concepts
        )
        r_resp = _call_ollama_generate(
            model, _CK_PROMPT_RELATIONS + "\n" + concept_lines,
            stream=False, temperature=temperature, fmt="json",
            num_predict=4096, num_ctx=num_ctx,
        )
        r_parsed = parse_relations_response(r_resp.get("response", ""))
        relations = r_parsed["relations"] if r_parsed else []

    decomposition_edges = derive_decomposition_edges(concepts)
    relations = validate_relations(relations, concepts)

    bottomup_mean = means.get("bottomup_ratio")
    if bottomup_mean is None:
        sequencing = "follow the decomposition order"
    elif bottomup_mean >= 0.6:
        sequencing = ("bottom-up: introduce sub-concepts/details first, then "
                      "roll up to the umbrella concept")
    elif bottomup_mean <= 0.4:
        sequencing = ("top-down: name the umbrella concept first, then carve "
                      "it into sub-concepts")
    else:
        sequencing = "mixed"

    concept_lines = "\n".join(
        f"- {c['name']}: {c.get('definition', '')}" for c in concepts
    )
    edge_lines = "\n".join(
        f"{e['from']} -> {e['to']}" for e in decomposition_edges
    ) or "(none)"
    relation_lines = "\n".join(
        f"{r['from']} -[{r.get('type', '') or 'related'}]-> {r['to']}"
        for r in relations
    ) or "(none)"

    gen_prompt = _LESSON_GEN_OUTLINE_PROMPT.format(
        generation_rules=generation_rules or "(no explicit style rules)",
        depth_mean=f"{(means.get('depth') or 0):.0f}",
        branch_mean=f"{(means.get('branch_factor') or 0.0):.1f}",
        sequencing=sequencing,
        concept_lines=concept_lines,
        edge_lines=edge_lines,
        relation_lines=relation_lines,
    )

    if gen_progress:
        push_log_event("progress", gen_progress[1], progress_pct=gen_progress[0])
    outline: list = []
    sequence: list = []
    try:
        g_resp = _call_ollama_generate(
            model, gen_prompt, stream=False, temperature=temperature,
            fmt="json", num_predict=4096, num_ctx=num_ctx,
        )
        parsed = parse_lesson_outline_response(g_resp.get("response", ""))
        if parsed is not None:
            outline = parsed["outline"]
            sequence = parsed["sequence"]
    except Exception as e:
        logging.getLogger("outline").warning(f"lesson-gen generation failed: {e}")

    if not sequence:
        sequence = [c["name"] for c in concepts]
    if not outline:
        outline = [{"name": name, "children": []} for name in sequence]

    return {
        "style_label": style_label,
        "title": unit_title,
        "outline": outline,
        "sequence": sequence,
        "concepts": concepts,
        "decomposition_edges": decomposition_edges,
        "relations": relations,
        "edges": decomposition_edges,
    }


def _run_lesson_gen(job_id: str, input_text: str, mode: str,
                    chunk_mode: str = "page") -> None:
    """Background worker: generate a styled teaching outline + push SSE progress.

    ``chunk_mode`` ("page" | "segment") controls how uploaded material is split
    for page-by-page vs segment-by-segment concept extraction.

    Steps (each pushes a ``progress`` SSE event prefixed ``教案生成 · ``):
      - read teacher style profile      → 10
      - material: extract concept decomposition (35), grounded relations (60)
        topic:    enumerate decomposition + relations (45)
      - style-aware outline generation  → 80
      - done                            → 100 (``success`` event)

    On hard failure the job is marked ``error`` and an ``error`` event is
    pushed. A bad/unparseable *generation* output is NOT a failure — it
    degrades to a flat outline (matching the previous synchronous behaviour).
    """
    logger = logging.getLogger("outline")
    try:
        # --- Step: teacher style profile ---
        push_log_event("progress", "教案生成 · 读取教师风格画像", progress_pct=10.0)
        profile = _aggregate_ckg_profile()
        means = profile.get("means", {})
        generation_rules = profile.get("generation_rules", "")
        style_label = profile.get("style_label", "")

        model, temperature, num_ctx = _load_ob_llm_config()

        # --- Step: concepts (decomposition) ---
        concepts, title, topic_relations = _lessongen_extract_concepts(
            input_text, mode, chunk_mode, model, temperature, num_ctx
        )
        if not concepts:
            raise RuntimeError("概念拆解失败：未能从资料中解析出概念。")

        # --- Step: relations + style-aware generation (single plan) ---
        rel_progress = (60.0, "教案生成 · 抽取拆解+讲述关联") if mode == "material" else None
        plan = _lessongen_build_plan(
            concepts, topic_relations, means, generation_rules, style_label,
            model, temperature, num_ctx, unit_title=title,
            rel_progress=rel_progress,
            gen_progress=(80.0, "教案生成 · 风格化生成教案"),
        )

        _lesson_gen_jobs[job_id] = {"state": "done", "result": plan}
        push_log_event("success", "教案生成 · 完成", progress_pct=100.0)
        logger.info(f"lesson-gen job {job_id} done (mode={mode})")
    except Exception as e:
        _lesson_gen_jobs[job_id] = {"state": "error", "detail": str(e)}
        push_log_event(
            "error",
            f"教案生成 · 失败：{e}",
            traceback=traceback.format_exc(),
        )
        logger.error(f"lesson-gen job {job_id} failed: {e}", exc_info=True)


def _run_lesson_gen_analyze(job_id: str, input_text: str,
                            chunk_mode: str = "page") -> None:
    """Phase-1 worker: extract + merge concepts from material (page-by-page),
    then report the count + a suggested unit split. No generation yet — the
    frontend decides how many units to split into before phase 2."""
    logger = logging.getLogger("outline")
    try:
        model, temperature, num_ctx = _load_ob_llm_config()
        concepts, title, _ = _lessongen_extract_concepts(
            input_text, "material", chunk_mode, model, temperature, num_ctx
        )
        if not concepts:
            raise RuntimeError("概念拆解失败：未能从资料中解析出概念。")
        n = len(concepts)
        _lesson_gen_jobs[job_id] = {
            "state": "done",
            "result": {
                "phase": "analyze",
                "concepts": concepts,
                "title": title,
                "concept_count": n,
                "needs_split": n > _UNIT_SPLIT_THRESHOLD,
                "suggested_units": _suggest_unit_count(n),
                "threshold": _UNIT_SPLIT_THRESHOLD,
            },
        }
        push_log_event("success", f"教案生成 · 分析完成（{n} 个概念）", progress_pct=58.0)
        logger.info(f"lesson-gen analyze job {job_id} done ({n} concepts)")
    except Exception as e:
        _lesson_gen_jobs[job_id] = {"state": "error", "detail": str(e)}
        push_log_event("error", f"教案生成 · 失败：{e}", traceback=traceback.format_exc())
        logger.error(f"lesson-gen analyze job {job_id} failed: {e}", exc_info=True)


def _run_lesson_gen_units(job_id: str, concepts: list, title: str,
                          units: int) -> None:
    """Phase-2 worker: split the concept set into ``units`` in-order units, then
    build a styled lesson plan per unit. Stores ``{"units": [plan, ...]}``."""
    logger = logging.getLogger("outline")
    try:
        push_log_event("progress", "教案生成 · 读取教师风格画像", progress_pct=58.0)
        profile = _aggregate_ckg_profile()
        means = profile.get("means", {})
        generation_rules = profile.get("generation_rules", "")
        style_label = profile.get("style_label", "")
        model, temperature, num_ctx = _load_ob_llm_config()

        parts = _partition_concepts_in_order(concepts, units)
        k = len(parts)
        plans = []
        for i, part in enumerate(parts):
            base = 60.0 + (i / max(1, k)) * 38.0  # 60→98 across units
            push_log_event(
                "progress", f"教案生成 · 生成单元 {i + 1}/{k}", progress_pct=base,
            )
            # Don't bake a language-specific label ("单元 N") into the title —
            # the front-end localizes the tab/heading from ``unit_index`` so it
            # adapts to zh/en/ja. Split units are positional, so title stays "".
            plan = _lessongen_build_plan(
                part, None, means, generation_rules, style_label,
                model, temperature, num_ctx, unit_title="",
            )
            plan["unit_index"] = i + 1
            plans.append(plan)

        _lesson_gen_jobs[job_id] = {
            "state": "done",
            "result": {"phase": "units", "title": title, "unit_count": k, "units": plans},
        }
        push_log_event("success", "教案生成 · 完成", progress_pct=100.0)
        logger.info(f"lesson-gen units job {job_id} done ({k} units)")
    except Exception as e:
        _lesson_gen_jobs[job_id] = {"state": "error", "detail": str(e)}
        push_log_event("error", f"教案生成 · 失败：{e}", traceback=traceback.format_exc())
        logger.error(f"lesson-gen units job {job_id} failed: {e}", exc_info=True)


@app.post("/api/outline/lesson-gen")
async def generate_lesson(request: Request):
    """Start an asynchronous styled teaching-outline generation job.

    body: ``{"input_text": "...", "mode": "material" | "topic"}``

    Validates the request (non-empty input, valid mode, non-empty teacher
    profile), then launches ``_run_lesson_gen`` in a daemon thread and returns
    ``202 {"job_id": ...}`` immediately. Progress is broadcast over the shared
    SSE log stream (events prefixed ``教案生成 · ``); the final result is read
    from ``GET /api/outline/lesson-gen/result/{job_id}``.

    A 400 is returned (and no job started) for empty input, bad mode, or an
    empty corpus profile.
    """
    body = await request.json()
    input_text = str(body.get("input_text", "")).strip()
    mode = str(body.get("mode", "material")).strip() or "material"
    chunk_mode = str(body.get("chunk_mode", "page")).strip() or "page"
    if not input_text:
        raise HTTPException(status_code=400, detail="input_text 不能为空")
    if mode not in ("material", "topic"):
        raise HTTPException(status_code=400, detail="mode 必须是 material 或 topic")
    if chunk_mode not in ("page", "segment"):
        raise HTTPException(status_code=400, detail="chunk_mode 必须是 page 或 segment")

    # Cheap fail-fast guard before starting any job/thread.
    profile = _aggregate_ckg_profile()
    if not profile.get("lecture_count"):
        raise HTTPException(
            status_code=400,
            detail="尚无已分析课程，请先在知识图谱页分析视频生成风格画像。",
        )

    job_id = uuid.uuid4().hex
    _lesson_gen_jobs[job_id] = {"state": "running"}

    thread = threading.Thread(
        target=_run_lesson_gen,
        args=(job_id, input_text, mode, chunk_mode),
        daemon=True,
    )
    thread.start()

    return JSONResponse(status_code=202, content={"job_id": job_id})


def _lesson_gen_profile_guard():
    """Raise 400 if no teacher profile exists yet (shared by lesson-gen endpoints)."""
    profile = _aggregate_ckg_profile()
    if not profile.get("lecture_count"):
        raise HTTPException(
            status_code=400,
            detail="尚无已分析课程，请先在知识图谱页分析视频生成风格画像。",
        )


@app.post("/api/outline/lesson-gen/analyze")
async def lesson_gen_analyze(request: Request):
    """Phase 1 (material only): extract concepts and report a suggested unit
    split. Returns ``202 {job_id}``; the result carries ``{concepts, title,
    concept_count, needs_split, suggested_units}``."""
    body = await request.json()
    input_text = str(body.get("input_text", "")).strip()
    chunk_mode = str(body.get("chunk_mode", "page")).strip() or "page"
    if not input_text:
        raise HTTPException(status_code=400, detail="input_text 不能为空")
    if chunk_mode not in ("page", "segment"):
        raise HTTPException(status_code=400, detail="chunk_mode 必须是 page 或 segment")
    _lesson_gen_profile_guard()

    job_id = uuid.uuid4().hex
    _lesson_gen_jobs[job_id] = {"state": "running"}
    threading.Thread(
        target=_run_lesson_gen_analyze,
        args=(job_id, input_text, chunk_mode),
        daemon=True,
    ).start()
    return JSONResponse(status_code=202, content={"job_id": job_id})


@app.post("/api/outline/lesson-gen/generate-units")
async def lesson_gen_generate_units(request: Request):
    """Phase 2: split the (phase-1) concept set into ``units`` in-order units and
    generate a styled lesson plan per unit. Returns ``202 {job_id}``; the result
    carries ``{units: [...], unit_count}``."""
    body = await request.json()
    concepts = body.get("concepts") or []
    title = str(body.get("title", "")).strip()
    try:
        units = int(body.get("units", 1) or 1)
    except (TypeError, ValueError):
        units = 1
    if not isinstance(concepts, list) or not concepts:
        raise HTTPException(status_code=400, detail="concepts 不能为空")
    units = max(1, min(units, len(concepts)))
    _lesson_gen_profile_guard()

    job_id = uuid.uuid4().hex
    _lesson_gen_jobs[job_id] = {"state": "running"}
    threading.Thread(
        target=_run_lesson_gen_units,
        args=(job_id, concepts, title, units),
        daemon=True,
    ).start()
    return JSONResponse(status_code=202, content={"job_id": job_id})


@app.get("/api/outline/lesson-gen/result/{job_id}")
async def get_lesson_gen_result(job_id: str):
    """Return the current status/result of a lesson-gen job.

    ``done`` → ``{"state": "done", "result": {...}}``; ``running`` →
    ``{"state": "running"}``; ``error`` → ``{"state": "error", "detail": ...}``.
    Unknown job_id → 404.
    """
    job = _lesson_gen_jobs.get(job_id)
    if job is None:
        raise HTTPException(status_code=404, detail="未知 job_id")
    return job


# ---------------------------------------------------------------------------
# Saved lesson plans (教案) — CRUD on the lesson_plans table
# ---------------------------------------------------------------------------


@app.get("/api/outline/lesson-plans")
async def list_lesson_plans():
    """List saved lesson plans (id / name / created_at), newest first."""
    conn = get_db(str(DB_PATH))
    try:
        rows = conn.execute(
            "SELECT id, name, created_at FROM lesson_plans ORDER BY id DESC"
        ).fetchall()
    finally:
        conn.close()
    return [
        {"id": r["id"], "name": r["name"], "created_at": r["created_at"]}
        for r in rows
    ]


@app.post("/api/outline/lesson-plans")
async def save_lesson_plan(request: Request):
    """Save a lesson plan: body ``{name, payload}``. Returns ``{id}``."""
    body = await request.json()
    name = str(body.get("name", "")).strip()
    payload = body.get("payload")
    if not name:
        raise HTTPException(status_code=400, detail="教案名称不能为空")
    if payload is None:
        raise HTTPException(status_code=400, detail="payload 不能为空")
    created_at = datetime.now(timezone.utc).isoformat()
    payload_json = json.dumps(payload, ensure_ascii=False)
    conn = get_db(str(DB_PATH))
    try:
        cur = conn.execute(
            "INSERT INTO lesson_plans (name, created_at, payload_json) "
            "VALUES (?, ?, ?)",
            (name, created_at, payload_json),
        )
        conn.execute("COMMIT")
        new_id = cur.lastrowid
    finally:
        conn.close()
    return {"id": new_id, "name": name, "created_at": created_at}


@app.get("/api/outline/lesson-plans/{plan_id}")
async def get_lesson_plan(plan_id: int):
    """Return one saved plan's full payload (``{id, name, created_at, payload}``)."""
    conn = get_db(str(DB_PATH))
    try:
        row = conn.execute(
            "SELECT id, name, created_at, payload_json FROM lesson_plans WHERE id = ?",
            (plan_id,),
        ).fetchone()
    finally:
        conn.close()
    if row is None:
        raise HTTPException(status_code=404, detail="教案不存在")
    try:
        payload = json.loads(row["payload_json"] or "{}")
    except (json.JSONDecodeError, ValueError):
        payload = {}
    return {"id": row["id"], "name": row["name"],
            "created_at": row["created_at"], "payload": payload}


@app.put("/api/outline/lesson-plans/{plan_id}")
async def rename_lesson_plan(plan_id: int, request: Request):
    """Rename a saved plan: body ``{name}``."""
    body = await request.json()
    name = str(body.get("name", "")).strip()
    if not name:
        raise HTTPException(status_code=400, detail="教案名称不能为空")
    conn = get_db(str(DB_PATH))
    try:
        cur = conn.execute(
            "UPDATE lesson_plans SET name = ? WHERE id = ?", (name, plan_id)
        )
        conn.execute("COMMIT")
    finally:
        conn.close()
    if cur.rowcount == 0:
        raise HTTPException(status_code=404, detail="教案不存在")
    return {"id": plan_id, "name": name}


@app.delete("/api/outline/lesson-plans/{plan_id}")
async def delete_lesson_plan(plan_id: int):
    """Delete a saved plan (idempotent — deleting a missing id still returns ok)."""
    conn = get_db(str(DB_PATH))
    try:
        conn.execute("DELETE FROM lesson_plans WHERE id = ?", (plan_id,))
        conn.execute("COMMIT")
    finally:
        conn.close()
    return {"status": "deleted", "id": plan_id}


# ---------------------------------------------------------------------------
# Lesson-gen — file upload → plain text (feeds the material flow)
# ---------------------------------------------------------------------------

_SUPPORTED_UPLOAD_EXTS = (".pdf", ".pptx", ".docx", ".txt", ".md")
_UNSUPPORTED_MSG = "不支持的文件类型，仅支持 pdf/pptx/docx/txt/md"

# Page/slide boundary marker. PDF pages and PPTX slides are joined with this so
# the lesson-gen worker can split the material back into per-page chunks for
# page-by-page concept extraction (avoiding context-window truncation).
_PAGE_DELIM = "\n\n<<<PAGE>>>\n\n"


def _decode_text_bytes(data: bytes) -> str:
    """Decode raw text bytes, tolerating non-UTF-8 input."""
    try:
        return data.decode("utf-8")
    except UnicodeDecodeError:
        try:
            return data.decode("latin-1")
        except Exception:
            return data.decode("utf-8", errors="ignore")


def _extract_pdf(data: bytes) -> str:
    import io
    import pypdf

    reader = pypdf.PdfReader(io.BytesIO(data))
    parts = []
    for page in reader.pages:
        try:
            parts.append(page.extract_text() or "")
        except Exception:
            parts.append("")
    # Join with the page delimiter so the worker can do per-page extraction.
    return _PAGE_DELIM.join(parts)


def _extract_pptx(data: bytes) -> str:
    import io
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    slides = []
    for slide in prs.slides:
        parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                parts.append(shape.text)
            elif getattr(shape, "text", None):
                parts.append(shape.text)
        slides.append("\n".join(p for p in parts if p))
    # One chunk per slide (the PPTX analog of a PDF page).
    return _PAGE_DELIM.join(slides)


def _extract_docx(data: bytes) -> str:
    import io
    import docx

    doc = docx.Document(io.BytesIO(data))
    return "\n".join(p.text for p in doc.paragraphs)


def _extract_text_from_upload(filename: str, data: bytes) -> str:
    """Extract plain text from an uploaded file by extension.

    Parsing libraries are imported inside the per-type helpers so a missing
    optional dependency only affects that file type, not the whole app.

    Raises ``ValueError`` for unsupported extensions (the endpoint maps it to
    a 400). Parse errors propagate as the underlying exception (the endpoint
    maps them to 422).
    """
    name = (filename or "").lower()
    ext = "." + name.rsplit(".", 1)[-1] if "." in name else ""
    if ext in (".txt", ".md"):
        return _decode_text_bytes(data)
    if ext == ".pdf":
        return _extract_pdf(data)
    if ext == ".pptx":
        return _extract_pptx(data)
    if ext == ".docx":
        return _extract_docx(data)
    raise ValueError(_UNSUPPORTED_MSG)


@app.post("/api/outline/lesson-gen/extract-file")
async def lesson_gen_extract_file(file: UploadFile = File(...)):
    """Extract plain text from an uploaded PDF/PPTX/DOCX/TXT/MD file.

    Returns ``{"filename", "text", "char_count"}``. The extracted text is
    meant to be dropped into the lesson-gen *material* input by the frontend;
    no generation happens here.

    Errors degrade gracefully: unsupported extension → 400; missing parsing
    library → 400 with a friendly hint; corrupt file / parse failure → 422.
    """
    filename = file.filename or ""
    data = await file.read()
    try:
        text = _extract_text_from_upload(filename, data)
    except ValueError as e:
        raise HTTPException(status_code=400, detail=str(e))
    except ImportError as e:
        raise HTTPException(
            status_code=400,
            detail=f"服务器缺少解析该文件类型所需的库：{e}",
        )
    except Exception as e:
        raise HTTPException(
            status_code=422,
            detail=f"文件解析失败（可能已损坏或格式不符）：{e}",
        )
    return {"filename": filename, "text": text, "char_count": len(text)}


@app.get("/api/outline/ckg/{video_id}")
async def get_ckg(video_id: int):
    """Return the stored CKG concepts for a video.

    Returns ``{"video_id", "model", "created_at", "concepts": [...]}``.
    Returns 404 if no CKG has been extracted for the video yet.
    """
    conn = get_db(str(DB_PATH))
    try:
        row = conn.execute(
            "SELECT video_id, graph_json, model, created_at, "
            "depth, branch_factor, convergence_count, relation_density, "
            "bottomup_ratio "
            "FROM course_ckg WHERE video_id = ?",
            (video_id,),
        ).fetchone()
    finally:
        conn.close()

    if row is None:
        raise HTTPException(
            status_code=404,
            detail="No CKG found — run extraction first",
        )

    try:
        graph = json.loads(row["graph_json"] or "{}")
    except (json.JSONDecodeError, ValueError):
        graph = {}

    # New (schema_version 2): decomposition backbone + grounded relations.
    # Legacy rows only had "edges" (prerequisite); fall back to it so old data
    # still renders. ``edges`` is kept as a back-compat alias of the
    # decomposition backbone for existing consumers.
    decomposition_edges = graph.get("decomposition_edges")
    if decomposition_edges is None:
        decomposition_edges = graph.get("edges", [])
    relations = graph.get("relations", [])

    return {
        "video_id": row["video_id"],
        "model": row["model"],
        "created_at": row["created_at"],
        "schema_version": graph.get("schema_version", 1),
        "title": graph.get("title", ""),
        "concepts": graph.get("concepts", []),
        "decomposition_edges": decomposition_edges,
        "relations": relations,
        "edges": decomposition_edges,
        # Decomposition-style parameters (issue 03). depth = 拆解深度,
        # branch_factor = 拆解宽度, relation_density = 关联密度,
        # convergence_count = 汇聚数. bottomup_ratio is an AUXILIARY signal.
        # density / avg_path_length / clustering are RETIRED.
        "depth": row["depth"],
        "branch_factor": row["branch_factor"],
        "convergence_count": row["convergence_count"],
        "relation_density": row["relation_density"],
        "bottomup_ratio": row["bottomup_ratio"],
    }


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------

if __name__ == "__main__":
    import uvicorn

    uvicorn.run("main_outline:app", host="0.0.0.0", port=8001)
