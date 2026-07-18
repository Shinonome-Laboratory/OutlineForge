"""Tests for outline module.

Covers: project skeleton (Slice 01), config panel API (Slice 11),
and future slices.
"""

import asyncio
import json
import logging
import sqlite3
import sys
from pathlib import Path
from unittest.mock import patch, MagicMock

import pytest

# Ensure the project root is on sys.path so we can import main_outline
PROJECT_ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(PROJECT_ROOT))

# ---------------------------------------------------------------------------
# Paths
# ---------------------------------------------------------------------------

DB_PATH = Path(r"d:\Project\All for Style\00-data\corpus.db")


# ---------------------------------------------------------------------------
# Fixtures
# ---------------------------------------------------------------------------

@pytest.fixture
def client():
    """Return a TestClient for the FastAPI app defined in main_outline.py."""
    from main_outline import app  # noqa: F811
    from fastapi.testclient import TestClient

    with TestClient(app) as tc:
        yield tc


# ---------------------------------------------------------------------------
# Test: FastAPI app startup
# ---------------------------------------------------------------------------

def test_app_starts(client):
    """The FastAPI application can be instantiated and responds to requests."""
    resp = client.get("/")
    assert resp.status_code == 200


# ---------------------------------------------------------------------------
# Test: GET / returns outline.html (Content-Type text/html)
# ---------------------------------------------------------------------------

def test_root_returns_html_content_type(client):
    """The root endpoint returns Content-Type text/html."""
    resp = client.get("/")
    assert "text/html" in resp.headers.get("content-type", "")


def test_root_body_contains_html_doctype(client):
    """The response body starts with an HTML doctype declaration."""
    resp = client.get("/")
    assert "<!DOCTYPE html>" in resp.text or "<!doctype html>" in resp.text.lower()


# ---------------------------------------------------------------------------
# Test: course_topics table created with correct columns
# ---------------------------------------------------------------------------

EXPECTED_COURSE_TOPICS_COLUMNS = {
    "id",
    "video_id",
    "start_para_index",
    "end_para_index",
    "start_time",
    "end_time",
    "topic_name",
    "subtree_json",
}


def test_course_topics_table_exists():
    """The course_topics table exists in corpus.db after init_db runs."""
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row

    tables = {
        row["name"]
        for row in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table'"
        ).fetchall()
    }
    conn.close()

    assert "course_topics" in tables, "course_topics table missing from corpus.db"


def test_course_topics_columns_correct():
    """The course_topics table has exactly the expected columns."""
    conn = sqlite3.connect(str(DB_PATH))
    cursor = conn.execute("PRAGMA table_info(course_topics)")
    actual_columns = {row[1] for row in cursor.fetchall()}
    conn.close()

    assert actual_columns == EXPECTED_COURSE_TOPICS_COLUMNS, (
        f"course_topics columns mismatch:\n"
        f"  expected: {sorted(EXPECTED_COURSE_TOPICS_COLUMNS)}\n"
        f"  got:      {sorted(actual_columns)}"
    )


# ---------------------------------------------------------------------------
# Test: 5 ob_ config keys exist with correct defaults
# ---------------------------------------------------------------------------

EXPECTED_OB_CONFIG = {
    "ob_llm_model": "qwen2.5:14b-instruct",
    "ob_llm_temperature": "0.0",
    "ob_max_topics": "auto",
    "ob_prompt_round1": (
        "You are a classroom instruction content analyst. Below is the complete "
        "transcript of a lecture, organized by paragraph index.\n"
        "\n"
        "Analyze the teaching content and complete the following tasks:\n"
        "\n"
        "1. Generate a concise course title (no more than 10 words).\n"
        "2. Segment the full text into topic sections. Each topic represents a "
        "semantically coherent unit of instruction. Segmentation principles:\n"
        "   - Split when the teaching subject clearly shifts\n"
        "   - No hard limit on total topic count; let the content naturally determine it\n"
        "\n"
        "Output strictly in the following JSON format. Do not output any other text:\n"
        "\n"
        '{\n'
        '  "course_name": "Course Title",\n'
        '  "topics": [\n'
        '    { "name": "Topic Name", "start_para": start_paragraph_index, '
        '"end_para": end_paragraph_index }\n'
        '  ]\n'
        '}\n'
        "\n"
        "Full transcript below:"
    ),
    "ob_prompt_round2": (
        'You are a subject knowledge structure analyst. Below is the teaching '
        'content for the topic "{topic_name}" from a lecture transcript.\n'
        "\n"
        "Analyze the knowledge structure within this text and generate a knowledge "
        "subtree for this topic. Let the depth be determined by the actual content "
        "structure — shallow for simple content, deep for complex content.\n"
        "\n"
        "Output strictly in the following JSON format. Do not output any other text:\n"
        "\n"
        '{\n'
        '  "topic": "{topic_name}",\n'
        '  "subtree": {\n'
        '    "children": [\n'
        '      { "name": "sub-topic name", "children": [...] }\n'
        "    ]\n"
        "  }\n"
        "}\n"
        "\n"
        "If the topic content is too simple to subdivide further, children can be "
        "an empty array.\n"
        "\n"
        "Topic text below:"
    ),
}


def test_ob_config_keys_exist():
    """All seeded ob_ config keys exist in the config table, non-empty.

    Keys are derived from ``OB_CONFIG_DEFAULTS`` (single source of truth) so
    the test can't rot when a key is added. Exact values are NOT asserted:
    every ob_ key is user-editable via the config panel, and the prompts are
    additionally auto-migrated by ``init_db`` — pinning values makes the test
    fail as soon as the user tweaks any setting.
    """
    from main_outline import OB_CONFIG_DEFAULTS

    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row
    rows = conn.execute(
        "SELECT key, value FROM config WHERE key LIKE 'ob_%'"
    ).fetchall()
    conn.close()

    actual = {row["key"]: row["value"] for row in rows}

    for key, _default in OB_CONFIG_DEFAULTS:
        assert key in actual, (
            f"Missing config key: {key}\n  got: {sorted(actual.keys())}"
        )
        assert (actual[key] or "").strip(), f"Empty config value: {key}"


# ---------------------------------------------------------------------------
# Test: init_db is idempotent
# ---------------------------------------------------------------------------

def test_init_db_idempotent():
    """Running init_db twice does not error (idempotent via IF NOT EXISTS
    and INSERT OR IGNORE)."""
    from main_outline import init_db

    # First call
    init_db(str(DB_PATH))
    # Second call — must not raise
    init_db(str(DB_PATH))

    # Verify data is still intact
    conn = sqlite3.connect(str(DB_PATH))
    conn.row_factory = sqlite3.Row

    # Check course_topics table still exists
    tables = {
        row["name"]
        for row in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table'"
        ).fetchall()
    }
    assert "course_topics" in tables

    # Check ob_ config keys still have correct count (6, not duplicated):
    # model / temperature / num_ctx / prompt_round1 / prompt_round2 ... and any
    # other ob_ defaults. Adding ob_llm_num_ctx raised this from 5 to 6.
    count = conn.execute(
        "SELECT COUNT(*) FROM config WHERE key LIKE 'ob_%'"
    ).fetchone()[0]
    conn.close()

    assert count == 6, f"Expected 6 ob_ config keys after idempotent init, got {count}"


# ============================================================================
# Slice 11 — Config panel API
# ============================================================================

# ---------------------------------------------------------------------------
# GET /api/outline/config
# ---------------------------------------------------------------------------


def test_get_config_returns_all_ob_keys(client):
    """GET /api/outline/config returns a JSON object with all ob_ prefix keys."""
    resp = client.get("/api/outline/config")
    assert resp.status_code == 200
    data = resp.json()
    assert isinstance(data, dict)
    # All keys should start with ob_
    for key in data:
        assert key.startswith("ob_"), f"Non-ob_ key returned: {key}"
    # The 5 default keys must be present
    for expected_key in EXPECTED_OB_CONFIG:
        assert expected_key in data, f"Missing config key: {expected_key}"


# ---------------------------------------------------------------------------
# PUT /api/outline/config
# ---------------------------------------------------------------------------


def test_put_config_updates_existing_key(client):
    """PUT /api/outline/config updates an existing ob_ config key."""
    new_temp = "0.7"
    resp = client.put(
        "/api/outline/config",
        json={"ob_llm_temperature": new_temp},
    )
    assert resp.status_code == 200
    data = resp.json()
    assert data.get("ob_llm_temperature") == new_temp

    # Verify persistence via GET
    resp2 = client.get("/api/outline/config")
    assert resp2.json()["ob_llm_temperature"] == new_temp

    # Restore default value (use hardcoded default, not current DB value)
    client.put("/api/outline/config",
               json={"ob_llm_temperature": EXPECTED_OB_CONFIG["ob_llm_temperature"]})


def test_put_config_upserts_multiple_keys(client):
    """PUT /api/outline/config can update multiple keys in one call."""
    payload = {
        "ob_llm_temperature": "0.3",
        "ob_max_topics": "5",
    }
    resp = client.put("/api/outline/config", json=payload)
    assert resp.status_code == 200
    data = resp.json()
    assert data["ob_llm_temperature"] == "0.3"
    assert data["ob_max_topics"] == "5"

    # Verify via GET
    resp2 = client.get("/api/outline/config")
    assert resp2.json()["ob_llm_temperature"] == "0.3"
    assert resp2.json()["ob_max_topics"] == "5"

    # Restore default values (use hardcoded defaults, not current DB values)
    client.put("/api/outline/config", json={
        "ob_llm_temperature": EXPECTED_OB_CONFIG["ob_llm_temperature"],
        "ob_max_topics": EXPECTED_OB_CONFIG["ob_max_topics"],
    })


def test_put_config_rejects_non_ob_prefix_key(client):
    """PUT /api/outline/config rejects keys not starting with ob_."""
    resp = client.put(
        "/api/outline/config",
        json={"malicious_key": "evil"},
    )
    assert resp.status_code == 422
    detail = resp.json().get("detail", "")
    assert "ob_" in str(detail).lower() or "prefix" in str(detail).lower()


def test_put_config_rejects_mixed_keys(client):
    """PUT /api/outline/config rejects the whole batch if any key lacks ob_ prefix."""
    payload = {
        "ob_llm_temperature": "0.5",
        "not_ob_key": "bad",
    }
    resp = client.put("/api/outline/config", json=payload)
    assert resp.status_code == 422


def test_put_config_empty_body_returns_400(client):
    """PUT /api/outline/config with empty JSON body returns 400."""
    resp = client.put("/api/outline/config", json={})
    assert resp.status_code == 400


# ---------------------------------------------------------------------------
# GET /api/outline/ollama/models
# ---------------------------------------------------------------------------


def test_get_ollama_models_returns_list(client):
    """GET /api/outline/ollama/models returns a list of model name strings."""
    mock_models = [
        {"name": "qwen2.5:14b-instruct", "modified_at": "2025-01-01T00:00:00Z"},
        {"name": "llama3:8b", "modified_at": "2025-01-02T00:00:00Z"},
    ]

    with patch("main_outline._fetch_ollama_tags", return_value=mock_models):
        resp = client.get("/api/outline/ollama/models")
        assert resp.status_code == 200
        data = resp.json()
        assert isinstance(data, list)
        assert len(data) == 2
        assert "qwen2.5:14b-instruct" in data
        assert "llama3:8b" in data


def test_get_ollama_models_handles_ollama_unavailable(client):
    """GET /api/outline/ollama/models returns empty list + 503 when Ollama is down."""
    import httpx

    with patch("main_outline._fetch_ollama_tags",
               side_effect=httpx.ConnectError("refused")):
        resp = client.get("/api/outline/ollama/models")
        assert resp.status_code == 503
        data = resp.json()
        assert isinstance(data, list)
        assert data == []


def test_get_ollama_models_extracts_name_only(client):
    """Model list only contains 'name' field stripped from each model entry."""
    mock_models = [
        {"name": "mistral:7b", "modified_at": "x", "size": 123},
        {"name": "gemma2:9b", "modified_at": "y", "details": {}},
    ]

    with patch("main_outline._fetch_ollama_tags", return_value=mock_models):
        resp = client.get("/api/outline/ollama/models")
        assert resp.status_code == 200
        data = resp.json()
        # Must be flat list of name strings
        assert all(isinstance(m, str) for m in data)
        assert data == ["mistral:7b", "gemma2:9b"]


# ============================================================================
# Slice 04 — 视频列表 + 段落浏览（后端 API）
# ============================================================================

# ---------------------------------------------------------------------------
# GET /api/outline/videos
# ---------------------------------------------------------------------------


def test_get_videos_returns_json_array(client):
    """GET /api/outline/videos returns a JSON array."""
    resp = client.get("/api/outline/videos")
    assert resp.status_code == 200, f"Expected 200, got {resp.status_code}"
    data = resp.json()
    assert isinstance(data, list), f"Expected list, got {type(data).__name__}"


def test_get_videos_each_item_has_required_keys(client):
    """Each video item contains id, name, duration, and paragraph_count."""
    resp = client.get("/api/outline/videos")
    data = resp.json()
    assert len(data) > 0, "Expected at least one video"

    for item in data:
        for key in ("id", "name", "duration", "paragraph_count"):
            assert key in item, f"Missing key '{key}' in video item {item}"


def test_get_videos_all_have_status_done(client):
    """All returned videos have paragraphs (i.e. they are processed/done)."""
    resp = client.get("/api/outline/videos")
    data = resp.json()

    for item in data:
        assert item["paragraph_count"] > 0, (
            f"Video {item['id']} has paragraph_count={item['paragraph_count']}, "
            f"expected > 0 for done videos"
        )


def test_get_videos_returns_known_videos(client):
    """The videos endpoint returns the 3 known done videos with correct counts."""
    resp = client.get("/api/outline/videos")
    data = resp.json()

    # Build a lookup by id
    by_id = {item["id"]: item for item in data}

    # Known done videos from corpus.db
    assert 25 in by_id, "Video 25 (ccna-network lesson) missing"
    assert 29 in by_id, "Video 29 (WOLF-LAB CCNA-day2-6-UDP) missing"
    assert 31 in by_id, "Video 31 (WOLF-LAB CCNA-day6-18-VLAN) missing"

    # Paragraph counts from corpus.db
    assert by_id[25]["paragraph_count"] == 29, (
        f"Video 25 expected 29 paragraphs, got {by_id[25]['paragraph_count']}"
    )
    assert by_id[29]["paragraph_count"] == 343, (
        f"Video 29 expected 343 paragraphs, got {by_id[29]['paragraph_count']}"
    )
    assert by_id[31]["paragraph_count"] == 218, (
        f"Video 31 expected 218 paragraphs, got {by_id[31]['paragraph_count']}"
    )


# ---------------------------------------------------------------------------
# GET /api/outline/video/{video_id}/paragraphs
# ---------------------------------------------------------------------------


def test_get_paragraphs_returns_json_array(client):
    """GET /api/outline/video/25/paragraphs returns a JSON array."""
    resp = client.get("/api/outline/video/25/paragraphs")
    assert resp.status_code == 200, f"Expected 200, got {resp.status_code}"
    data = resp.json()
    assert isinstance(data, list), f"Expected list, got {type(data).__name__}"


def test_get_paragraphs_each_item_has_required_keys(client):
    """Each paragraph item contains paragraph_index, start_time, end_time, text."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()
    assert len(data) > 0, "Expected at least one paragraph"

    for item in data:
        for key in ("paragraph_index", "start_time", "end_time", "text"):
            assert key in item, f"Missing key '{key}' in paragraph item {item}"


def test_get_paragraphs_text_truncated_to_100_chars(client):
    """The text field is truncated to at most 100 characters for preview."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()

    for item in data:
        text = item["text"]
        assert len(text) <= 100, (
            f"Paragraph {item['paragraph_index']} text length is {len(text)}, "
            f"expected <= 100"
        )


def test_get_paragraphs_ordered_by_index(client):
    """Paragraphs are returned in ascending paragraph_index order."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()

    indices = [item["paragraph_index"] for item in data]
    assert indices == sorted(indices), "Paragraphs are not sorted by paragraph_index"


def test_get_paragraphs_returns_correct_count(client):
    """Video 25 returns exactly 29 paragraphs (matches corpus.db)."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()
    assert len(data) == 29, f"Expected 29 paragraphs for video 25, got {len(data)}"


def test_get_paragraphs_first_item_matches_db(client):
    """The first paragraph of video 25 has correct index and times."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()
    first = data[0]

    assert first["paragraph_index"] == 1, (
        f"Expected paragraph_index=1, got {first['paragraph_index']}"
    )
    assert first["start_time"] == 0.0, (
        f"Expected start_time=0.0, got {first['start_time']}"
    )
    assert first["end_time"] == 12.56, (
        f"Expected end_time=12.56, got {first['end_time']}"
    )


# ---------------------------------------------------------------------------
# GET /api/outline/video/{video_id}/paragraphs — 404 cases
# ---------------------------------------------------------------------------


def test_get_paragraphs_nonexistent_video_returns_404(client):
    """Requesting paragraphs for a non-existent video returns 404."""
    resp = client.get("/api/outline/video/99999/paragraphs")
    assert resp.status_code == 404, (
        f"Expected 404 for non-existent video, got {resp.status_code}"
    )


def test_get_paragraphs_nonexistent_video_has_detail(client):
    """The 404 response includes a detail message."""
    resp = client.get("/api/outline/video/99999/paragraphs")
    data = resp.json()
    assert "detail" in data, "404 response should contain 'detail' key"


# ---------------------------------------------------------------------------
# Edge case tests — Slice 04 (additional)
# ---------------------------------------------------------------------------


def test_get_paragraphs_time_fields_are_numeric(client):
    """start_time and end_time are numeric (int or float)."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()

    for item in data:
        assert isinstance(item["start_time"], (int, float)), (
            f"start_time should be numeric, got {type(item['start_time']).__name__}"
        )
        assert isinstance(item["end_time"], (int, float)), (
            f"end_time should be numeric, got {type(item['end_time']).__name__}"
        )


def test_get_paragraphs_text_is_non_empty(client):
    """Each paragraph's text field is a non-empty string."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()

    for item in data:
        assert isinstance(item["text"], str), (
            f"text should be str, got {type(item['text']).__name__}"
        )
        assert len(item["text"]) > 0, (
            f"Paragraph {item['paragraph_index']} text is empty"
        )


def test_get_paragraphs_negative_video_id_returns_404(client):
    """GET /api/outline/video/-1/paragraphs returns 404."""
    resp = client.get("/api/outline/video/-1/paragraphs")
    assert resp.status_code == 404


def test_get_paragraphs_zero_video_id_returns_404(client):
    """GET /api/outline/video/0/paragraphs returns 404."""
    resp = client.get("/api/outline/video/0/paragraphs")
    assert resp.status_code == 404


def test_get_paragraphs_large_video_count(client):
    """Video 29 with 343 paragraphs returns the correct count."""
    resp = client.get("/api/outline/video/29/paragraphs")
    data = resp.json()
    assert len(data) == 343, f"Expected 343 paragraphs for video 29, got {len(data)}"


def test_get_paragraphs_first_and_last_indices_match(client):
    """First paragraph has index 1 and last paragraph index equals count."""
    resp = client.get("/api/outline/video/25/paragraphs")
    data = resp.json()

    assert data[0]["paragraph_index"] == 1
    assert data[-1]["paragraph_index"] == len(data)


def test_get_paragraphs_non_string_video_id_returns_422(client):
    """GET /api/outline/video/abc/paragraphs returns 422 (validation error)."""
    resp = client.get("/api/outline/video/abc/paragraphs")
    assert resp.status_code == 422


# ============================================================================
# Slice 02 — SSE 全透明控制台
# ============================================================================

# ---------------------------------------------------------------------------
# SSE event formatting (unit tests — no server needed)
# ---------------------------------------------------------------------------


def test_event_to_sse_formats_info_event_correctly():
    """_event_to_sse() produces properly formatted SSE for an info event."""
    from main_outline import LogEvent, _event_to_sse

    event = LogEvent(category="info", message="Test info message")
    output = _event_to_sse(event)

    # Must start with event: line
    assert output.startswith("event: info\n")
    # Must contain data: line
    assert "\ndata: " in output
    # Must end with double newline
    assert output.endswith("\n\n")

    # Parse the data JSON
    lines = output.strip().split("\n")
    data_line = [l for l in lines if l.startswith("data: ")][0]
    data = json.loads(data_line[6:])  # strip "data: " prefix
    assert data["msg"] == "Test info message"
    assert data["category"] == "info"
    assert "ts" in data


def test_event_to_sse_handles_all_categories():
    """_event_to_sse() works for all five SSE categories."""
    from main_outline import LogEvent, _event_to_sse

    for category in ("info", "debug", "success", "error", "progress"):
        event = LogEvent(category=category, message=f"Test {category}")
        output = _event_to_sse(event)
        assert output.startswith(f"event: {category}\n")

        data_line = [l for l in output.strip().split("\n") if l.startswith("data: ")][0]
        data = json.loads(data_line[6:])
        assert data["category"] == category


def test_event_to_sse_includes_traceback_when_present():
    """_event_to_sse() includes traceback field when LogEvent has traceback set."""
    from main_outline import LogEvent, _event_to_sse

    tb_text = (
        'Traceback (most recent call last):\n'
        '  File "test.py", line 1, in <module>\n'
        '    1/0\n'
        'ZeroDivisionError: division by zero'
    )
    event = LogEvent(category="error", message="Something failed", traceback=tb_text)
    output = _event_to_sse(event)

    assert "event: error" in output
    data_line = [l for l in output.strip().split("\n") if l.startswith("data: ")][0]
    data = json.loads(data_line[6:])
    assert data["traceback"] == tb_text


def test_event_to_sse_includes_progress_when_present():
    """_event_to_sse() includes progress field when LogEvent has progress_pct set."""
    from main_outline import LogEvent, _event_to_sse

    event = LogEvent(category="progress", message="Processing...", progress_pct=50.0)
    output = _event_to_sse(event)

    data_line = [l for l in output.strip().split("\n") if l.startswith("data: ")][0]
    data = json.loads(data_line[6:])
    assert data["progress"] == 50.0


def test_event_to_sse_omits_traceback_when_none():
    """_event_to_sse() does not include a traceback key when traceback is None."""
    from main_outline import LogEvent, _event_to_sse

    event = LogEvent(category="info", message="All good")
    output = _event_to_sse(event)

    data_line = [l for l in output.strip().split("\n") if l.startswith("data: ")][0]
    data = json.loads(data_line[6:])
    assert "traceback" not in data


# ---------------------------------------------------------------------------
# Log level → SSE category mapping
# ---------------------------------------------------------------------------


def test_level_to_category_maps_correctly():
    """_level_to_category() maps Python log levels to SSE categories."""
    from main_outline import _level_to_category

    assert _level_to_category(logging.DEBUG) == "debug"
    assert _level_to_category(logging.INFO) == "info"
    assert _level_to_category(logging.WARNING) == "info"
    assert _level_to_category(logging.ERROR) == "error"
    assert _level_to_category(logging.CRITICAL) == "error"


# ---------------------------------------------------------------------------
# SSE endpoint — content type
# ---------------------------------------------------------------------------


def test_sse_endpoint_returns_event_stream_content_type():
    """GET /api/stream/logs/outline returns a StreamingResponse with
    text/event-stream media type.

    We test by calling the endpoint function directly and inspecting the
    returned Starlette StreamingResponse, avoiding the hang caused by
    the infinite SSE stream in TestClient.
    """
    import asyncio
    from main_outline import stream_logs

    resp = asyncio.run(stream_logs())
    assert resp.media_type == "text/event-stream"
    assert "no-cache" in resp.headers.get("cache-control", "")
    assert "keep-alive" in resp.headers.get("connection", "")


# ---------------------------------------------------------------------------
# POST /api/outline/trigger-error
# ---------------------------------------------------------------------------


def test_trigger_error_endpoint_returns_200(client):
    """POST /api/outline/trigger-error returns 200 and confirms error triggered."""
    resp = client.post("/api/outline/trigger-error")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "error_triggered"


# ---------------------------------------------------------------------------
# SSELogHandler is attached to the outline logger
# ---------------------------------------------------------------------------


def test_sse_log_handler_configured(client):
    """After lifespan startup, the 'outline' logger has SSELogHandler attached."""
    from main_outline import SSELogHandler

    logger = logging.getLogger("outline")
    handlers = logger.handlers
    assert any(isinstance(h, SSELogHandler) for h in handlers), (
        f"SSELogHandler not found among outline logger handlers: {handlers}"
    )


# ---------------------------------------------------------------------------
# SSE async generator yields events from queue
# ---------------------------------------------------------------------------


@pytest.mark.asyncio
async def test_sse_event_generator_yields_events():
    """The async SSE generator yields formatted events from the queue."""
    from main_outline import LogEvent, sse_event_generator

    q = asyncio.Queue()
    await q.put(LogEvent(category="info", message="Hello from queue"))

    gen = sse_event_generator(q)
    event_text = await gen.__anext__()

    assert "event: info" in event_text
    assert "Hello from queue" in event_text


# ---------------------------------------------------------------------------
# push_log_event helper
# ---------------------------------------------------------------------------


def test_push_log_event_adds_to_queue():
    """push_log_event() puts a LogEvent into the module-level queue."""
    from main_outline import push_log_event, get_log_queue, LogEvent

    # Reset the module-level queue to a clean state for this test
    import main_outline
    main_outline._log_queue = asyncio.Queue()

    push_log_event("info", "Test push")

    q = get_log_queue()
    assert not q.empty()
    event = q.get_nowait()
    assert isinstance(event, LogEvent)
    assert event.category == "info"
    assert event.message == "Test push"


def test_push_log_event_with_progress():
    """push_log_event() accepts progress_pct for progress events."""
    from main_outline import push_log_event, get_log_queue

    import main_outline
    main_outline._log_queue = asyncio.Queue()

    push_log_event("progress", "Half done", progress_pct=50.0)

    q = get_log_queue()
    event = q.get_nowait()
    assert event.category == "progress"
    assert event.progress_pct == 50.0


def test_push_log_event_with_traceback():
    """push_log_event() accepts traceback for error events."""
    from main_outline import push_log_event, get_log_queue

    import main_outline
    main_outline._log_queue = asyncio.Queue()

    tb = "ZeroDivisionError: division by zero"
    push_log_event("error", "Bad error", traceback=tb)

    q = get_log_queue()
    event = q.get_nowait()
    assert event.category == "error"
    assert event.traceback == tb


# ---------------------------------------------------------------------------
# LogEvent dataclass defaults
# ---------------------------------------------------------------------------


def test_log_event_defaults():
    """LogEvent has sensible defaults: ts auto-generated, traceback None,
    progress_pct None."""
    from main_outline import LogEvent

    event = LogEvent(category="debug", message="Test defaults")
    assert event.ts  # auto-generated timestamp
    assert event.traceback is None
    assert event.progress_pct is None
    assert event.category == "debug"
    assert event.message == "Test defaults"


# ============================================================================
# Slice 06 — 第一轮 LLM 分析（话题切分）
# ============================================================================

# ---------------------------------------------------------------------------
# JSON parsing helpers (unit tests — no server needed)
# ---------------------------------------------------------------------------

VALID_LLM_RESPONSE = """{"course_name": "计算机网络基础", "topics": [{"name": "UDP概述", "start_para": 1, "end_para": 50}, {"name": "UDP头部字段", "start_para": 51, "end_para": 120}]}"""

VALID_LLM_RESPONSE_SINGLE_TOPIC = """{"course_name": "VLAN入门", "topics": [{"name": "VLAN概述", "start_para": 1, "end_para": 218}]}"""

INVALID_LLM_RESPONSE = "This is not JSON at all, just some random text from the model..."

INVALID_LLM_RESPONSE_WRONG_STRUCTURE = """{"course_title": "Wrong key", "sections": [{"title": "Section 1", "start": 1, "end": 50}]}"""

INVALID_LLM_RESPONSE_TRUNCATED = """{"course_name": "Intro", "topics": [{"name": "Topic 1", "start_para": 1"""


def test_parse_llm_response_valid():
    """parse_round1_response() correctly parses a valid LLM JSON response."""
    from main_outline import parse_round1_response

    result = parse_round1_response(VALID_LLM_RESPONSE)
    assert result is not None
    assert result["course_name"] == "计算机网络基础"
    assert len(result["topics"]) == 2
    assert result["topics"][0] == {"name": "UDP概述", "start_para": 1, "end_para": 50}
    assert result["topics"][1] == {"name": "UDP头部字段", "start_para": 51, "end_para": 120}


def test_parse_llm_response_single_topic():
    """parse_round1_response() works with a single topic."""
    from main_outline import parse_round1_response

    result = parse_round1_response(VALID_LLM_RESPONSE_SINGLE_TOPIC)
    assert result is not None
    assert result["course_name"] == "VLAN入门"
    assert len(result["topics"]) == 1


def test_parse_llm_response_invalid_json():
    """parse_round1_response() returns None for non-JSON text."""
    from main_outline import parse_round1_response

    result = parse_round1_response(INVALID_LLM_RESPONSE)
    assert result is None


def test_parse_llm_response_wrong_structure():
    """parse_round1_response() returns None when JSON has wrong keys."""
    from main_outline import parse_round1_response

    result = parse_round1_response(INVALID_LLM_RESPONSE_WRONG_STRUCTURE)
    assert result is None


def test_parse_llm_response_truncated():
    """parse_round1_response() returns None for truncated JSON."""
    from main_outline import parse_round1_response

    result = parse_round1_response(INVALID_LLM_RESPONSE_TRUNCATED)
    assert result is None


def test_parse_llm_response_empty_string():
    """parse_round1_response() returns None for empty string."""
    from main_outline import parse_round1_response

    result = parse_round1_response("")
    assert result is None


def test_parse_llm_response_strips_markdown_fences():
    """parse_round1_response() strips ```json fences before parsing."""
    from main_outline import parse_round1_response

    fenced = '```json\n' + VALID_LLM_RESPONSE + '\n```'
    result = parse_round1_response(fenced)
    assert result is not None
    assert result["course_name"] == "计算机网络基础"


def test_parse_llm_response_strips_leading_text():
    """parse_round1_response() extracts JSON even when there is leading text."""
    from main_outline import parse_round1_response

    messy = 'Here is the analysis:\n\n' + VALID_LLM_RESPONSE + '\n\nHope this helps!'
    result = parse_round1_response(messy)
    assert result is not None
    assert result["course_name"] == "计算机网络基础"


# ---------------------------------------------------------------------------
# Worker thread: full-text concatenation (unit test — no server)
# ---------------------------------------------------------------------------


def test_concat_paragraphs_to_full_text():
    """concat_paragraphs() joins paragraph texts in order with paragraph markers."""
    from main_outline import concat_paragraphs

    paragraphs = [
        {"paragraph_index": 1, "text": "Hello world"},
        {"paragraph_index": 2, "text": "This is paragraph two"},
        {"paragraph_index": 3, "text": "Third paragraph here"},
    ]
    full_text, para_map = concat_paragraphs(paragraphs)

    # Should contain paragraph markers
    assert "[段落 1]" in full_text
    assert "[段落 2]" in full_text
    assert "[段落 3]" in full_text
    # Should contain all text
    assert "Hello world" in full_text
    assert "This is paragraph two" in full_text
    assert "Third paragraph here" in full_text
    # para_map should be empty dict for preview data (no start_para mapping needed)
    assert isinstance(para_map, dict)


# ---------------------------------------------------------------------------
# POST /api/outline/analyze/{video_id}
# ---------------------------------------------------------------------------


def test_analyze_endpoint_returns_202(client):
    """POST /api/outline/analyze/{video_id} returns 202 Accepted.

    We mock the Ollama call so the background thread does not make a real
    HTTP request, and we mock the paragraph fetch to return controlled data.
    """
    from unittest.mock import patch, MagicMock
    import time

    mock_ollama_response = {
        "response": VALID_LLM_RESPONSE,
    }

    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 10.0,
         "text": "First paragraph text here"},
        {"paragraph_index": 2, "start_time": 10.0, "end_time": 20.0,
         "text": "Second paragraph text here"},
        {"paragraph_index": 3, "start_time": 20.0, "end_time": 30.0,
         "text": "Third paragraph text here"},
    ]

    # Reset cancel flag before test
    import main_outline
    main_outline._analysis_cancel_flags.pop(25, None)
    main_outline._topic_segmentation_results.pop(25, None)

    with patch("main_outline._fetch_paragraphs_full_text",
               return_value=mock_paragraphs):
        with patch("main_outline._call_ollama_generate",
                   return_value=mock_ollama_response):
            resp = client.post("/api/outline/analyze/25")
            assert resp.status_code == 202, (
                f"Expected 202, got {resp.status_code}: {resp.text}"
            )
            data = resp.json()
            assert data["status"] == "started"
            assert data["video_id"] == 25

            # Wait for thread to complete (max 5 seconds)
            deadline = time.time() + 5.0
            result = None
            while time.time() < deadline:
                result = main_outline._topic_segmentation_results.get(25)
                if result is not None:
                    break
                time.sleep(0.1)

            assert result is not None, (
                "Analysis result was not stored within 5 seconds"
            )
            assert result["course_name"] == "计算机网络基础"
            assert len(result["topics"]) == 2


def test_analyze_endpoint_nonexistent_video_returns_404(client):
    """POST /api/outline/analyze/{video_id} returns 404 for non-existent video."""
    resp = client.post("/api/outline/analyze/99999")
    assert resp.status_code == 404


def test_analyze_endpoint_already_running_returns_409(client):
    """POST /api/outline/analyze returns 409 when analysis is already running for
    that video. We simulate by pre-populating the cancel flag with False (meaning
    a thread is running)."""
    import main_outline

    # Simulate a running analysis
    main_outline._analysis_cancel_flags[25] = False

    try:
        resp = client.post("/api/outline/analyze/25")
        assert resp.status_code == 409, (
            f"Expected 409 for already-running analysis, got {resp.status_code}"
        )
    finally:
        main_outline._analysis_cancel_flags.pop(25, None)


# ---------------------------------------------------------------------------
# POST /api/outline/stop/{video_id}
# ---------------------------------------------------------------------------


def test_stop_endpoint_sets_cancel_flag(client):
    """POST /api/outline/stop/{video_id} returns 200 and sets the cancel flag."""
    import main_outline

    # Set up: pretend analysis is running
    main_outline._analysis_cancel_flags[25] = False

    try:
        resp = client.post("/api/outline/stop/25")
        assert resp.status_code == 200
        data = resp.json()
        assert data["status"] == "cancelling"

        # The flag should now be True
        assert main_outline._analysis_cancel_flags.get(25) is True
    finally:
        main_outline._analysis_cancel_flags.pop(25, None)


def test_stop_endpoint_idempotent(client):
    """POST /api/outline/stop returns 200 even if no analysis was running
    for that video."""
    import main_outline

    main_outline._analysis_cancel_flags.pop(25, None)

    resp = client.post("/api/outline/stop/25")
    assert resp.status_code == 200
    data = resp.json()
    assert data["status"] == "cancelling"


def test_stop_endpoint_nonexistent_video_returns_200(client):
    """POST /api/outline/stop returns 200 even for non-existent video
    (stop is a fire-and-forget signal)."""
    resp = client.post("/api/outline/stop/99999")
    assert resp.status_code == 200


# ---------------------------------------------------------------------------
# Cancel flag: worker thread checks flag between steps
# ---------------------------------------------------------------------------


def test_cancel_flag_stops_worker_before_ollama():
    """When cancel flag is set, the worker thread does not call Ollama
    and pushes a 'cancelled' info event."""
    import main_outline

    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 10.0,
         "text": "Some text"},
    ]

    # Set cancel flag before starting
    main_outline._analysis_cancel_flags[25] = True
    main_outline._topic_segmentation_results.pop(25, None)

    from unittest.mock import patch, MagicMock

    mock_ollama = MagicMock()

    with patch("main_outline._fetch_paragraphs_full_text",
               return_value=mock_paragraphs):
        with patch("main_outline._call_ollama_generate", mock_ollama):
            main_outline._run_round1_analysis(25)

            # Ollama should NOT have been called
            mock_ollama.assert_not_called()

            # Result should NOT be stored
            assert 25 not in main_outline._topic_segmentation_results


# ---------------------------------------------------------------------------
# JSON parse failure: sends error SSE event with raw output
# ---------------------------------------------------------------------------


def test_parse_failure_pushes_error_event():
    """When JSON parsing fails, the worker pushes an error event with
    first 500 chars of raw LLM output."""
    import main_outline

    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 10.0,
         "text": "Some text for LLM analysis"},
    ]

    mock_ollama_response = {
        "response": INVALID_LLM_RESPONSE,
    }

    main_outline._analysis_cancel_flags.pop(25, None)
    main_outline._topic_segmentation_results.pop(25, None)

    from unittest.mock import patch

    with patch("main_outline._fetch_paragraphs_full_text",
               return_value=mock_paragraphs):
        with patch("main_outline._call_ollama_generate",
                   return_value=mock_ollama_response):
            main_outline._run_round1_analysis(25)

            # Result should NOT be stored (parse failed)
            assert 25 not in main_outline._topic_segmentation_results


# ---------------------------------------------------------------------------
# i18n labels for Slice 06 frontend
# ---------------------------------------------------------------------------


def test_i18n_has_analyze_button_label():
    """The I18N.zh dictionary contains labels for the analyze button."""
    import json

    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")

    # Check i18n keys exist — keys are defined with single quotes in the
    # JavaScript I18N object, and some also appear in data-i18n attributes
    # with double quotes.
    assert "'analyze.btn'" in content or '"analyze.btn"' in content, (
        "Missing i18n key analyze.btn"
    )
    assert "'analyze.running'" in content or '"analyze.running"' in content, (
        "Missing i18n key analyze.running"
    )
    assert "'analyze.success'" in content or '"analyze.success"' in content, (
        "Missing i18n key analyze.success"
    )
    assert "'analyze.error'" in content or '"analyze.error"' in content, (
        "Missing i18n key analyze.error"
    )
    assert "'analyze.stopping'" in content or '"analyze.stopping"' in content, (
        "Missing i18n key analyze.stopping"
    )


# ============================================================================
# Slice 07 — 第二轮 LLM 分析（子树生成）+ 落盘
# ============================================================================

# ---------------------------------------------------------------------------
# parse_round2_response — unit tests
# ---------------------------------------------------------------------------

VALID_LLM_RESPONSE_ROUND2 = """{"topic": "UDP概述", "subtree": {"children": [{"name": "UDP定义", "children": []}, {"name": "UDP特点", "children": [{"name": "无连接", "children": []}, {"name": "不可靠", "children": []}]}]}}"""

VALID_LLM_RESPONSE_ROUND2_EMPTY_CHILDREN = """{"topic": "简单概念", "subtree": {"children": []}}"""

VALID_LLM_RESPONSE_ROUND2_DEEP = """{"topic": "校验和计算", "subtree": {"children": [{"name": "计算原理", "children": [{"name": "二进制加法", "children": [{"name": "进位处理", "children": []}]}]}, {"name": "例题演示", "children": []}]}}"""

INVALID_LLM_RESPONSE_ROUND2 = "Sorry, I cannot analyze this topic."

INVALID_LLM_RESPONSE_ROUND2_WRONG_STRUCTURE = """{"title": "UDP", "nodes": [{"label": "Node1"}]}"""

INVALID_LLM_RESPONSE_ROUND2_TRUNCATED = """{"topic": "UDP", "subtree": {"children": [{"name": "Def"""


def test_parse_round2_response_valid():
    """parse_round2_response() correctly parses a valid round-2 LLM JSON response."""
    from main_outline import parse_round2_response

    result = parse_round2_response(VALID_LLM_RESPONSE_ROUND2)
    assert result is not None
    assert result["topic"] == "UDP概述"
    assert "subtree" in result
    assert "children" in result["subtree"]
    assert len(result["subtree"]["children"]) == 2
    assert result["subtree"]["children"][0]["name"] == "UDP定义"
    assert result["subtree"]["children"][1]["name"] == "UDP特点"


def test_parse_round2_response_empty_children():
    """parse_round2_response() accepts a subtree with empty children array."""
    from main_outline import parse_round2_response

    result = parse_round2_response(VALID_LLM_RESPONSE_ROUND2_EMPTY_CHILDREN)
    assert result is not None
    assert result["subtree"]["children"] == []


def test_parse_round2_response_deeply_nested():
    """parse_round2_response() handles deeply nested subtree structures."""
    from main_outline import parse_round2_response

    result = parse_round2_response(VALID_LLM_RESPONSE_ROUND2_DEEP)
    assert result is not None
    # Verify deep nesting preserved
    assert result["subtree"]["children"][0]["children"][0]["children"][0]["name"] == "进位处理"


def test_parse_round2_response_invalid_json():
    """parse_round2_response() returns None for non-JSON text."""
    from main_outline import parse_round2_response

    result = parse_round2_response(INVALID_LLM_RESPONSE_ROUND2)
    assert result is None


def test_parse_round2_response_wrong_structure():
    """parse_round2_response() returns None when JSON has wrong keys (no 'subtree')."""
    from main_outline import parse_round2_response

    result = parse_round2_response(INVALID_LLM_RESPONSE_ROUND2_WRONG_STRUCTURE)
    assert result is None


def test_parse_round2_response_truncated():
    """parse_round2_response() returns None for truncated JSON."""
    from main_outline import parse_round2_response

    result = parse_round2_response(INVALID_LLM_RESPONSE_ROUND2_TRUNCATED)
    assert result is None


def test_parse_round2_response_empty_string():
    """parse_round2_response() returns None for empty string."""
    from main_outline import parse_round2_response

    result = parse_round2_response("")
    assert result is None


def test_parse_round2_response_strips_markdown_fences():
    """parse_round2_response() strips ```json fences before parsing."""
    from main_outline import parse_round2_response

    fenced = '```json\n' + VALID_LLM_RESPONSE_ROUND2 + '\n```'
    result = parse_round2_response(fenced)
    assert result is not None
    assert result["topic"] == "UDP概述"


def test_parse_round2_response_strips_leading_text():
    """parse_round2_response() extracts JSON even when there is leading/trailing text."""
    from main_outline import parse_round2_response

    messy = 'Here is the subtree analysis:\n\n' + VALID_LLM_RESPONSE_ROUND2 + '\n\nDone.'
    result = parse_round2_response(messy)
    assert result is not None
    assert result["topic"] == "UDP概述"


# ---------------------------------------------------------------------------
# Round 2 flow — INSERT into course_topics (integration tests)
# ---------------------------------------------------------------------------


def test_round2_inserts_rows_into_course_topics():
    """_run_round2_analysis INSERTs one row per topic into course_topics."""
    import main_outline
    from unittest.mock import patch, MagicMock

    # Set up round1 result in module-level dict
    main_outline._topic_segmentation_results[25] = {
        "course_name": "计算机网络基础",
        "topics": [
            {"name": "UDP概述", "start_para": 1, "end_para": 10},
            {"name": "UDP头部字段", "start_para": 11, "end_para": 20},
        ],
    }

    # Mock paragraphs (10 per topic to satisfy the range)
    mock_paragraphs = []
    for i in range(1, 21):
        mock_paragraphs.append({
            "paragraph_index": i,
            "start_time": float(i * 10.0),
            "end_time": float(i * 10.0 + 5.0),
            "text": f"Paragraph {i} text content for LLM analysis",
        })

    # Two different round2 responses for two topics
    round2_responses = [
        {"response": '{"topic": "UDP概述", "subtree": {"children": [{"name": "定义", "children": []}]}}'},
        {"response": '{"topic": "UDP头部字段", "subtree": {"children": [{"name": "源端口", "children": []}, {"name": "目的端口", "children": []}]}}'},
    ]

    # Clean up any existing rows for video 25
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    main_outline._analysis_cancel_flags.pop(25, None)

    call_count = [0]
    def mock_ollama(model, prompt, stream=False):
        idx = call_count[0]
        call_count[0] += 1
        return round2_responses[idx]

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       side_effect=mock_ollama):
                main_outline._run_round2_analysis(25)

        # Verify: 2 rows inserted
        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT * FROM course_topics WHERE video_id = 25 ORDER BY start_para_index"
        ).fetchall()
        conn.close()

        assert len(rows) == 2, f"Expected 2 rows, got {len(rows)}"

        # Row 1
        r1 = rows[0]
        assert r1["video_id"] == 25
        assert r1["start_para_index"] == 1
        assert r1["end_para_index"] == 10
        assert r1["start_time"] == 10.0  # first para in range (idx=1, time=10.0)
        assert r1["end_time"] == 105.0  # last para (idx=10, time=105.0)
        assert r1["topic_name"] == "UDP概述"

        # Verify subtree_json
        subtree1 = json.loads(r1["subtree_json"])
        assert subtree1["course_name"] == "计算机网络基础"
        assert "subtree" in subtree1
        assert len(subtree1["subtree"]["children"]) == 1
        assert subtree1["subtree"]["children"][0]["name"] == "定义"

        # Row 2
        r2 = rows[1]
        assert r2["video_id"] == 25
        assert r2["start_para_index"] == 11
        assert r2["end_para_index"] == 20
        assert r2["start_time"] == 110.0  # first para in range (idx=11, time=110.0)
        assert r2["end_time"] == 205.0  # last para (idx=20, time=205.0)
        assert r2["topic_name"] == "UDP头部字段"

        subtree2 = json.loads(r2["subtree_json"])
        assert len(subtree2["subtree"]["children"]) == 2

    finally:
        # Clean up
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_subtree_json_format_correct():
    """Each course_topics row has subtree_json with course_name and subtree keys."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "VLAN基础",
        "topics": [
            {"name": "VLAN概念", "start_para": 1, "end_para": 5},
        ],
    }

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i * 5.0),
         "end_time": float(i * 5.0 + 3.0),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]

    mock_response = {
        "response": '{"topic": "VLAN概念", "subtree": {"children": [{"name": "802.1Q", "children": []}]}}'
    }

    main_outline._analysis_cancel_flags.pop(25, None)

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_round2_analysis(25)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT subtree_json FROM course_topics WHERE video_id = 25"
        ).fetchone()
        conn.close()

        assert row is not None
        subtree = json.loads(row["subtree_json"])

        # Must have course_name and subtree
        assert subtree["course_name"] == "VLAN基础"
        assert "subtree" in subtree
        assert "children" in subtree["subtree"]
        assert isinstance(subtree["subtree"]["children"], list)
        assert len(subtree["subtree"]["children"]) == 1
        assert subtree["subtree"]["children"][0]["name"] == "802.1Q"
        assert subtree["subtree"]["children"][0]["children"] == []

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_time_fields_derived_from_paragraphs():
    """start_time and end_time are correctly derived from first and last paragraph
    in the topic's range."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "测试课程",
        "topics": [
            {"name": "第一节", "start_para": 3, "end_para": 7},
        ],
    }

    # Paragraphs 1-10 with known times
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i * 12.0),
         "end_time": float(i * 12.0 + 8.0),
         "text": f"Paragraph {i}"}
        for i in range(1, 11)
    ]

    mock_response = {
        "response": '{"topic": "第一节", "subtree": {"children": []}}'
    }

    main_outline._analysis_cancel_flags.pop(25, None)

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_round2_analysis(25)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT start_para_index, end_para_index, start_time, end_time "
            "FROM course_topics WHERE video_id = 25"
        ).fetchone()
        conn.close()

        assert row is not None
        # Topic covers paragraphs 3-7
        assert row["start_para_index"] == 3
        assert row["end_para_index"] == 7
        # start_time = paragraph 3's start_time = 3 * 12 = 36.0
        assert row["start_time"] == 36.0
        # end_time = paragraph 7's end_time = 7 * 12 + 8 = 92.0
        assert row["end_time"] == 92.0

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_single_topic_failure_does_not_block_others():
    """When one topic's JSON parse fails, that topic is skipped but remaining
    topics are still processed and inserted."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "测试课程",
        "topics": [
            {"name": "话题A", "start_para": 1, "end_para": 3},
            {"name": "话题B", "start_para": 4, "end_para": 6},
            {"name": "话题C", "start_para": 7, "end_para": 9},
        ],
    }

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 10)
    ]

    # Topic B returns invalid JSON, others return valid
    round2_responses = [
        {"response": '{"topic": "话题A", "subtree": {"children": [{"name": "A1", "children": []}]}}'},
        {"response": 'NOT VALID JSON AT ALL'},
        {"response": '{"topic": "话题C", "subtree": {"children": [{"name": "C1", "children": []}]}}'},
    ]

    call_count = [0]
    def mock_ollama(model, prompt, stream=False):
        idx = call_count[0]
        call_count[0] += 1
        return round2_responses[idx]

    main_outline._analysis_cancel_flags.pop(25, None)

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       side_effect=mock_ollama):
                main_outline._run_round2_analysis(25)

        # Only topics A and C should be inserted
        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT topic_name FROM course_topics WHERE video_id = 25 "
            "ORDER BY start_para_index"
        ).fetchall()
        conn.close()

        topic_names = [r["topic_name"] for r in rows]
        assert len(topic_names) == 2, f"Expected 2 topics inserted, got {len(topic_names)}: {topic_names}"
        assert "话题A" in topic_names
        assert "话题C" in topic_names
        assert "话题B" not in topic_names

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_cancel_between_topics():
    """When cancel flag is set between topics, processing stops.
    Already-completed topics remain in the database."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "测试课程",
        "topics": [
            {"name": "话题1", "start_para": 1, "end_para": 3},
            {"name": "话题2", "start_para": 4, "end_para": 6},
            {"name": "话题3", "start_para": 7, "end_para": 9},
        ],
    }

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 10)
    ]

    mock_response = {
        "response": '{"topic": "dummy", "subtree": {"children": []}}'
    }

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    # Set cancel after first topic completes, but the worker checks the flag
    # at the start of each iteration, so we set it BEFORE calling and
    # the mock ollama will set the flag when the second topic starts.
    # Simpler: set the flag to True after the first call completes.
    # We use a side_effect to simulate cancel.

    main_outline._analysis_cancel_flags[25] = False  # running

    call_count = [0]
    def mock_ollama_with_cancel(model, prompt, stream=False):
        idx = call_count[0]
        call_count[0] += 1
        # After the first call (index 0 completes), set cancel flag
        if idx == 0:
            # Let the first topic complete, then cancel before the next
            main_outline._analysis_cancel_flags[25] = True
        return mock_response

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       side_effect=mock_ollama_with_cancel):
                main_outline._run_round2_analysis(25)

        # Only the first topic should be inserted (topic 2 and 3 are skipped)
        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT topic_name FROM course_topics WHERE video_id = 25 "
            "ORDER BY start_para_index"
        ).fetchall()
        conn.close()

        topic_names = [r["topic_name"] for r in rows]
        # After the first one completes, the flag is set, and at the beginning
        # of the next loop iteration the flag is checked — so only 1 topic
        # should be inserted.
        assert len(topic_names) == 1, (
            f"Expected 1 topic (cancel after first), got {len(topic_names)}: {topic_names}"
        )
        assert topic_names[0] == "话题1"

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_no_round1_result():
    """_run_round2_analysis pushes error and returns early when no round1 result
    exists for the video_id."""
    import main_outline
    from unittest.mock import patch

    # Ensure no round1 result
    main_outline._topic_segmentation_results.pop(25, None)
    main_outline._analysis_cancel_flags.pop(25, None)

    # Clear the log queue for clean assertions
    main_outline._log_queue = None

    with patch("main_outline._fetch_paragraphs_full_text") as mock_fetch:
        main_outline._run_round2_analysis(25)
        # Should not have fetched paragraphs at all
        mock_fetch.assert_not_called()


def test_round2_empty_topics_list():
    """_run_round2_analysis handles empty topics list gracefully."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "空课程",
        "topics": [],
    }
    main_outline._analysis_cancel_flags.pop(25, None)

    with patch("main_outline._fetch_paragraphs_full_text") as mock_fetch:
        main_outline._run_round2_analysis(25)
        mock_fetch.assert_not_called()

    main_outline._topic_segmentation_results.pop(25, None)


def test_round2_no_paragraphs_for_topic():
    """When a topic's paragraph range yields no paragraphs, it is skipped
    and other topics are still processed."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "测试",
        "topics": [
            {"name": "有效话题", "start_para": 1, "end_para": 3},
            {"name": "空话题", "start_para": 100, "end_para": 200},
        ],
    }

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 4)
    ]

    mock_response = {
        "response": '{"topic": "有效话题", "subtree": {"children": []}}'
    }

    main_outline._analysis_cancel_flags.pop(25, None)

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_round2_analysis(25)

        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT topic_name FROM course_topics WHERE video_id = 25"
        ).fetchall()
        conn.close()

        topic_names = [r["topic_name"] for r in rows]
        assert topic_names == ["有效话题"], (
            f"Expected only '有效话题', got {topic_names}"
        )

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


def test_round2_ollama_call_failure_per_topic():
    """When Ollama raises an exception for one topic, that topic is skipped
    and other topics continue to be processed."""
    import main_outline
    from unittest.mock import patch

    main_outline._topic_segmentation_results[25] = {
        "course_name": "测试",
        "topics": [
            {"name": "话题X", "start_para": 1, "end_para": 3},
            {"name": "话题Y", "start_para": 4, "end_para": 6},
        ],
    }

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 7)
    ]

    # First topic: Ollama fails. Second topic: succeeds.
    def mock_ollama_fail_first(model, prompt, stream=False):
        if "话题X" in prompt:
            raise Exception("Ollama connection refused")
        return {"response": '{"topic": "话题Y", "subtree": {"children": []}}'}

    main_outline._analysis_cancel_flags.pop(25, None)

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       side_effect=mock_ollama_fail_first):
                main_outline._run_round2_analysis(25)

        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT topic_name FROM course_topics WHERE video_id = 25"
        ).fetchall()
        conn.close()

        topic_names = [r["topic_name"] for r in rows]
        assert topic_names == ["话题Y"], (
            f"Expected only '话题Y' (话题X failed), got {topic_names}"
        )

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


# ---------------------------------------------------------------------------
# Full pipeline: round 1 → round 2 (end-to-end)
# ---------------------------------------------------------------------------


def test_full_pipeline_runs_round1_then_round2():
    """_run_full_analysis runs round1 (stores in memory) then round2
    (INSERTs into course_topics)."""
    import main_outline
    from unittest.mock import patch
    import time

    main_outline._analysis_cancel_flags.pop(25, None)
    main_outline._topic_segmentation_results.pop(25, None)

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Paragraph {i} content for analysis"}
        for i in range(1, 6)
    ]

    mock_round1_response = {
        "response": '{"course_name": "测试课程", "topics": [{"name": "话题1", "start_para": 1, "end_para": 3}, {"name": "话题2", "start_para": 4, "end_para": 5}]}'
    }

    mock_round2_responses = [
        {"response": '{"topic": "话题1", "subtree": {"children": [{"name": "子1", "children": []}]}}'},
        {"response": '{"topic": "话题2", "subtree": {"children": []}}'},
    ]

    call_count = [0]
    def mock_ollama(model, prompt, stream=False):
        idx = call_count[0]
        call_count[0] += 1
        # First call is round1
        if idx == 0:
            return mock_round1_response
        # Subsequent calls are round2
        return mock_round2_responses[idx - 1]

    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_topics WHERE video_id = 25")
    conn.execute("COMMIT")
    conn.close()

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       side_effect=mock_ollama):
                main_outline._run_full_analysis(25)

        # Round1 result should be stored
        assert 25 in main_outline._topic_segmentation_results
        assert main_outline._topic_segmentation_results[25]["course_name"] == "测试课程"

        # Round2: 2 rows in course_topics
        conn = main_outline.get_db(str(DB_PATH))
        rows = conn.execute(
            "SELECT topic_name, start_para_index, end_para_index, subtree_json "
            "FROM course_topics WHERE video_id = 25 ORDER BY start_para_index"
        ).fetchall()
        conn.close()

        assert len(rows) == 2
        assert rows[0]["topic_name"] == "话题1"
        assert rows[0]["start_para_index"] == 1
        assert rows[0]["end_para_index"] == 3
        assert rows[1]["topic_name"] == "话题2"
        assert rows[1]["start_para_index"] == 4
        assert rows[1]["end_para_index"] == 5

        # Verify subtree_json has course_name
        subtree1 = json.loads(rows[0]["subtree_json"])
        assert subtree1["course_name"] == "测试课程"

        # Cancel flag should be cleaned up
        assert 25 not in main_outline._analysis_cancel_flags

    finally:
        main_outline._topic_segmentation_results.pop(25, None)
        main_outline._analysis_cancel_flags.pop(25, None)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_topics WHERE video_id = 25")
        conn.execute("COMMIT")
        conn.close()


# ---------------------------------------------------------------------------
# Slice 09 — GET /api/outline/topics/{video_id}
# ---------------------------------------------------------------------------


def test_get_topics_returns_empty_list_for_nonexistent_video(client):
    """GET /api/outline/topics/{video_id} returns [] when video has no topics."""
    resp = client.get("/api/outline/topics/99999")
    assert resp.status_code == 200
    assert resp.json() == []


def test_get_topics_returns_data_for_video_with_topics(client):
    """GET /api/outline/topics/{video_id} returns topics with expected fields."""
    import main_outline

    # Insert a test topic row directly (disable FK for synthetic video_id)
    conn = main_outline.get_db(str(DB_PATH))
    try:
        conn.execute("PRAGMA foreign_keys = OFF")
        conn.execute(
            "INSERT INTO course_topics "
            "(video_id, start_para_index, end_para_index, start_time, end_time, "
            "topic_name, subtree_json) "
            "VALUES (99999, 0, 5, 10.5, 45.3, 'Test Topic', '{}')"
        )
        conn.execute("PRAGMA foreign_keys = ON")
        conn.execute("COMMIT")

        resp = client.get("/api/outline/topics/99999")
        assert resp.status_code == 200
        data = resp.json()
        assert isinstance(data, list)
        assert len(data) == 1

        topic = data[0]
        assert topic["topic_name"] == "Test Topic"
        assert topic["start_time"] == 10.5
        assert topic["end_time"] == 45.3
        assert topic["duration"] == 34.8
        assert topic["start_para_index"] == 0
        assert topic["end_para_index"] == 5
        assert topic["video_id"] == 99999
        assert "id" in topic
        assert "subtree_json" in topic

    finally:
        conn.execute("DELETE FROM course_topics WHERE video_id = 99999")
        conn.execute("COMMIT")
        conn.close()


def test_get_topics_returns_multiple_topics_ordered_by_start_para(client):
    """Multiple topics are returned sorted by start_para_index."""
    import main_outline

    conn = main_outline.get_db(str(DB_PATH))
    try:
        conn.execute("PRAGMA foreign_keys = OFF")
        conn.execute(
            "INSERT INTO course_topics "
            "(video_id, start_para_index, end_para_index, start_time, end_time, "
            "topic_name, subtree_json) "
            "VALUES (99998, 10, 15, 100.0, 150.0, 'Topic B', '{}')"
        )
        conn.execute(
            "INSERT INTO course_topics "
            "(video_id, start_para_index, end_para_index, start_time, end_time, "
            "topic_name, subtree_json) "
            "VALUES (99998, 0, 9, 0.0, 99.0, 'Topic A', '{}')"
        )
        conn.execute("PRAGMA foreign_keys = ON")
        conn.execute("COMMIT")

        resp = client.get("/api/outline/topics/99998")
        assert resp.status_code == 200
        data = resp.json()
        assert len(data) == 2
        # Should be ordered by start_para_index, so Topic A first
        assert data[0]["topic_name"] == "Topic A"
        assert data[1]["topic_name"] == "Topic B"

    finally:
        conn.execute("DELETE FROM course_topics WHERE video_id = 99998")
        conn.execute("COMMIT")
        conn.close()


# ============================================================================
# Slice 10 — PUT /api/outline/topics/{topic_id} (思维导图编辑回写)
# ============================================================================

# ---------------------------------------------------------------------------
# Helper: insert a test topic row and return its id
# ---------------------------------------------------------------------------


def _insert_test_topic(video_id=99999, topic_name="Test Topic",
                       subtree_json='{"course_name":"Test","subtree":{"children":[]}}'):
    """Insert a test row into course_topics and return its id."""
    import main_outline

    conn = main_outline.get_db(str(DB_PATH))
    try:
        conn.execute("PRAGMA foreign_keys = OFF")
        cur = conn.execute(
            "INSERT INTO course_topics "
            "(video_id, start_para_index, end_para_index, start_time, end_time, "
            "topic_name, subtree_json) "
            "VALUES (?, 0, 5, 10.0, 50.0, ?, ?)",
            (video_id, topic_name, subtree_json),
        )
        conn.execute("PRAGMA foreign_keys = ON")
        conn.execute("COMMIT")
        return cur.lastrowid
    finally:
        conn.close()


def _delete_test_topics(*ids):
    """Delete test rows from course_topics by id."""
    import main_outline

    conn = main_outline.get_db(str(DB_PATH))
    try:
        for tid in ids:
            conn.execute("DELETE FROM course_topics WHERE id = ?", (tid,))
        conn.execute("COMMIT")
    finally:
        conn.close()


# ---------------------------------------------------------------------------
# PUT /api/outline/topics/{topic_id} — success cases
# ---------------------------------------------------------------------------


def test_put_topic_update_topic_name_only(client):
    """PUT updates topic_name while leaving other columns intact."""
    tid = _insert_test_topic(video_id=99999, topic_name="Original Name")
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"topic_name": "Updated Name"},
        )
        assert resp.status_code == 200, (
            f"Expected 200, got {resp.status_code}: {resp.text}"
        )
        data = resp.json()
        assert data["topic_name"] == "Updated Name"
        assert data["id"] == tid
        assert data["video_id"] == 99999
        # Other fields unchanged
        assert data["start_para_index"] == 0
        assert data["end_para_index"] == 5

        # Verify persistence via GET
        resp2 = client.get(f"/api/outline/topics/{data['video_id']}")
        topics = resp2.json()
        assert any(t["id"] == tid and t["topic_name"] == "Updated Name"
                   for t in topics)
    finally:
        _delete_test_topics(tid)


def test_put_topic_update_subtree_json_only(client):
    """PUT updates subtree_json while leaving topic_name intact."""
    new_subtree = json.dumps({
        "course_name": "Test",
        "subtree": {"children": [{"name": "NewChild", "children": []}]},
    })
    tid = _insert_test_topic(video_id=99999, topic_name="Keep This Name")
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"subtree_json": new_subtree},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["topic_name"] == "Keep This Name"
        assert data["subtree_json"] == new_subtree
    finally:
        _delete_test_topics(tid)


def test_put_topic_update_both_fields(client):
    """PUT updates both topic_name and subtree_json simultaneously."""
    new_subtree = json.dumps({
        "course_name": "UpdatedCourse",
        "subtree": {"children": [{"name": "Child1", "children": []}]},
    })
    tid = _insert_test_topic(video_id=99999, topic_name="Old Name")
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"topic_name": "New Name", "subtree_json": new_subtree},
        )
        assert resp.status_code == 200
        data = resp.json()
        assert data["topic_name"] == "New Name"
        assert data["subtree_json"] == new_subtree
    finally:
        _delete_test_topics(tid)


def test_put_topic_response_contains_all_fields(client):
    """The PUT response contains all expected row fields."""
    tid = _insert_test_topic(video_id=99999)
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"topic_name": "All Fields Check"},
        )
        assert resp.status_code == 200
        data = resp.json()
        for key in ("id", "video_id", "start_para_index", "end_para_index",
                     "start_time", "end_time", "topic_name", "subtree_json"):
            assert key in data, f"Missing key '{key}' in response: {data}"
    finally:
        _delete_test_topics(tid)


# ---------------------------------------------------------------------------
# PUT /api/outline/topics/{topic_id} — error cases
# ---------------------------------------------------------------------------


def test_put_topic_nonexistent_id_returns_404(client):
    """PUT to a topic id that does not exist returns 404."""
    resp = client.put(
        "/api/outline/topics/9999999",
        json={"topic_name": "Ghost"},
    )
    assert resp.status_code == 404, (
        f"Expected 404, got {resp.status_code}: {resp.text}"
    )
    detail = resp.json().get("detail", "")
    assert "not found" in detail.lower()


def test_put_topic_empty_body_returns_400(client):
    """PUT with empty JSON body returns 400."""
    tid = _insert_test_topic(video_id=99999)
    try:
        resp = client.put(f"/api/outline/topics/{tid}", json={})
        assert resp.status_code == 400, (
            f"Expected 400, got {resp.status_code}: {resp.text}"
        )
    finally:
        _delete_test_topics(tid)


def test_put_topic_unexpected_field_returns_422(client):
    """PUT with a field other than topic_name or subtree_json returns 422."""
    tid = _insert_test_topic(video_id=99999)
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"video_id": 42, "topic_name": "Hack"},
        )
        assert resp.status_code == 422, (
            f"Expected 422, got {resp.status_code}: {resp.text}"
        )
    finally:
        _delete_test_topics(tid)


def test_put_topic_idempotent_update(client):
    """PUT with the same values is idempotent (no error on repeated calls)."""
    tid = _insert_test_topic(video_id=99999, topic_name="Idem")
    try:
        payload = {"topic_name": "Idem2"}
        # First update
        resp1 = client.put(f"/api/outline/topics/{tid}", json=payload)
        assert resp1.status_code == 200
        # Second update with same payload
        resp2 = client.put(f"/api/outline/topics/{tid}", json=payload)
        assert resp2.status_code == 200
        assert resp2.json()["topic_name"] == "Idem2"
    finally:
        _delete_test_topics(tid)


def test_put_topic_negative_id_returns_404(client):
    """PUT /api/outline/topics/-1 returns 404."""
    resp = client.put("/api/outline/topics/-1",
                      json={"topic_name": "Negative"})
    assert resp.status_code == 404


def test_put_topic_zero_id_returns_404(client):
    """PUT /api/outline/topics/0 returns 404."""
    resp = client.put("/api/outline/topics/0",
                      json={"topic_name": "Zero"})
    assert resp.status_code == 404


def test_put_topic_non_integer_id_returns_422(client):
    """PUT /api/outline/topics/abc returns 422 (FastAPI path validation)."""
    resp = client.put("/api/outline/topics/abc",
                      json={"topic_name": "String"})
    assert resp.status_code == 422


# ---------------------------------------------------------------------------
# PUT with subtree_json containing special characters / unicode
# ---------------------------------------------------------------------------


def test_put_topic_unicode_subtree_json(client):
    """PUT accepts subtree_json with CJK characters."""
    subtree = json.dumps({
        "course_name": "计算机网络基础",
        "subtree": {
            "children": [
                {"name": "UDP概述", "children": [
                    {"name": "无连接传输", "children": []},
                ]},
            ],
        },
    }, ensure_ascii=False)
    tid = _insert_test_topic(video_id=99999)
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"subtree_json": subtree},
        )
        assert resp.status_code == 200
        data = resp.json()
        parsed = json.loads(data["subtree_json"])
        assert parsed["course_name"] == "计算机网络基础"
        assert parsed["subtree"]["children"][0]["name"] == "UDP概述"
        assert parsed["subtree"]["children"][0]["children"][0]["name"] == "无连接传输"
    finally:
        _delete_test_topics(tid)


def test_put_topic_subtree_json_with_special_chars(client):
    """PUT handles subtree_json with quotes, newlines, and special characters."""
    subtree = json.dumps({
        "course_name": "Test's \"Special\" Course",
        "subtree": {
            "children": [
                {"name": "Line1\nLine2", "children": []},
                {"name": "Backslash \\ test", "children": []},
            ],
        },
    })
    tid = _insert_test_topic(video_id=99999)
    try:
        resp = client.put(
            f"/api/outline/topics/{tid}",
            json={"subtree_json": subtree},
        )
        assert resp.status_code == 200
        data = resp.json()
        parsed = json.loads(data["subtree_json"])
        assert parsed["course_name"] == "Test's \"Special\" Course"
        assert parsed["subtree"]["children"][0]["name"] == "Line1\nLine2"
        assert parsed["subtree"]["children"][1]["name"] == "Backslash \\ test"
    finally:
        _delete_test_topics(tid)


# ============================================================================
# CKG — concept extraction end-to-end (issue 01)
# ============================================================================

# A real video id that exists in corpus.db (FK target for course_ckg writes).
CKG_TEST_VIDEO_ID = 132

VALID_CONCEPTS_RESPONSE = (
    '{"concepts": ['
    '{"name": "Gradient Descent", "definition": "An iterative optimization '
    'algorithm to minimize a cost function.", "first_para": 3}, '
    '{"name": "Learning Rate", "definition": "The step size controlling how '
    'far each gradient descent update moves.", "first_para": 5}, '
    '{"name": "Cost Function", "definition": "A function measuring prediction '
    'error to be minimized.", "first_para": 2}'
    ']}'
)

# Same concept name appearing twice (different case/spacing) -> must de-dup.
DUPLICATE_CONCEPTS_RESPONSE = (
    '{"concepts": ['
    '{"name": "Gradient Descent", "definition": "first", "first_para": 3}, '
    '{"name": "gradient descent", "definition": "dup lower", "first_para": 9}, '
    '{"name": "Learning Rate", "definition": "lr", "first_para": 5}'
    ']}'
)

# Decomposition tree (issue 01): title root → concepts each carrying a parent.
# "Supervised Learning" is carved from the title; "Classification"/"Regression"
# are carved from "Supervised Learning".
DECOMP_CONCEPTS_RESPONSE = (
    '{"title": "What is Machine Learning", "concepts": ['
    '{"name": "Supervised Learning", "definition": "learning from labeled data", '
    '"first_para": 2, "parent": "What is Machine Learning"}, '
    '{"name": "Classification", "definition": "predicting a discrete label", '
    '"first_para": 4, "parent": "Supervised Learning"}, '
    '{"name": "Regression", "definition": "predicting a continuous value", '
    '"first_para": 6, "parent": "supervised learning"}'
    ']}'
)

# Grounded cross-relations (讲述关联, issue 02). Endpoints reference the DECOMP
# concept names; `type` is free-labelled by the LLM (here arbitrary strings).
VALID_RELATIONS_RESPONSE = (
    '{"relations": ['
    '{"from": "Supervised Learning", "to": "Classification", '
    '"type": "explains-via", "evidence": "[段落 4]"}, '
    '{"from": "Classification", "to": "Regression", '
    '"type": "contrasts-with", "evidence": "para 6"}'
    ']}'
)


def _cleanup_ckg(video_id):
    import main_outline
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute("DELETE FROM course_ckg WHERE video_id = ?", (video_id,))
    conn.execute("COMMIT")
    conn.close()


# ---------------------------------------------------------------------------
# Schema + config seed
# ---------------------------------------------------------------------------

def test_course_ckg_table_exists_with_columns():
    """course_ckg is created by init_db with the expected columns."""
    from main_outline import init_db
    init_db(str(DB_PATH))

    conn = sqlite3.connect(str(DB_PATH))
    tables = {
        row[0]
        for row in conn.execute(
            "SELECT name FROM sqlite_master WHERE type='table'"
        ).fetchall()
    }
    assert "course_ckg" in tables

    cols = {row[1] for row in conn.execute("PRAGMA table_info(course_ckg)")}
    conn.close()
    for required in ("video_id", "graph_json", "model", "created_at"):
        assert required in cols, f"course_ckg missing column {required}"


def test_ck_prompt_concepts_seeded():
    """config key ck_prompt_concepts is seeded with the English Step-1 prompt."""
    from main_outline import init_db
    init_db(str(DB_PATH))

    conn = sqlite3.connect(str(DB_PATH))
    row = conn.execute(
        "SELECT value FROM config WHERE key = 'ck_prompt_concepts'"
    ).fetchone()
    conn.close()
    assert row is not None, "ck_prompt_concepts not seeded"
    assert "TEACHING CONCEPT" in row[0]
    assert "first_para" in row[0]
    # New decomposition-tree prompt (issue 01): must ask for a parent field.
    assert '"parent"' in row[0]
    assert "DECOMPOS" in row[0].upper()


def test_ck_prompt_concepts_seed_idempotent():
    """Running init_db twice keeps exactly one ck_prompt_concepts row."""
    from main_outline import init_db
    init_db(str(DB_PATH))
    init_db(str(DB_PATH))

    conn = sqlite3.connect(str(DB_PATH))
    count = conn.execute(
        "SELECT COUNT(*) FROM config WHERE key = 'ck_prompt_concepts'"
    ).fetchone()[0]
    conn.close()
    assert count == 1


# ---------------------------------------------------------------------------
# parse_concepts_response - unit tests
# ---------------------------------------------------------------------------

def test_parse_concepts_valid():
    from main_outline import parse_concepts_response
    result = parse_concepts_response(VALID_CONCEPTS_RESPONSE)
    assert result is not None
    assert len(result["concepts"]) == 3
    first = result["concepts"][0]
    assert first["name"] == "Gradient Descent"
    assert first["definition"].startswith("An iterative")
    assert first["first_para"] == 3


def test_parse_concepts_dedup_by_name():
    from main_outline import parse_concepts_response
    result = parse_concepts_response(DUPLICATE_CONCEPTS_RESPONSE)
    assert result is not None
    names = [c["name"] for c in result["concepts"]]
    assert len(names) == 2, f"Expected dedup to 2 concepts, got {names}"
    assert "Gradient Descent" in names
    assert "Learning Rate" in names


def test_parse_concepts_strips_markdown_fences():
    from main_outline import parse_concepts_response
    fenced = "```json\n" + VALID_CONCEPTS_RESPONSE + "\n```"
    result = parse_concepts_response(fenced)
    assert result is not None
    assert len(result["concepts"]) == 3


def test_parse_concepts_strips_leading_text():
    from main_outline import parse_concepts_response
    messy = "Sure, here you go:\n" + VALID_CONCEPTS_RESPONSE + "\nDone."
    result = parse_concepts_response(messy)
    assert result is not None
    assert len(result["concepts"]) == 3


def test_parse_concepts_invalid_json_returns_none():
    from main_outline import parse_concepts_response
    assert parse_concepts_response("not json at all") is None


def test_parse_concepts_empty_string_returns_none():
    from main_outline import parse_concepts_response
    assert parse_concepts_response("") is None


def test_parse_concepts_wrong_structure_returns_none():
    from main_outline import parse_concepts_response
    assert parse_concepts_response('{"topics": []}') is None


def test_parse_concepts_empty_list_returns_none():
    from main_outline import parse_concepts_response
    assert parse_concepts_response('{"concepts": []}') is None


def test_parse_concepts_captures_parent_and_title():
    """Decomposition prompt output: each concept carries parent; title passes through."""
    from main_outline import parse_concepts_response
    result = parse_concepts_response(DECOMP_CONCEPTS_RESPONSE)
    assert result is not None
    assert result["title"] == "What is Machine Learning"
    by_name = {c["name"]: c for c in result["concepts"]}
    assert by_name["Supervised Learning"]["parent"] == "What is Machine Learning"
    assert by_name["Classification"]["parent"] == "Supervised Learning"


def test_parse_concepts_parent_defaults_empty():
    """Concepts with no parent field default parent to '' (legacy/title roots)."""
    from main_outline import parse_concepts_response
    result = parse_concepts_response(VALID_CONCEPTS_RESPONSE)
    assert result is not None
    assert all(c["parent"] == "" for c in result["concepts"])


# ---------------------------------------------------------------------------
# derive_decomposition_edges — backbone from parent pointers
# ---------------------------------------------------------------------------

def test_derive_decomposition_edges_parent_to_child():
    from main_outline import parse_concepts_response, derive_decomposition_edges
    concepts = parse_concepts_response(DECOMP_CONCEPTS_RESPONSE)["concepts"]
    edges = derive_decomposition_edges(concepts)
    # Title is not a concept → "Supervised Learning" is a root (no incoming edge).
    # Its two children resolve (case-insensitively) to it.
    assert {"from": "Supervised Learning", "to": "Classification"} in edges
    assert {"from": "Supervised Learning", "to": "Regression"} in edges
    assert len(edges) == 2


def test_derive_decomposition_edges_drops_unknown_and_self():
    from main_outline import derive_decomposition_edges
    concepts = [
        {"name": "A", "parent": "Nonexistent"},   # unknown parent → root
        {"name": "B", "parent": "A"},              # resolves
        {"name": "C", "parent": "C"},              # self-loop → dropped
    ]
    edges = derive_decomposition_edges(concepts)
    assert edges == [{"from": "A", "to": "B"}]


# ---------------------------------------------------------------------------
# Worker: persistence to course_ckg
# ---------------------------------------------------------------------------

def test_ckg_extraction_persists_dedup_concepts():
    """_run_ckg_extraction writes one course_ckg row with de-duplicated
    concepts (each carrying name/definition/first_para)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i} content"}
        for i in range(1, 11)
    ]
    mock_response = {"response": DUPLICATE_CONCEPTS_RESPONSE}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT video_id, graph_json, model, created_at "
            "FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        assert row is not None, "no course_ckg row written"
        graph = json.loads(row["graph_json"])
        concepts = graph["concepts"]
        assert len(concepts) == 2, f"expected dedup to 2, got {len(concepts)}"
        for c in concepts:
            assert "name" in c and "definition" in c and "first_para" in c
        assert row["model"]
        assert row["created_at"]
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_extraction_parse_failure_no_write_no_crash():
    """A malformed LLM response must not write to course_ckg nor crash."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 1.0, "text": "x"},
    ]
    mock_response = {"response": "I cannot help with that."}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT 1 FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()
        assert row is None, "parse failure should not write a course_ckg row"
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_cancel_flag_skips_ollama():
    """When the cancel flag is set, the worker does not call Ollama."""
    import main_outline

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 1.0, "text": "x"},
    ]
    main_outline._ckg_cancel_flags[vid] = True

    mock_ollama = MagicMock()
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", mock_ollama):
                main_outline._run_ckg_extraction(vid)
        mock_ollama.assert_not_called()
    finally:
        main_outline._ckg_cancel_flags.pop(vid, None)


# ---------------------------------------------------------------------------
# Trigger + read API
# ---------------------------------------------------------------------------

def test_extract_ckg_endpoint_returns_202(client):
    """POST /api/outline/ckg/{video_id} returns 202 and eventually persists."""
    import main_outline
    import time

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]
    mock_response = {"response": VALID_CONCEPTS_RESPONSE}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                resp = client.post(f"/api/outline/ckg/{vid}")
                assert resp.status_code == 202, resp.text
                assert resp.json()["status"] == "started"

                deadline = time.time() + 5.0
                row = None
                while time.time() < deadline:
                    conn = main_outline.get_db(str(DB_PATH))
                    row = conn.execute(
                        "SELECT graph_json FROM course_ckg WHERE video_id = ?",
                        (vid,),
                    ).fetchone()
                    conn.close()
                    if row is not None:
                        break
                    time.sleep(0.1)
                assert row is not None, "extraction did not persist within 5s"
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_extract_ckg_nonexistent_video_returns_404(client):
    resp = client.post("/api/outline/ckg/99999999")
    assert resp.status_code == 404


def test_get_ckg_returns_concepts(client):
    """GET /api/outline/ckg/{video_id} returns stored concepts."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    graph_json = json.dumps({"concepts": [
        {"name": "A", "definition": "def a", "first_para": 1},
        {"name": "B", "definition": "def b", "first_para": 4},
    ]}, ensure_ascii=False)

    _cleanup_ckg(vid)
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at) "
        "VALUES (?, ?, ?, ?)",
        (vid, graph_json, "test-model", "2026-06-25T00:00:00+00:00"),
    )
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get(f"/api/outline/ckg/{vid}")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        assert data["model"] == "test-model"
        assert len(data["concepts"]) == 2
        assert data["concepts"][0]["name"] == "A"
        assert data["concepts"][1]["first_para"] == 4
    finally:
        _cleanup_ckg(vid)


def test_get_ckg_missing_returns_404(client):
    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    resp = client.get(f"/api/outline/ckg/{vid}")
    assert resp.status_code == 404


# ---------------------------------------------------------------------------
# Front-end i18n: CKG analyze button labels
# ---------------------------------------------------------------------------

def test_i18n_has_ckg_concept_labels():
    """outline.html defines i18n keys for the CKG concept panel/button."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.analyze_btn", "ckg.concepts_title",
                "ckg.concepts_empty", "ckg.first_para"):
        assert (f"'{key}'" in content or f'"{key}"' in content), \
            f"Missing i18n key {key}"


# ============================================================================
# CKG — prerequisite edges + single-lecture DAG view (issue 02)
# ============================================================================

VALID_EDGES_RESPONSE = (
    '{"edges": ['
    '{"from": "Cost Function", "to": "Gradient Descent"}, '
    '{"from": "Gradient Descent", "to": "Learning Rate"}'
    ']}'
)


# ---------------------------------------------------------------------------
# ck_prompt_edges is RETIRED for CKG extraction (ADR 0003): the CKG backbone is
# now the title→concept decomposition tree. parse_edges_response/validate_edges/
# break_cycles remain (used by lesson generation), so their unit tests stay.
# ---------------------------------------------------------------------------


def test_ck_prompt_edges_not_exposed_by_prompt_editor(client):
    """The retired ck_prompt_edges key is rejected by the CK prompt editor."""
    from main_outline import init_db
    init_db(str(DB_PATH))
    resp = client.put(
        "/api/outline/ckg/prompts",
        json={"ck_prompt_edges": "anything"},
    )
    assert resp.status_code == 422


# ---------------------------------------------------------------------------
# parse_edges_response — unit tests
# ---------------------------------------------------------------------------

def test_parse_edges_valid():
    from main_outline import parse_edges_response
    result = parse_edges_response(VALID_EDGES_RESPONSE)
    assert result is not None
    assert len(result["edges"]) == 2
    assert result["edges"][0] == {"from": "Cost Function", "to": "Gradient Descent"}


def test_parse_edges_strips_markdown_fences():
    from main_outline import parse_edges_response
    fenced = "```json\n" + VALID_EDGES_RESPONSE + "\n```"
    result = parse_edges_response(fenced)
    assert result is not None
    assert len(result["edges"]) == 2


def test_parse_edges_strips_leading_text():
    from main_outline import parse_edges_response
    messy = "Here are the edges:\n" + VALID_EDGES_RESPONSE + "\nThanks."
    result = parse_edges_response(messy)
    assert result is not None
    assert len(result["edges"]) == 2


def test_parse_edges_invalid_json_returns_none():
    from main_outline import parse_edges_response
    assert parse_edges_response("no json here") is None


def test_parse_edges_empty_string_returns_none():
    from main_outline import parse_edges_response
    assert parse_edges_response("") is None


def test_parse_edges_wrong_structure_returns_none():
    from main_outline import parse_edges_response
    assert parse_edges_response('{"concepts": []}') is None


def test_parse_edges_empty_list_returns_empty_edges():
    """An empty edge list is valid (a lecture may legitimately have no edges)."""
    from main_outline import parse_edges_response
    result = parse_edges_response('{"edges": []}')
    assert result == {"edges": []}


def test_parse_edges_skips_malformed_edges():
    from main_outline import parse_edges_response
    raw = ('{"edges": ['
           '{"from": "A", "to": "B"}, '
           '{"from": "", "to": "C"}, '
           '{"to": "D"}, '
           '"not a dict"'
           ']}')
    result = parse_edges_response(raw)
    assert result == {"edges": [{"from": "A", "to": "B"}]}


# ---------------------------------------------------------------------------
# Edge validation — drop edges referencing unknown concepts
# ---------------------------------------------------------------------------

def test_validate_edges_drops_unknown_concepts():
    from main_outline import validate_edges
    concepts = [
        {"name": "Cost Function", "definition": "", "first_para": 1},
        {"name": "Gradient Descent", "definition": "", "first_para": 2},
    ]
    edges = [
        {"from": "Cost Function", "to": "Gradient Descent"},  # keep
        {"from": "Cost Function", "to": "Backpropagation"},   # drop: unknown to
        {"from": "Mystery", "to": "Gradient Descent"},        # drop: unknown from
    ]
    valid = validate_edges(edges, concepts)
    assert valid == [{"from": "Cost Function", "to": "Gradient Descent"}]


def test_validate_edges_matches_by_normalized_name():
    from main_outline import validate_edges
    concepts = [
        {"name": "Cost Function", "definition": "", "first_para": 1},
        {"name": "Gradient Descent", "definition": "", "first_para": 2},
    ]
    edges = [{"from": "cost function", "to": "  GRADIENT DESCENT "}]
    valid = validate_edges(edges, concepts)
    # Normalized to the canonical concept names.
    assert valid == [{"from": "Cost Function", "to": "Gradient Descent"}]


def test_validate_edges_drops_self_loops():
    from main_outline import validate_edges
    concepts = [{"name": "A", "definition": "", "first_para": 1}]
    edges = [{"from": "A", "to": "A"}]
    assert validate_edges(edges, concepts) == []


# ---------------------------------------------------------------------------
# break_cycles — code-guaranteed DAG
# ---------------------------------------------------------------------------

def test_break_cycles_keeps_acyclic_graph_intact():
    from main_outline import break_cycles
    edges = [
        {"from": "A", "to": "B"},
        {"from": "B", "to": "C"},
        {"from": "A", "to": "C"},
    ]
    result = break_cycles(edges)
    assert result == edges  # already a DAG, nothing dropped


def test_break_cycles_breaks_a_cycle():
    from main_outline import break_cycles
    # A -> B -> C -> A forms a cycle; the closing edge must be dropped.
    edges = [
        {"from": "A", "to": "B"},
        {"from": "B", "to": "C"},
        {"from": "C", "to": "A"},
    ]
    result = break_cycles(edges)
    assert len(result) == 2
    assert _is_dag(result)
    # The first two edges (added before the cycle closes) survive.
    assert {"from": "A", "to": "B"} in result
    assert {"from": "B", "to": "C"} in result


def test_break_cycles_breaks_two_node_cycle():
    from main_outline import break_cycles
    edges = [{"from": "A", "to": "B"}, {"from": "B", "to": "A"}]
    result = break_cycles(edges)
    assert result == [{"from": "A", "to": "B"}]
    assert _is_dag(result)


def test_break_cycles_complex_graph_is_dag():
    from main_outline import break_cycles
    edges = [
        {"from": "A", "to": "B"},
        {"from": "B", "to": "C"},
        {"from": "C", "to": "D"},
        {"from": "D", "to": "B"},  # closes B->C->D->B
        {"from": "A", "to": "D"},
    ]
    result = break_cycles(edges)
    assert _is_dag(result)
    assert {"from": "D", "to": "B"} not in result


def _is_dag(edges):
    """Helper: True if the directed graph defined by edges is acyclic."""
    from collections import defaultdict
    adj = defaultdict(list)
    nodes = set()
    for e in edges:
        adj[e["from"]].append(e["to"])
        nodes.add(e["from"])
        nodes.add(e["to"])
    WHITE, GREY, BLACK = 0, 1, 2
    color = {n: WHITE for n in nodes}

    def visit(n):
        color[n] = GREY
        for m in adj[n]:
            if color[m] == GREY:
                return False
            if color[m] == WHITE and not visit(m):
                return False
        color[n] = BLACK
        return True

    return all(color[n] != WHITE or visit(n) for n in nodes)


# ---------------------------------------------------------------------------
# ck_prompt_relations seed + parse_relations_response + validate_relations
# ---------------------------------------------------------------------------

def test_ck_prompt_relations_seeded():
    """config key ck_prompt_relations is seeded with the cross-relation prompt."""
    from main_outline import init_db
    init_db(str(DB_PATH))
    conn = sqlite3.connect(str(DB_PATH))
    row = conn.execute(
        "SELECT value FROM config WHERE key = 'ck_prompt_relations'"
    ).fetchone()
    conn.close()
    assert row is not None, "ck_prompt_relations not seeded"
    # Cross-branch, free-labelled type, from/to/type output.
    assert "CROSS-RELATION" in row[0].upper() or "CROSS-BRANCH" in row[0].upper()
    assert '"from"' in row[0] and '"to"' in row[0] and '"type"' in row[0]
    # Must signal the type is NOT a fixed taxonomy.
    assert "NOT limited" in row[0] or "your own words" in row[0].lower()


def test_ck_prompt_relations_seed_idempotent():
    from main_outline import init_db
    init_db(str(DB_PATH))
    init_db(str(DB_PATH))
    conn = sqlite3.connect(str(DB_PATH))
    count = conn.execute(
        "SELECT COUNT(*) FROM config WHERE key = 'ck_prompt_relations'"
    ).fetchone()[0]
    conn.close()
    assert count == 1


def test_parse_relations_valid():
    from main_outline import parse_relations_response
    result = parse_relations_response(VALID_RELATIONS_RESPONSE)
    assert result is not None
    assert len(result["relations"]) == 2
    first = result["relations"][0]
    assert first["from"] == "Supervised Learning"
    assert first["to"] == "Classification"
    assert first["type"] == "explains-via"
    assert first["evidence"] == "[段落 4]"


def test_parse_relations_freeform_type_preserved():
    """An arbitrary, non-anchor type label is stored verbatim (emergent)."""
    from main_outline import parse_relations_response
    raw = ('{"relations": [{"from": "A", "to": "B", '
           '"type": "draws-an-analogy-between", "evidence": ""}]}')
    result = parse_relations_response(raw)
    assert result["relations"][0]["type"] == "draws-an-analogy-between"


def test_parse_relations_missing_type_defaults_empty():
    from main_outline import parse_relations_response
    result = parse_relations_response('{"relations": [{"from": "A", "to": "B"}]}')
    assert result["relations"][0]["type"] == ""
    assert result["relations"][0]["evidence"] == ""


def test_parse_relations_strips_markdown_fences():
    from main_outline import parse_relations_response
    fenced = "```json\n" + VALID_RELATIONS_RESPONSE + "\n```"
    result = parse_relations_response(fenced)
    assert result is not None and len(result["relations"]) == 2


def test_parse_relations_empty_list_valid():
    from main_outline import parse_relations_response
    assert parse_relations_response('{"relations": []}') == {"relations": []}


def test_parse_relations_invalid_returns_none():
    from main_outline import parse_relations_response
    assert parse_relations_response("no json here") is None
    assert parse_relations_response("") is None
    assert parse_relations_response('{"edges": []}') is None


def test_validate_relations_drops_unknown_and_self_preserves_type():
    from main_outline import validate_relations
    concepts = [{"name": "Gradient Descent"}, {"name": "Cost Function"}]
    relations = [
        {"from": "cost function", "to": " GRADIENT DESCENT ", "type": "builds-on", "evidence": "x"},
        {"from": "Cost Function", "to": "Nonexistent", "type": "motivates"},  # unknown → drop
        {"from": "A", "to": "A", "type": "self"},                              # self → drop
    ]
    valid = validate_relations(relations, concepts)
    assert valid == [
        {"from": "Cost Function", "to": "Gradient Descent", "type": "builds-on", "evidence": "x"}
    ]


# ---------------------------------------------------------------------------
# Worker: single decomposition call — graph_json carries the new structure
# ---------------------------------------------------------------------------

def test_ckg_extraction_persists_decomposition_graph():
    """_run_ckg_extraction writes schema_version-2 graph_json: concepts (with
    parent), decomposition_edges derived from parents, empty relations."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 11)
    ]
    # Single LLM call now (decomposition). Must accept temperature + kwargs.
    mock_response = {"response": DECOMP_CONCEPTS_RESPONSE}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT graph_json FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        assert row is not None, "no course_ckg row written"
        graph = json.loads(row["graph_json"])
        assert graph["schema_version"] == 2
        assert graph["title"] == "What is Machine Learning"
        assert "decomposition_edges" in graph and "relations" in graph
        assert graph["relations"] == []
        assert len(graph["concepts"]) == 3
        assert all("parent" in c and "first_para" in c for c in graph["concepts"])
        names = {c["name"] for c in graph["concepts"]}
        decomp = graph["decomposition_edges"]
        assert len(decomp) == 2
        for e in decomp:
            assert e["from"] in names and e["to"] in names
        assert _is_dag(decomp)
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_extraction_unresolved_parents_persist_concepts_no_edges():
    """Concepts whose parents don't resolve to a known concept (e.g. all point
    at the title) still persist, with an empty decomposition backbone."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]
    # All three concepts are carved directly from the title → no parent resolves
    # to another concept → no decomposition edges, but concepts are kept.
    flat_response = {"response": (
        '{"title": "T", "concepts": ['
        '{"name": "A", "definition": "a", "first_para": 1, "parent": "T"}, '
        '{"name": "B", "definition": "b", "first_para": 2, "parent": "T"}, '
        '{"name": "C", "definition": "c", "first_para": 3, "parent": "T"}'
        ']}'
    )}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=flat_response):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT graph_json FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        assert row is not None
        graph = json.loads(row["graph_json"])
        assert len(graph["concepts"]) == 3
        assert graph["decomposition_edges"] == []
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_extraction_persists_grounded_relations():
    """The 2nd call extracts grounded cross-relations, validated against the
    concept set, and persisted into graph_json['relations']."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 11)
    ]
    responses = [
        {"response": DECOMP_CONCEPTS_RESPONSE},
        {"response": VALID_RELATIONS_RESPONSE},
    ]

    def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
        return responses.pop(0)

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", side_effect=mock_ollama):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT graph_json FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        graph = json.loads(row["graph_json"])
        names = {c["name"] for c in graph["concepts"]}
        rels = graph["relations"]
        assert len(rels) == 2
        for r in rels:
            assert r["from"] in names and r["to"] in names
            assert "type" in r
        # free-labelled types preserved verbatim
        types = {r["type"] for r in rels}
        assert "explains-via" in types and "contrasts-with" in types
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_extraction_bad_relations_still_persists_backbone():
    """If the relation step fails to parse, concepts + decomposition still
    persist with an empty relations list (no data loss)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]
    responses = [
        {"response": DECOMP_CONCEPTS_RESPONSE},
        {"response": "I cannot produce relations."},
    ]

    def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
        return responses.pop(0)

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", side_effect=mock_ollama):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT graph_json FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        graph = json.loads(row["graph_json"])
        assert len(graph["concepts"]) == 3
        assert len(graph["decomposition_edges"]) == 2
        assert graph["relations"] == []
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


# ---------------------------------------------------------------------------
# GET endpoint returns the decomposition structure (+ back-compat edges alias)
# ---------------------------------------------------------------------------

def test_get_ckg_returns_decomposition(client):
    """GET returns concepts, decomposition_edges, relations, and an edges alias."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    graph_json = json.dumps({
        "schema_version": 2,
        "title": "T",
        "concepts": [
            {"name": "A", "definition": "def a", "first_para": 1, "parent": "T"},
            {"name": "B", "definition": "def b", "first_para": 4, "parent": "A"},
        ],
        "decomposition_edges": [{"from": "A", "to": "B"}],
        "relations": [],
    }, ensure_ascii=False)

    _cleanup_ckg(vid)
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at) "
        "VALUES (?, ?, ?, ?)",
        (vid, graph_json, "test-model", "2026-06-25T00:00:00+00:00"),
    )
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get(f"/api/outline/ckg/{vid}")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        assert len(data["concepts"]) == 2
        assert "edges" in data
        assert data["edges"] == [{"from": "A", "to": "B"}]
    finally:
        _cleanup_ckg(vid)


def test_get_ckg_edges_defaults_to_empty_when_absent(client):
    """Legacy rows (concepts only, no edges key) still return edges: []."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    graph_json = json.dumps({
        "concepts": [{"name": "A", "definition": "a", "first_para": 1}],
    }, ensure_ascii=False)

    _cleanup_ckg(vid)
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at) "
        "VALUES (?, ?, ?, ?)",
        (vid, graph_json, "test-model", "2026-06-25T00:00:00+00:00"),
    )
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get(f"/api/outline/ckg/{vid}")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        # Legacy row: no edges / decomposition_edges / relations / schema_version.
        assert data["schema_version"] == 1
        assert len(data["concepts"]) == 1
        assert data["decomposition_edges"] == []
        assert data["relations"] == []
        assert data["edges"] == []
    finally:
        _cleanup_ckg(vid)


# ---------------------------------------------------------------------------
# Front-end: DAG view (ECharts graph) + i18n
# ---------------------------------------------------------------------------

def test_i18n_has_ckg_dag_labels():
    """outline.html defines i18n keys for the DAG view across all 3 languages."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.dag_title", "ckg.edges_count", "ckg.view_list",
                "ckg.view_dag"):
        # Must appear at least 3 times (zh / en / ja i18n blocks).
        assert content.count(f"'{key}'") >= 3, f"i18n key {key} not in all 3 langs"


def test_frontend_renders_dag_with_echarts_graph():
    """The DAG render uses ECharts graph type with layout:'none', explicit x
    coords (first_para-based horizontal layout) and arrow edge symbols."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "renderCKGGraph" in content, "missing DAG render function"
    assert "type: 'graph'" in content or 'type: "graph"' in content
    assert "layout: 'none'" in content or "layout:'none'" in content
    # Arrow on the target end of each edge.
    assert "edgeSymbol" in content and "arrow" in content
    # first_para drives the x coordinate mapping.
    assert "first_para" in content


def test_frontend_renders_grounded_relations():
    """renderCKGGraph takes a relations arg and draws them as a distinct second
    link set (own style + a type label), separate from decomposition edges."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # render fn accepts relations; data carries relations; links tag their kind.
    assert "renderCKGGraph(concepts, edges, relations)" in content \
        or "function renderCKGGraph(concepts, edges, relations)" in content
    assert "_ckgData.relations" in content, "relations not threaded into render"
    assert "data.relations" in content, "relations not read from API payload"
    assert "_kind: 'relation'" in content, "relation links not distinguished"
    # Relation links show their free-labelled type as the edge label.
    assert "r.type" in content


def test_i18n_has_ckg_relations_labels():
    """relations_count + prompt_relations_label exist in all 3 languages."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.relations_count", "ckg.prompt_relations_label"):
        assert content.count(f"'{key}'") >= 3, f"i18n key {key} not in all 3 langs"


# ============================================================================
# CKG — directed-graph topology params + delivery direction (issue 03)
# ============================================================================

def _concepts(*names_with_fp):
    """Build concept dicts. Each arg is (name, first_para)."""
    return [
        {"name": n, "definition": f"def {n}", "first_para": fp}
        for n, fp in names_with_fp
    ]


def _edge(a, b):
    return {"from": a, "to": b}


def test_dag_topology_linear_chain():
    """Linear chain A->B->C->D: depth=3 (longest path edge count),
    branch_factor=1 (every non-sink out-degree 1), convergence_count=0,
    density = 3 / (4*3) = 0.25."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3), ("D", 4))
    edges = [_edge("A", "B"), _edge("B", "C"), _edge("C", "D")]
    t = compute_dag_topology(concepts, edges)
    assert t["depth"] == 3
    assert t["branch_factor"] == 1.0
    assert t["convergence_count"] == 0
    assert abs(t["density"] - 0.25) < 1e-9


def test_dag_topology_flat_divergent():
    """Flat star A->B, A->C, A->D: depth=1 (small), branch_factor=3 (large),
    convergence_count=0, density = 3 / (4*3) = 0.25."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3), ("D", 4))
    edges = [_edge("A", "B"), _edge("A", "C"), _edge("A", "D")]
    t = compute_dag_topology(concepts, edges)
    assert t["depth"] == 1
    assert t["branch_factor"] == 3.0
    assert t["convergence_count"] == 0
    assert abs(t["density"] - 0.25) < 1e-9


def test_dag_topology_convergence_diamond():
    """Diamond A->B, A->C, B->D, C->D. D in-degree 2 -> convergence_count=1.
    Non-sinks A,B,C out-degrees 2,1,1 -> branch_factor = 4/3. depth=2."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3), ("D", 4))
    edges = [_edge("A", "B"), _edge("A", "C"), _edge("B", "D"), _edge("C", "D")]
    t = compute_dag_topology(concepts, edges)
    assert t["depth"] == 2
    assert t["convergence_count"] == 1
    assert abs(t["branch_factor"] - (4 / 3)) < 1e-9


def test_dag_topology_empty_graph():
    """No nodes -> all params 0, no crash."""
    from main_outline import compute_dag_topology
    t = compute_dag_topology([], [])
    assert t["depth"] == 0
    assert t["branch_factor"] == 0.0
    assert t["convergence_count"] == 0
    assert t["density"] == 0.0
    assert t["avg_path_length"] == 0.0
    assert t["clustering"] == 0.0


def test_dag_topology_no_edges():
    """Nodes but no edges -> depth 0, branch_factor 0, convergence 0, density 0."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2))
    t = compute_dag_topology(concepts, [])
    assert t["depth"] == 0
    assert t["branch_factor"] == 0.0
    assert t["convergence_count"] == 0
    assert t["density"] == 0.0
    assert t["avg_path_length"] == 0.0


def test_dag_topology_avg_path_on_undirected_projection():
    """Chain A-B-C undirected: 6 ordered pairs, distances sum=8, avg=8/6."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3))
    edges = [_edge("A", "B"), _edge("B", "C")]
    t = compute_dag_topology(concepts, edges)
    assert abs(t["avg_path_length"] - (8 / 6)) < 1e-9


def test_dag_topology_clustering_triangle():
    """Undirected triangle (A->B, B->C, A->C) has clustering=1.0."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3))
    edges = [_edge("A", "B"), _edge("B", "C"), _edge("A", "C")]
    t = compute_dag_topology(concepts, edges)
    assert abs(t["clustering"] - 1.0) < 1e-9


def test_dag_topology_clustering_zero_on_tree():
    """A tree (no triangles) has clustering=0; UI shows it honestly."""
    from main_outline import compute_dag_topology
    concepts = _concepts(("A", 1), ("B", 2), ("C", 3), ("D", 4))
    edges = [_edge("A", "B"), _edge("A", "C"), _edge("B", "D")]
    t = compute_dag_topology(concepts, edges)
    assert t["clustering"] == 0.0


def test_bottomup_ratio_all_bottomup():
    """Every edge teaches prerequisite first -> ratio 1.0."""
    from main_outline import compute_bottomup_ratio
    concepts = _concepts(("A", 1), ("B", 5), ("C", 9))
    edges = [_edge("A", "B"), _edge("B", "C")]
    assert compute_bottomup_ratio(concepts, edges) == 1.0


def test_bottomup_ratio_mixed_and_equal_excluded():
    """3 edges: bottom-up (1<5), top-down (9>5), equal (3==3) excluded -> 0.5."""
    from main_outline import compute_bottomup_ratio
    concepts = _concepts(("A", 1), ("B", 5), ("C", 9), ("D", 3), ("E", 3))
    edges = [_edge("A", "B"), _edge("C", "B"), _edge("D", "E")]
    assert compute_bottomup_ratio(concepts, edges) == 0.5


def test_bottomup_ratio_none_when_no_valid_edges():
    """No edges, or only equal-first_para edges -> None (UI shows dash)."""
    from main_outline import compute_bottomup_ratio
    assert compute_bottomup_ratio(_concepts(("A", 1)), []) is None
    eq = _concepts(("A", 3), ("B", 3))
    assert compute_bottomup_ratio(eq, [_edge("A", "B")]) is None


def test_compute_relation_density():
    from main_outline import compute_relation_density
    concepts = [{"name": "A"}, {"name": "B"}, {"name": "C"}, {"name": "D"}]
    relations = [{"from": "A", "to": "B"}, {"from": "C", "to": "D"}]
    assert compute_relation_density(concepts, relations) == 2 / 4
    assert compute_relation_density(concepts, []) == 0.0
    assert compute_relation_density([], relations) == 0.0


def test_compute_convergence_counts_relations():
    """A concept reached by its decomposition parent AND a relation has
    in-degree 2 → it is a convergence point (pure trees yield 0)."""
    from main_outline import compute_convergence_count
    concepts = [{"name": "A"}, {"name": "B"}, {"name": "C"}]
    # Decomposition tree: A→B, A→C (every node in-degree ≤ 1 → 0 convergence).
    decomp = [{"from": "A", "to": "B"}, {"from": "A", "to": "C"}]
    assert compute_convergence_count(concepts, decomp, []) == 0
    # Add a relation B→C: now C has in-degree 2 (parent A + relation B).
    relations = [{"from": "B", "to": "C", "type": "builds-on"}]
    assert compute_convergence_count(concepts, decomp, relations) == 1


def test_compute_convergence_ignores_unknown_and_self():
    from main_outline import compute_convergence_count
    concepts = [{"name": "A"}, {"name": "B"}]
    decomp = [{"from": "A", "to": "B"}]
    relations = [
        {"from": "A", "to": "B"},          # B in-degree → 2 (A twice) → converge
        {"from": "X", "to": "B"},          # unknown source → ignored
        {"from": "B", "to": "B"},          # self → ignored
    ]
    assert compute_convergence_count(concepts, decomp, relations) == 1


def test_ckg_topology_columns_added_by_migration():
    """init_db adds topology columns to a pre-existing old course_ckg; idempotent."""
    import sqlite3 as _sq
    from main_outline import init_db
    conn = _sq.connect(str(DB_PATH))
    conn.execute("DROP TABLE IF EXISTS course_ckg")
    conn.execute(
        "CREATE TABLE course_ckg ("
        "video_id INTEGER PRIMARY KEY, graph_json TEXT, model TEXT, created_at TEXT)"
    )
    conn.commit()
    conn.close()

    init_db(str(DB_PATH))
    init_db(str(DB_PATH))

    conn = _sq.connect(str(DB_PATH))
    cols = {row[1] for row in conn.execute("PRAGMA table_info(course_ckg)")}
    conn.close()
    for c in ("depth", "branch_factor", "convergence_count", "relation_density",
              "density", "avg_path_length", "clustering", "bottomup_ratio"):
        assert c in cols, f"migration did not add column {c}"


def test_ckg_extraction_writes_topology_columns():
    """_run_ckg_extraction computes & persists the decomposition-style params
    (issue 03): depth/branch on the backbone, relation_density, convergence,
    bottomup (aux). density/avg_path_length/clustering are retired → NULL."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 11)
    ]
    # Single decomposition call. Backbone: Supervised Learning → {Classification,
    # Regression}. (Issue 03 reworks the parameter set; here we only assert the
    # worker still computes & persists topology on the decomposition edges.)
    mock_response = {"response": DECOMP_CONCEPTS_RESPONSE}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)

    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate",
                       return_value=mock_response):
                main_outline._run_ckg_extraction(vid)

        conn = main_outline.get_db(str(DB_PATH))
        row = conn.execute(
            "SELECT depth, branch_factor, convergence_count, relation_density, "
            "density, avg_path_length, clustering, bottomup_ratio "
            "FROM course_ckg WHERE video_id = ?", (vid,)
        ).fetchone()
        conn.close()

        assert row is not None
        assert row["depth"] == 1             # longest path: SL → child (1 edge)
        assert row["branch_factor"] == 2.0   # SL is the only non-sink, out-degree 2
        assert row["convergence_count"] == 0  # no relations → no convergence
        assert row["relation_density"] == 0.0  # 0 relations / 3 concepts
        assert row["bottomup_ratio"] == 1.0  # parent taught before both children
        # retired params persisted as NULL
        assert row["density"] is None
        assert row["avg_path_length"] is None
        assert row["clustering"] is None
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_get_ckg_returns_topology_params(client):
    """GET returns the decomposition-style params (issue 03); retired params
    (density/avg_path_length/clustering) are NOT in the response."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    graph_json = json.dumps({
        "schema_version": 2,
        "concepts": [
            {"name": "A", "definition": "a", "first_para": 1, "parent": ""},
            {"name": "B", "definition": "b", "first_para": 4, "parent": "A"},
            {"name": "C", "definition": "c", "first_para": 7, "parent": "B"},
        ],
        "decomposition_edges": [{"from": "A", "to": "B"}, {"from": "B", "to": "C"}],
        "relations": [],
    }, ensure_ascii=False)

    _cleanup_ckg(vid)
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at, "
        "depth, branch_factor, convergence_count, relation_density, "
        "bottomup_ratio) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?)",
        (vid, graph_json, "test-model", "2026-06-25T00:00:00+00:00",
         2, 1.0, 0, 0.0, 1.0),
    )
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get(f"/api/outline/ckg/{vid}")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        for k in ("depth", "branch_factor", "convergence_count",
                  "relation_density", "bottomup_ratio"):
            assert k in data, f"API response missing {k}"
        assert data["depth"] == 2
        assert data["relation_density"] == 0.0
        assert data["bottomup_ratio"] == 1.0
        # retired params must be gone from the contract
        for k in ("density", "avg_path_length", "clustering"):
            assert k not in data, f"retired param {k} still in API response"
    finally:
        _cleanup_ckg(vid)


def test_i18n_has_ckg_topology_param_labels():
    """outline.html defines result-area param i18n keys in all 3 languages."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.result_title", "ckg.param.depth", "ckg.param.branch",
                "ckg.param.relation_density", "ckg.param.convergence",
                "ckg.param.bottomup"):
        assert content.count(f"'{key}'") >= 3, \
            f"i18n key {key} not in all 3 langs"


def _parse_i18n_blocks(content: str):
    """Return {lang: set(keys)} for the three I18N language blocks of a single
    -file HTML app. Blocks open with ``<lang>: {`` and close at the dedented
    ``};``. Keys are quoted ('x') or bare (x) identifiers before a colon."""
    import re
    blocks = {}
    for lang in ("zh", "en", "ja"):
        m = re.search(r"\n\s*" + lang + r":\s*\{", content)
        if not m:
            continue
        start = m.end()
        # find the matching close: first dedented ``},`` or ``};`` line. Inner
        # blocks close with ``},``; only the last (ja) closes with ``};``.
        rest = content[start:]
        endm = re.search(r"\n\s*\}[,;]", rest)
        body = rest[: endm.start()] if endm else rest
        keys = set()
        for line in body.split("\n"):
            km = re.match(r"\s*'([^']+)'\s*:", line) or re.match(
                r"\s*([A-Za-z0-9_.]+)\s*:", line)
            if km:
                keys.add(km.group(1))
        blocks[lang] = keys
    return blocks


def test_outline_i18n_full_parity_across_languages():
    """Every i18n key defined in zh must also exist in en AND ja, so no UI text
    ever silently falls back to Chinese. Guards the whole-project tri-lingual
    requirement against future zh-only additions."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    blocks = _parse_i18n_blocks(html_path.read_text(encoding="utf-8"))
    assert {"zh", "en", "ja"} <= set(blocks)
    miss_en = sorted(blocks["zh"] - blocks["en"])
    miss_ja = sorted(blocks["zh"] - blocks["ja"])
    assert not miss_en, f"keys missing in en: {miss_en}"
    assert not miss_ja, f"keys missing in ja: {miss_ja}"


def test_dbviewer_i18n_full_parity_across_languages():
    """The DB-viewer sub-app must also have full zh/en/ja key parity."""
    html_path = Path(r"d:\Project\All for Style\02-outline\db-viewer-02outline.html")
    blocks = _parse_i18n_blocks(html_path.read_text(encoding="utf-8"))
    assert {"zh", "en", "ja"} <= set(blocks)
    assert not (blocks["zh"] - blocks["en"]), blocks["zh"] - blocks["en"]
    assert not (blocks["zh"] - blocks["ja"]), blocks["zh"] - blocks["ja"]


def test_lessongen_unit_labels_localized_and_live():
    """Unit tabs/headings must localize: the backend bakes no '单元 N' string,
    and the front-end derives the label from unit_index via t(), re-running on
    language switch."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "function _lessonGenPlanLabel" in content
    assert "t('lessongen.unit_n', { n: n })" in content
    assert "function refreshLessonGenUnitsI18n()" in content
    assert "refreshLessonGenUnitsI18n();" in content  # wired into applyI18N
    for key in ("lessongen.unit_n",):
        assert content.count(f"'{key}'") >= 3, f"{key} not in all 3 langs"
    # the "基于 N 节课分析" caption re-translates too
    assert "t('profile.based_on', { n: _profileLectureCount })" in content


def test_ck_sidebar_localized_from_style_key():
    """Workbench CK sidebar localizes its badge/description from style_key and
    re-renders on language switch (was backend Chinese, frozen on switch)."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "function _applyCKProfile" in content
    assert "function refreshCKProfileI18n()" in content
    assert "refreshCKProfileI18n();" in content  # wired into applyI18N
    assert "t('profile.stylelabel.' + data.style_key)" in content
    assert "t('profile.desc.' + data.style_key" in content


def test_ckg_graph_has_edge_legend():
    """The 拆解+关联图 box carries an always-on legend explaining the blue solid
    (decomposition) and yellow dashed (grounded relation) edges, driven by
    data-i18n so it switches zh/en/ja live."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert 'id="ckg-graph-legend"' in content
    assert 'data-i18n="ckg.legend.decomp"' in content
    assert 'data-i18n="ckg.legend.relation"' in content
    # the two swatches use the real edge colors
    assert "border-top:2px solid #a5b4fc" in content   # decomposition edge
    assert "border-top:2px dashed #f59e0b" in content   # grounded relation
    for key in ("ckg.legend.decomp", "ckg.legend.relation"):
        assert content.count(f"'{key}'") >= 3, f"{key} not in all 3 langs"


def test_title_tooltips_are_i18n_driven():
    """Hardcoded CJK title= tooltips were converted to data-i18n-title, and
    applyI18N gained a handler that rewrites the title attribute on lang switch."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # applyI18N handler present.
    assert "data-i18n-title" in content
    assert "setAttribute('title', t(key))" in content
    # Each title key is defined in all 3 languages and wired to an element.
    for key in ("title.more_analysis", "title.drag_divider",
                "title.toggle_console", "title.continue", "title.retry",
                "title.test_error"):
        assert content.count(f"'{key}'") >= 3, f"{key} not in all 3 langs"
        assert f'data-i18n-title="{key}"' in content, f"{key} not bound to element"


def test_dbviewer_row_count_regex_is_language_agnostic():
    """Row deletion must update the header count in any language — the old
    hardcoded /\\d+ 行/ regex broke for en ('N rows'). It now derives the
    suffix from the localized rows_label template."""
    html_path = Path(r"d:\Project\All for Style\02-outline\db-viewer-02outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "/(\\d+) 行/" not in content and "/\\d+ 行/" not in content
    assert "const suffix = t('rows_label', {n: ''});" in content
    assert "new RegExp('(\\\\d+)' + suffix)" in content


def test_frontend_has_ckg_result_area():
    """Middle control panel renders the decomposition-style params; the retired
    density/avgpath/clustering rows are gone."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "ckg-result-area" in content, "missing result-area container"
    assert "renderCKGParams" in content, "missing param render function"
    for el in ("ckg-param-depth", "ckg-param-branch",
               "ckg-param-relation-density", "ckg-param-convergence",
               "ckg-param-bottomup"):
        assert el in content, f"missing result-area element {el}"
    for gone in ("ckg-param-avgpath", "ckg-param-clustering"):
        assert gone not in content, f"retired result-area element {gone} still present"


# ============================================================================
# CKG — corpus-wide scatter view (issue 04)
# ============================================================================

# Three real video ids that exist in corpus.db (FK targets for course_ckg).
CORPUS_TEST_VIDEO_IDS = (132, 133, 134)


def _cleanup_corpus_ckg():
    import main_outline
    conn = main_outline.get_db(str(DB_PATH))
    for vid in CORPUS_TEST_VIDEO_IDS:
        conn.execute("DELETE FROM course_ckg WHERE video_id = ?", (vid,))
    conn.execute("COMMIT")
    conn.close()


def _insert_corpus_row(conn, vid, depth, bottomup_ratio):
    """Insert a minimal course_ckg row with given depth / bottomup_ratio."""
    graph_json = json.dumps(
        {"concepts": [{"name": "A", "definition": "a", "first_para": 1}],
         "edges": []},
        ensure_ascii=False,
    )
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at, "
        "depth, branch_factor, convergence_count, density, avg_path_length, "
        "clustering, bottomup_ratio) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
        (vid, graph_json, "test-model", "2026-06-25T00:00:00+00:00",
         depth, 1.0, 0, 0.1, 1.0, 0.0, bottomup_ratio),
    )


def test_get_corpus_ckg_returns_all_extracted(client):
    """GET /api/outline/ckg returns one entry per extracted video with the
    fields needed by the scatter view (video_id, name, depth, bottomup_ratio)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    _cleanup_corpus_ckg()
    conn = main_outline.get_db(str(DB_PATH))
    _insert_corpus_row(conn, CORPUS_TEST_VIDEO_IDS[0], depth=2, bottomup_ratio=0.8)
    _insert_corpus_row(conn, CORPUS_TEST_VIDEO_IDS[1], depth=5, bottomup_ratio=0.3)
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get("/api/outline/ckg")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        assert isinstance(data, list)
        rows = {r["video_id"]: r for r in data
                if r["video_id"] in CORPUS_TEST_VIDEO_IDS}
        # Only the two inserted videos are present (third was not inserted).
        assert set(rows.keys()) == {CORPUS_TEST_VIDEO_IDS[0],
                                    CORPUS_TEST_VIDEO_IDS[1]}
        for r in rows.values():
            for k in ("video_id", "name", "depth", "bottomup_ratio"):
                assert k in r, f"corpus entry missing {k}"
            assert r["name"], "video name should be non-empty"
        assert rows[CORPUS_TEST_VIDEO_IDS[0]]["depth"] == 2
        assert rows[CORPUS_TEST_VIDEO_IDS[0]]["bottomup_ratio"] == 0.8
        assert rows[CORPUS_TEST_VIDEO_IDS[1]]["depth"] == 5
    finally:
        _cleanup_corpus_ckg()


def test_get_corpus_ckg_excludes_unextracted(client):
    """A video with no course_ckg row must NOT appear in the corpus response."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    _cleanup_corpus_ckg()
    conn = main_outline.get_db(str(DB_PATH))
    _insert_corpus_row(conn, CORPUS_TEST_VIDEO_IDS[0], depth=3, bottomup_ratio=0.5)
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get("/api/outline/ckg")
        assert resp.status_code == 200, resp.text
        ids = {r["video_id"] for r in resp.json()}
        assert CORPUS_TEST_VIDEO_IDS[0] in ids
        # The other two were never inserted -> excluded.
        assert CORPUS_TEST_VIDEO_IDS[1] not in ids
        assert CORPUS_TEST_VIDEO_IDS[2] not in ids
    finally:
        _cleanup_corpus_ckg()


def test_get_corpus_ckg_handles_null_bottomup(client):
    """A row whose bottomup_ratio is NULL (no valid edges) must not crash the
    endpoint; the field is returned as null (front-end skips/marks it)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    _cleanup_corpus_ckg()
    conn = main_outline.get_db(str(DB_PATH))
    _insert_corpus_row(conn, CORPUS_TEST_VIDEO_IDS[0], depth=4, bottomup_ratio=0.6)
    _insert_corpus_row(conn, CORPUS_TEST_VIDEO_IDS[1], depth=1, bottomup_ratio=None)
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get("/api/outline/ckg")
        assert resp.status_code == 200, resp.text
        rows = {r["video_id"]: r for r in resp.json()
                if r["video_id"] in CORPUS_TEST_VIDEO_IDS}
        assert set(rows.keys()) == {CORPUS_TEST_VIDEO_IDS[0],
                                    CORPUS_TEST_VIDEO_IDS[1]}
        # null bottomup is preserved as None, not coerced to 0.
        assert rows[CORPUS_TEST_VIDEO_IDS[1]]["bottomup_ratio"] is None
    finally:
        _cleanup_corpus_ckg()


# NOTE: the corpus-scatter view was removed (single-teacher, can't separate
# content vs style + mixed old/new data). Its removal is asserted by
# test_ckg_scatter_view_removed above. The corpus endpoint /api/outline/ckg
# stays (now feeds the profile-card small-multiples).


# ============================================================================
# CKG prompt editor — read/write API + frontend (issue 05)
# ============================================================================

# ck_prompt_edges retired (ADR 0003) — only the decomposition prompt is editable.
CK_PROMPT_KEYS = ("ck_prompt_concepts",)


def test_get_ckg_prompts_returns_concepts_key(client):
    """GET /api/outline/ckg/prompts returns the editable CK prompt config value."""
    resp = client.get("/api/outline/ckg/prompts")
    assert resp.status_code == 200
    data = resp.json()
    assert isinstance(data, dict)
    for key in CK_PROMPT_KEYS:
        assert key in data, f"Missing CK prompt key: {key}"
        assert isinstance(data[key], str) and data[key].strip(), \
            f"CK prompt {key} should be a non-empty string"
    # The retired edges prompt must NOT be exposed.
    assert "ck_prompt_edges" not in data


def test_put_ckg_prompts_persists_new_value(client):
    """PUT /api/outline/ckg/prompts writes back to config; re-read returns it."""
    # Snapshot original so we can restore.
    original = client.get("/api/outline/ckg/prompts").json()
    try:
        new_concepts = "TEST concept prompt — edited by issue 05"
        resp = client.put(
            "/api/outline/ckg/prompts",
            json={"ck_prompt_concepts": new_concepts},
        )
        assert resp.status_code == 200
        assert resp.json()["ck_prompt_concepts"] == new_concepts

        # Persistence: re-read via GET
        again = client.get("/api/outline/ckg/prompts")
        assert again.json()["ck_prompt_concepts"] == new_concepts

        # And it actually landed in the config table.
        conn = sqlite3.connect(str(DB_PATH))
        row = conn.execute(
            "SELECT value FROM config WHERE key = 'ck_prompt_concepts'"
        ).fetchone()
        conn.close()
        assert row is not None and row[0] == new_concepts
    finally:
        client.put("/api/outline/ckg/prompts", json=original)


def test_put_ckg_prompts_rejects_retired_edges_key(client):
    """PUT including the retired ck_prompt_edges key is rejected (422)."""
    resp = client.put(
        "/api/outline/ckg/prompts",
        json={"ck_prompt_concepts": "ok", "ck_prompt_edges": "retired"},
    )
    assert resp.status_code == 422


def test_put_ckg_prompts_rejects_unknown_key(client):
    """PUT /api/outline/ckg/prompts rejects keys outside the CK prompt set."""
    resp = client.put(
        "/api/outline/ckg/prompts",
        json={"ob_prompt_round1": "should not be allowed here"},
    )
    assert resp.status_code == 422


def test_put_ckg_prompts_empty_body_returns_400(client):
    """PUT /api/outline/ckg/prompts with empty body returns 400."""
    resp = client.put("/api/outline/ckg/prompts", json={})
    assert resp.status_code == 400


def test_ckg_prompts_endpoint_does_not_leak_into_ob_config(client):
    """The ob_ config endpoint must NOT start returning ck_ prompt keys."""
    data = client.get("/api/outline/config").json()
    for key in CK_PROMPT_KEYS:
        assert key not in data, f"{key} must not appear in ob_ config endpoint"


def test_frontend_has_ckg_prompt_editor():
    """Middle panel exposes editors + save button for the two CK prompts and a
    model selector wired to the existing ob_llm_model config."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert 'id="ckg-prompt-concepts"' in content, "missing concepts textarea"
    # ck_prompt_edges retired (ADR 0003): its editor must be gone.
    assert 'id="ckg-prompt-edges"' not in content, "retired edges textarea still present"
    assert 'id="ckg-prompt-save"' in content, "missing save button"
    assert 'id="ckg-llm-model"' in content, "missing model selector"
    # Save handler hits the CK prompts endpoint.
    assert "/api/outline/ckg/prompts" in content, \
        "frontend must call the CK prompts endpoint"
    # Loader / saver functions present.
    assert "loadCKGPrompts" in content, "missing loadCKGPrompts function"
    assert "saveCKGPrompts" in content, "missing saveCKGPrompts function"


def test_i18n_has_ckg_prompt_editor_labels():
    """outline.html defines CK prompt-editor i18n keys across all 3 languages."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.prompt_concepts_label", "ckg.prompt_edges_label",
                "ckg.prompt_save", "ckg.prompt_saved", "ckg.model_label"):
        assert content.count(f"'{key}'") >= 3, \
            f"i18n key {key} not in all 3 langs"


# ============================================================================
# CKG fixes/enhancements: SSE broadcast, delete, edge prompt,
# auto-analysis & delete front-end (issue 05 / SSE rewrite)
# ============================================================================

# ---------------------------------------------------------------------------
# _broadcast fans every event out to every subscriber + the default queue
# ---------------------------------------------------------------------------

def test_broadcast_delivers_to_all_subscribers():
    """_broadcast (and push_log_event) must deliver the SAME event to every
    queue registered in _log_subscribers, not let one connection steal it."""
    import main_outline
    from main_outline import push_log_event, get_log_queue, LogEvent

    main_outline._log_queue = asyncio.Queue()
    q1 = asyncio.Queue()
    q2 = asyncio.Queue()
    main_outline._log_subscribers.clear()
    main_outline._log_subscribers.add(q1)
    main_outline._log_subscribers.add(q2)
    try:
        push_log_event("progress", "broadcast me", progress_pct=42.0)

        # Both per-connection queues received it (broadcast, not stolen).
        e1 = q1.get_nowait()
        e2 = q2.get_nowait()
        assert isinstance(e1, LogEvent) and e1.message == "broadcast me"
        assert isinstance(e2, LogEvent) and e2.message == "broadcast me"
        assert e1.progress_pct == 42.0 and e2.progress_pct == 42.0
        # Default queue still gets it too (keeps old SSE tests passing).
        d = get_log_queue().get_nowait()
        assert d.message == "broadcast me"
    finally:
        main_outline._log_subscribers.clear()


def test_status_polling_endpoint_removed(client):
    """The polling /status endpoint is gone (replaced by SSE broadcast)."""
    resp = client.get(f"/api/outline/ckg/{CKG_TEST_VIDEO_ID}/status")
    assert resp.status_code == 404


def test_ckg_progress_dict_removed():
    """The module-level _ckg_progress polling state is gone."""
    import main_outline
    assert not hasattr(main_outline, "_ckg_progress")


# ---------------------------------------------------------------------------
# DELETE /api/outline/ckg/{video_id} — idempotent delete
# ---------------------------------------------------------------------------

def test_delete_ckg_removes_existing_row(client):
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at) "
        "VALUES (?, ?, ?, ?)",
        (vid, json.dumps({"concepts": [{"name": "A"}], "edges": []}),
         "m", "2026-06-25T00:00:00+00:00"),
    )
    conn.execute("COMMIT")
    conn.close()
    try:
        resp = client.delete(f"/api/outline/ckg/{vid}")
        assert resp.status_code == 200, resp.text
        assert client.get(f"/api/outline/ckg/{vid}").status_code == 404
    finally:
        _cleanup_ckg(vid)


def test_delete_ckg_missing_row_is_idempotent(client):
    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    resp = client.delete(f"/api/outline/ckg/{vid}")
    assert resp.status_code == 200, resp.text


# ---------------------------------------------------------------------------
# Worker pushes SSE progress + success/error events (drives the front-end)
# ---------------------------------------------------------------------------

def test_ckg_worker_pushes_success_event_with_counts():
    """On completion the worker must push a 100% progress 'success' event whose
    message names course_ckg — the front-end keys off this over SSE."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]
    responses = [
        {"response": VALID_CONCEPTS_RESPONSE},
        {"response": VALID_EDGES_RESPONSE},
    ]

    def fake_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
        return responses.pop(0)

    events = []

    def fake_push(category, message, traceback=None, progress_pct=None):
        events.append({"category": category, "message": message,
                       "progress_pct": progress_pct})

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", side_effect=fake_ollama):
                with patch("main_outline.push_log_event", side_effect=fake_push):
                    main_outline._run_ckg_extraction(vid)
        success = [e for e in events if e["category"] == "success"]
        assert success, "worker did not push a success event"
        assert "course_ckg" in success[-1]["message"]
        assert success[-1]["progress_pct"] == 100.0
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


def test_ckg_worker_pushes_error_event_on_parse_failure():
    """A parse failure must push an 'error' SSE event (no progress dict)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID
    mock_paragraphs = [
        {"paragraph_index": 1, "start_time": 0.0, "end_time": 1.0, "text": "x"},
    ]

    def fake_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
        return {"response": "I cannot help with that."}

    events = []

    def fake_push(category, message, traceback=None, progress_pct=None):
        events.append(category)

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", side_effect=fake_ollama):
                with patch("main_outline.push_log_event", side_effect=fake_push):
                    main_outline._run_ckg_extraction(vid)
        assert "error" in events, "worker did not push an error event"
    finally:
        _cleanup_ckg(vid)
        main_outline._ckg_cancel_flags.pop(vid, None)


# ---------------------------------------------------------------------------
# Bug fix: custom ck_prompt_concepts (decomposition) is actually used. The
# worker makes two Ollama calls: decomposition (call 0) then relations (call 1).
# ---------------------------------------------------------------------------

def test_ckg_worker_uses_custom_decompose_prompt():
    """A custom ck_prompt_concepts must be passed to the first (decomposition)
    Ollama call. The second call is the grounded-relations step (issue 02)."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID

    sentinel = "CUSTOM_DECOMPOSE_PROMPT_SENTINEL_XYZ"
    conn = main_outline.get_db(str(DB_PATH))
    conn.execute(
        "INSERT INTO config (key, value) VALUES ('ck_prompt_concepts', ?) "
        "ON CONFLICT(key) DO UPDATE SET value = excluded.value",
        (sentinel,),
    )
    conn.execute("COMMIT")
    conn.close()

    mock_paragraphs = [
        {"paragraph_index": i, "start_time": float(i), "end_time": float(i + 1),
         "text": f"Para {i}"}
        for i in range(1, 6)
    ]
    calls = []

    def fake_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
        calls.append(prompt)
        # call 0 = decomposition, call 1 = relations
        if len(calls) == 1:
            return {"response": DECOMP_CONCEPTS_RESPONSE}
        return {"response": VALID_RELATIONS_RESPONSE}

    main_outline._ckg_cancel_flags.pop(vid, None)
    _cleanup_ckg(vid)
    try:
        with patch("main_outline._fetch_paragraphs_full_text",
                   return_value=mock_paragraphs):
            with patch("main_outline._call_ollama_generate", side_effect=fake_ollama):
                main_outline._run_ckg_extraction(vid)
        assert len(calls) == 2, f"expected 2 Ollama calls, got {len(calls)}"
        assert sentinel in calls[0], "custom decomposition prompt was not used"
    finally:
        _cleanup_ckg(vid)
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM config WHERE key = 'ck_prompt_concepts'")
        conn.execute("COMMIT")
        conn.close()
        init_db(str(DB_PATH))
        main_outline._ckg_cancel_flags.pop(vid, None)


# ---------------------------------------------------------------------------
# Front-end: i18n keys, new buttons, removed dedicated EventSource
# ---------------------------------------------------------------------------

def test_i18n_has_ckg_video_list_title():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert content.count("'ckg.video_list_title'") >= 3, \
        "ckg.video_list_title not in all 3 langs"


def test_i18n_has_ckg_auto_and_delete_keys():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # ckg.auto_btn / ckg.auto_stop replaced by the merged analyze button +
    # auto checkbox. New keys must be present in all three langs.
    for key in ("ckg.auto_label", "ckg.stop_btn", "ckg.delete_btn"):
        assert content.count(f"'{key}'") >= 3, f"i18n key {key} not in all 3 langs"


def test_ck_title_updated_to_mindmap_structure():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "思维导图结构" in content
    assert "Mind Map Structure" in content
    assert "マインドマップ構造" in content


def test_ckg_auto_and_delete_buttons_exist():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # The standalone auto button is replaced by an auto checkbox; the analyze
    # button is merged. A save button now sits left of delete.
    assert 'id="ckg-auto-btn"' not in content
    assert 'id="ckg-auto-check"' in content
    assert 'id="ckg-analyze-btn"' in content
    assert 'id="ckg-save-btn"' in content
    assert 'id="ckg-delete-btn"' in content


def test_ckg_uses_dedicated_eventsource():
    """There are now THREE EventSources on the outline log stream: the main
    console, a dedicated CKG one, and a dedicated lesson-gen one (broadcast
    makes extra connections safe)."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert content.count("new EventSource('/api/stream/logs/outline')") == 3


def test_ckg_frontend_has_no_status_polling():
    """The front-end no longer polls a /status endpoint; CKG extraction is
    SSE-driven. (The lone setInterval is the status-bar clock — not polling.)"""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "/status" not in content
    assert "ckgPollStatus" not in content
    assert "ckgStopPolling" not in content
    assert "_ckgPollTimer" not in content
    # The only setInterval left is the 1s status-bar clock.
    assert content.count("setInterval") == 1


def test_ckg_dag_uses_circle_nodes():
    """The single-course DAG renderer uses small circle nodes (symbol:'circle')
    with hover-only labels, so the global graph structure stays legible
    (labelled rectangles obscured the overview)."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "function renderCKGGraph" in content
    assert "symbol: 'circle'" in content


# ---------------------------------------------------------------------------
# Task 1: ECharts containers must resize after setOption so a freshly-shown
# (clientHeight==0) container renders at its real size instead of 0x0.
# ---------------------------------------------------------------------------

def test_ckg_graph_render_triggers_resize():
    """renderCKGGraph schedules a deferred resize so the chart fills the
    container that was just un-hidden (the restart-blank bug)."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    start = content.index("function renderCKGGraph")
    end = content.index("renderProfileMulti", start)
    body = content[start:end]
    assert "_ckgGraphChart.resize()" in body, \
        "renderCKGGraph must resize the chart after setOption"
    # A deferred trigger (setTimeout or requestAnimationFrame) is used so the
    # resize runs after the browser reflows the just-shown container.
    assert ("setTimeout" in body or "requestAnimationFrame" in body), \
        "renderCKGGraph must defer the resize via setTimeout/rAF"


def test_ckg_single_view_has_tree_above_graph():
    """单课「拆解图」视图 = 上方拆解树 + 下方 first_para 图，互为参照。"""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # tree container above the graph, both inside the dag view wrapper
    assert "ckg-dag-view" in content
    assert "ckg-tree-chart" in content
    assert "function renderCKGTree" in content
    assert "_buildDecompTreeData" in content
    # ECharts tree series for the hierarchy
    assert "type: 'tree'" in content or 'type: "tree"' in content
    # the dag view renders BOTH tree and graph
    assert "renderCKGTree(_ckgData.concepts" in content
    assert "renderCKGGraph(_ckgData.concepts" in content


def test_ckg_scatter_view_removed():
    """全库散点档已移除（按钮 / 容器 / 渲染函数 / i18n 全清）。"""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for gone in ("ckg-view-scatter-btn", "ckg-scatter-chart",
                 "renderCKGScatter", "_ckgScatterChart", "ckg.view_scatter"):
        assert gone not in content, f"残留散点引用: {gone}"


# ---------------------------------------------------------------------------
# Task 2: merged analyze button + auto checkbox
# ---------------------------------------------------------------------------

def test_ckg_auto_checkbox_present():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert 'id="ckg-auto-check"' in content
    assert 'type="checkbox"' in content
    assert 'data-i18n="ckg.auto_label"' in content


def test_i18n_has_new_ckg_button_keys():
    """New / repurposed i18n keys present in all three languages."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    for key in ("ckg.analyze_btn", "ckg.stop_btn", "ckg.auto_label",
                "ckg.save_btn", "ckg.saved"):
        assert content.count(f"'{key}'") >= 3, \
            f"i18n key {key} not in all 3 langs"


# ---------------------------------------------------------------------------
# Task 3: save button + PUT /api/outline/ckg/{video_id}
# ---------------------------------------------------------------------------

def test_ckg_save_button_present_left_of_delete():
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    save_idx = content.index('id="ckg-save-btn"')
    del_idx = content.index('id="ckg-delete-btn"')
    assert save_idx < del_idx, "save button must appear left of delete"


def test_put_ckg_persists_and_recomputes(client):
    """PUT /api/outline/ckg/{id} upserts concepts+edges and server-recomputes
    topology params; GET returns the same concepts/edges with params filled."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    payload = {
        "concepts": [
            {"name": "A", "definition": "a", "first_para": 1},
            {"name": "B", "definition": "b", "first_para": 3},
            {"name": "C", "definition": "c", "first_para": 5},
        ],
        "edges": [{"from": "A", "to": "B"}, {"from": "B", "to": "C"}],
    }
    try:
        resp = client.put(f"/api/outline/ckg/{vid}", json=payload)
        assert resp.status_code == 200, resp.text
        saved = resp.json()
        # Server recomputed the topology: a chain A->B->C has depth 2.
        assert saved["depth"] == 2
        assert "bottomup_ratio" in saved

        got = client.get(f"/api/outline/ckg/{vid}").json()
        names = {c["name"] for c in got["concepts"]}
        assert names == {"A", "B", "C"}
        assert {(e["from"], e["to"]) for e in got["edges"]} == \
            {("A", "B"), ("B", "C")}
        assert got["depth"] == 2
    finally:
        _cleanup_ckg(vid)


def test_put_ckg_empty_concepts_returns_400(client):
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    try:
        resp = client.put(f"/api/outline/ckg/{vid}",
                          json={"concepts": [], "edges": []})
        assert resp.status_code == 400, resp.text
    finally:
        _cleanup_ckg(vid)


def test_put_ckg_breaks_cycles_to_dag(client):
    """Edges containing a cycle are passed through break_cycles before storing,
    so the persisted/returned edges form a DAG."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    vid = CKG_TEST_VIDEO_ID
    _cleanup_ckg(vid)
    payload = {
        "concepts": [
            {"name": "A", "definition": "a", "first_para": 1},
            {"name": "B", "definition": "b", "first_para": 2},
            {"name": "C", "definition": "c", "first_para": 3},
        ],
        # A->B->C->A is a 3-cycle; one edge must be dropped.
        "edges": [{"from": "A", "to": "B"},
                  {"from": "B", "to": "C"},
                  {"from": "C", "to": "A"}],
    }
    try:
        resp = client.put(f"/api/outline/ckg/{vid}", json=payload)
        assert resp.status_code == 200, resp.text
        got = client.get(f"/api/outline/ckg/{vid}").json()
        assert len(got["edges"]) == 2
        assert _is_dag(got["edges"]), "stored edges must form a DAG"
    finally:
        _cleanup_ckg(vid)


# ============================================================================
# Task A — JSON robustness: Ollama format mode + num_predict + salvage parsing
# ============================================================================

def test_call_ollama_generate_passes_format_json():
    """When fmt='json', the request body carries top-level format='json' and
    options.num_predict reflects the num_predict argument."""
    import main_outline

    captured = {}

    class _FakeResp:
        status_code = 200
        def raise_for_status(self):
            pass
        def json(self):
            return {"response": "{}", "done": True, "context": []}

    class _FakeClient:
        def __init__(self, *a, **kw):
            pass
        def __enter__(self):
            return self
        def __exit__(self, *a):
            return False
        def post(self, url, json=None, **kw):
            captured["url"] = url
            captured["json"] = json
            return _FakeResp()

    with patch("main_outline.httpx.Client", _FakeClient):
        main_outline._call_ollama_generate(
            "m", "p", fmt="json", num_predict=4096
        )

    body = captured["json"]
    assert body["format"] == "json"
    assert body["options"]["num_predict"] == 4096


def test_call_ollama_generate_no_format_by_default():
    """Without fmt, no top-level format key is sent; default num_predict=2048."""
    import main_outline

    captured = {}

    class _FakeResp:
        status_code = 200
        def raise_for_status(self):
            pass
        def json(self):
            return {"response": "{}", "done": True, "context": []}

    class _FakeClient:
        def __init__(self, *a, **kw):
            pass
        def __enter__(self):
            return self
        def __exit__(self, *a):
            return False
        def post(self, url, json=None, **kw):
            captured["json"] = json
            return _FakeResp()

    with patch("main_outline.httpx.Client", _FakeClient):
        main_outline._call_ollama_generate("m", "p")

    body = captured["json"]
    assert "format" not in body
    assert body["options"]["num_predict"] == 2048


# --- salvage: broken concept JSON ------------------------------------------

# The real-machine failure: array starts with two valid objects, then the model
# drops the '{' and '"name":' for the next entry, corrupting the whole array.
BROKEN_CONCEPTS_RESPONSE = (
    '{"concepts": ['
    '{"name": "Function", "definition": "A reusable block of code.", "first_para": 1}, '
    '{"name": "Argument", "definition": "A value passed to a function.", "first_para": 2}, '
    '"AND Function", "definition": "logical and", "first_para": 3'
)


def test_parse_concepts_salvages_broken_json():
    """A truncated/corrupt concepts array still yields the leading valid
    concepts via brace-salvage instead of returning None."""
    from main_outline import parse_concepts_response
    result = parse_concepts_response(BROKEN_CONCEPTS_RESPONSE)
    assert result is not None
    names = [c["name"] for c in result["concepts"]]
    assert "Function" in names
    assert "Argument" in names
    assert len(result["concepts"]) == 2


def test_parse_concepts_unsalvageable_returns_none():
    """When nothing flat-object-shaped with a name can be recovered, return None."""
    from main_outline import parse_concepts_response
    assert parse_concepts_response("complete garbage no braces") is None


def test_parse_edges_salvages_broken_json():
    """Corrupt edges array still recovers the leading valid edge objects."""
    from main_outline import parse_edges_response
    broken = (
        '{"edges": ['
        '{"from": "A", "to": "B"}, '
        '{"from": "B", "to": "C"}, '
        '"D", "to": "E"'
    )
    result = parse_edges_response(broken)
    assert result is not None
    assert {"from": "A", "to": "B"} in result["edges"]
    assert {"from": "B", "to": "C"} in result["edges"]
    assert len(result["edges"]) == 2


# ============================================================================
# Task B — corpus endpoint concept_count + frontend list rendering
# ============================================================================

def test_get_corpus_ckg_includes_concept_count(client):
    """GET /api/outline/ckg returns concept_count (and edge_count) per entry,
    computed from the stored graph_json."""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    _cleanup_corpus_ckg()
    conn = main_outline.get_db(str(DB_PATH))
    graph_json = json.dumps({
        "concepts": [
            {"name": "A", "definition": "a", "first_para": 1},
            {"name": "B", "definition": "b", "first_para": 2},
            {"name": "C", "definition": "c", "first_para": 3},
        ],
        "edges": [{"from": "A", "to": "B"}],
    }, ensure_ascii=False)
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at, "
        "depth, branch_factor, convergence_count, density, avg_path_length, "
        "clustering, bottomup_ratio) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
        (CORPUS_TEST_VIDEO_IDS[0], graph_json, "test-model",
         "2026-06-25T00:00:00+00:00", 2, 1.0, 0, 0.1, 1.0, 0.0, 0.5),
    )
    conn.execute("COMMIT")
    conn.close()

    try:
        resp = client.get("/api/outline/ckg")
        assert resp.status_code == 200, resp.text
        rows = {r["video_id"]: r for r in resp.json()
                if r["video_id"] in CORPUS_TEST_VIDEO_IDS}
        entry = rows[CORPUS_TEST_VIDEO_IDS[0]]
        assert entry["concept_count"] == 3
        assert entry["edge_count"] == 1
    finally:
        _cleanup_corpus_ckg()


def test_frontend_ckg_video_list_uses_analysis_status():
    """loadCKGVideos must drive its status dot off the corpus CKG endpoint
    (green=analyzed/red=none) and show a concept-count subtitle."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # References the corpus endpoint to learn which videos are analyzed.
    assert "/api/outline/ckg'" in content or '/api/outline/ckg"' in content
    # Red-dot CSS variant for unanalyzed videos.
    assert ".status-dot.none" in content
    # i18n keys for concept-count + not-analyzed, all three languages.
    for key in ("ckg.concept_count_label", "ckg.not_analyzed"):
        assert content.count(f"'{key}'") >= 3, \
            f"i18n key {key} missing in all 3 langs"


# ============================================================================
# 地基任务 — 导航重构（风格画像卡 / 教案生成）+ DB Viewer 自动加载
# ============================================================================

def _outline_html():
    return Path(r"d:\Project\All for Style\02-outline\outline.html").read_text(
        encoding="utf-8"
    )


def _dbviewer_html():
    return Path(
        r"d:\Project\All for Style\02-outline\db-viewer-02outline.html"
    ).read_text(encoding="utf-8")


# ---- 任务 2/3：导航 tab + 页面壳 ----

def test_nav_has_profile_and_lessongen_tabs():
    """导航存在 data-page=profile 与 data-page=lessongen 两个 tab。"""
    content = _outline_html()
    assert 'data-page="profile"' in content
    assert 'data-page="lessongen"' in content


def test_nav_no_longer_has_dashboard_tab():
    """旧的 data-page=dashboard tab 已移除。"""
    content = _outline_html()
    assert 'data-page="dashboard"' not in content


def test_pages_profile_and_lessongen_exist():
    """页面壳 #page-profile 与 #page-lessongen 存在。"""
    content = _outline_html()
    assert 'id="page-profile"' in content
    assert 'id="page-lessongen"' in content


def test_page_dashboard_removed():
    """旧的 #page-dashboard 页面已移除。"""
    content = _outline_html()
    assert 'id="page-dashboard"' not in content


def test_tab_order_profile_before_lessongen_before_dbviewer():
    """tab 顺序：风格画像卡 在 教案生成 左边，教案生成 在 DB Viewer 左边。"""
    content = _outline_html()
    i_profile = content.index('data-page="profile"')
    i_lessongen = content.index('data-page="lessongen"')
    i_dbviewer = content.index('data-page="dbviewer"')
    assert i_profile < i_lessongen < i_dbviewer


def test_page_switch_logic_handles_profile_and_lessongen():
    """切换逻辑包含 page-profile / 'profile' 与 page-lessongen / 'lessongen'。"""
    content = _outline_html()
    assert "page-profile" in content
    assert "'profile'" in content
    assert "page-lessongen" in content
    assert "'lessongen'" in content
    # 不再切换已删除的 dashboard 页
    assert "page-dashboard" not in content


# ---- i18n 三语 ----

def test_i18n_tab_profile_all_three_langs():
    content = _outline_html()
    assert content.count("'tab_profile'") >= 3, "tab_profile 三语不全"


def test_i18n_tab_lessongen_all_three_langs():
    content = _outline_html()
    assert content.count("'tab_lessongen'") >= 3, "tab_lessongen 三语不全"


def test_i18n_dashboard_keys_removed():
    """旧 i18n 键已移除。"""
    content = _outline_html()
    assert "'tab_dashboard'" not in content
    assert "'dashboard_placeholder'" not in content
    assert "'dashboard_hint'" not in content


# ---- 任务 1：DB Viewer 初始化自动加载 ----

def test_dbviewer_auto_calls_load_default_on_init():
    """db-viewer 在初始化时（DOMContentLoaded 或脚本末尾）自动调用 loadDefault。"""
    content = _dbviewer_html()
    assert (
        "DOMContentLoaded" in content
        and "loadDefault" in content
    ) or content.rstrip().rstrip("</script>").rstrip().endswith("loadDefault();"), (
        "db-viewer 未在初始化时自动调用 loadDefault()"
    )


# ============================================================================
# 风格画像卡 — GET /api/outline/ckg/profile (aggregate over whole course_ckg)
# ============================================================================

import statistics as _stats
from contextlib import contextmanager


@contextmanager
def _isolated_ckg_table():
    """Snapshot the whole course_ckg table, clear it, yield, then restore.

    profile aggregates the *entire* table, so a test must own its full
    contents. We dump every row, DELETE all, hand control to the test, and
    on teardown wipe whatever the test inserted and re-insert the originals.
    """
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))
    conn = main_outline.get_db(str(DB_PATH))
    cols = [r[1] for r in conn.execute("PRAGMA table_info(course_ckg)")]
    saved = [tuple(r) for r in conn.execute(
        f"SELECT {', '.join(cols)} FROM course_ckg"
    ).fetchall()]
    conn.execute("DELETE FROM course_ckg")
    conn.execute("COMMIT")
    conn.close()
    try:
        yield
    finally:
        conn = main_outline.get_db(str(DB_PATH))
        conn.execute("DELETE FROM course_ckg")
        if saved:
            ph = ", ".join(["?"] * len(cols))
            conn.executemany(
                f"INSERT INTO course_ckg ({', '.join(cols)}) VALUES ({ph})",
                saved,
            )
        conn.execute("COMMIT")
        conn.close()


def _insert_ckg_row(conn, *, video_id, concepts_n, depth, branch_factor,
                    convergence_count, density, avg_path_length, clustering,
                    bottomup_ratio):
    graph_json = json.dumps(
        {"concepts": [{"name": f"c{i}"} for i in range(concepts_n)], "edges": []},
        ensure_ascii=False,
    )
    conn.execute(
        "INSERT INTO course_ckg (video_id, graph_json, model, created_at, "
        "depth, branch_factor, convergence_count, density, avg_path_length, "
        "clustering, bottomup_ratio) "
        "VALUES (?, ?, ?, ?, ?, ?, ?, ?, ?, ?, ?)",
        (video_id, graph_json, "test-model", "2026-06-25T00:00:00+00:00",
         depth, branch_factor, convergence_count, density, avg_path_length,
         clustering, bottomup_ratio),
    )


def test_ckg_profile_aggregates_mean_sd(client):
    """profile returns lecture_count + per-param mean/sd; null bottomup_ratio
    rows are skipped for that param only."""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        # Three lectures; the middle one has null bottomup_ratio.
        _insert_ckg_row(conn, video_id=901, concepts_n=10, depth=4,
                        branch_factor=2.0, convergence_count=1, density=0.2,
                        avg_path_length=2.0, clustering=0.1, bottomup_ratio=0.6)
        _insert_ckg_row(conn, video_id=902, concepts_n=20, depth=6,
                        branch_factor=4.0, convergence_count=3, density=0.4,
                        avg_path_length=3.0, clustering=0.3, bottomup_ratio=None)
        _insert_ckg_row(conn, video_id=903, concepts_n=30, depth=8,
                        branch_factor=6.0, convergence_count=5, density=0.6,
                        avg_path_length=4.0, clustering=0.5, bottomup_ratio=0.8)
        conn.execute("COMMIT")
        conn.close()

        resp = client.get("/api/outline/ckg/profile")
        assert resp.status_code == 200, resp.text
        data = resp.json()

        assert data["lecture_count"] == 3
        params = data["params"]

        # depth over all 3 rows: 4,6,8
        assert params["depth"]["mean"] == pytest.approx(6.0)
        assert params["depth"]["sd"] == pytest.approx(_stats.pstdev([4, 6, 8]))
        # node_count over all 3 rows: 10,20,30
        assert params["node_count"]["mean"] == pytest.approx(20.0)
        # bottomup_ratio: null row skipped -> only 0.6, 0.8
        assert params["bottomup_ratio"]["mean"] == pytest.approx(0.7)
        assert params["bottomup_ratio"]["sd"] == pytest.approx(
            _stats.pstdev([0.6, 0.8]))

        for key in ("style_label", "description", "generation_rules"):
            assert key in data and isinstance(data[key], str) and data[key]
        # every required param present (issue 03 param set; density/avg_path/
        # clustering retired, relation_density added)
        for p in ("depth", "branch_factor", "relation_density",
                  "convergence_count", "bottomup_ratio", "node_count"):
            assert p in params
            assert "mean" in params[p] and "sd" in params[p]
        for retired in ("density", "avg_path_length", "clustering"):
            assert retired not in params, f"retired param {retired} still aggregated"


def test_ckg_profile_exposes_style_and_direction_keys(client):
    """Profile returns machine-readable style_key/direction_key so the front-end
    can compose a *localized* description/rules rather than show backend Chinese
    strings. Keys must be from the known vocabularies."""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        # depth<5, branch<3 -> compact; bottomup mean ~0.9 -> bottomup.
        _insert_ckg_row(conn, video_id=921, concepts_n=8, depth=2,
                        branch_factor=1.7, convergence_count=2, density=0.1,
                        avg_path_length=1.5, clustering=0.0, bottomup_ratio=0.9)
        _insert_ckg_row(conn, video_id=922, concepts_n=7, depth=1,
                        branch_factor=1.8, convergence_count=1, density=0.1,
                        avg_path_length=1.0, clustering=0.0, bottomup_ratio=0.95)
        conn.execute("COMMIT")
        conn.close()

        data = client.get("/api/outline/ckg/profile").json()
        assert data["style_key"] in ("deep", "divergent", "balanced", "compact")
        assert data["direction_key"] in ("bottomup", "topdown", "mixed", "unknown")
        assert data["style_key"] == "compact"
        assert data["direction_key"] == "bottomup"


def test_profile_descriptor_rounds_mean_numbers():
    """Mean-based descriptor strings must not leak long floats (e.g.
    1.5044247787610618); depth/node_count are rounded like the other params."""
    import main_outline
    means = {
        "depth": 1.5044247787610618,
        "branch_factor": 1.7345,
        "node_count": 7.79646017699115,
        "avg_path_length": 2.0,
        "bottomup_ratio": 0.93,
    }
    d = main_outline._build_profile_descriptor(means)
    blob = d["style_label"] + d["description"] + d["generation_rules"]
    assert "1.5044247787610618" not in blob
    assert "7.79646017699115" not in blob
    assert d["style_key"] == "compact" and d["direction_key"] == "bottomup"


def test_profile_card_param_tips_and_compose_i18n():
    """Front-end exposes a per-param explanation caption (data-tip-key) for all
    six params in 3 languages, plus client-side text composition that reacts to
    language switches (refreshProfileTextI18n wired into applyI18N)."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    # Six tip keys, each defined in all 3 langs.
    for key in ("profile.tip.depth", "profile.tip.branch",
                "profile.tip.relation_density", "profile.tip.convergence",
                "profile.tip.nodes", "profile.tip.bottomup"):
        assert content.count(f"'{key}'") >= 3, f"tip key {key} not in all 3 langs"
    # Style/direction/description/rules templates defined in all 3 langs.
    for key in ("profile.stylelabel.compact", "profile.dirlabel.bottomup",
                "profile.desc.compact", "profile.dirphrase.bottomup",
                "profile.genrule.depth_shallow", "profile.genrule.branch_few"):
        assert content.count(f"'{key}'") >= 3, f"template key {key} not in all 3 langs"
    # Composition + live re-translation are wired up.
    assert "function _composeProfileText()" in content
    assert "function refreshProfileTextI18n()" in content
    assert "refreshProfileTextI18n();" in content  # called from applyI18N
    assert "data-tip-key=" in content
    # Param NAMES themselves must also re-translate on language switch — they are
    # rendered via a label key and refreshed, not baked into innerHTML once.
    assert "data-param-label-key=" in content
    assert "[data-param-label-key]" in content


def test_ckg_profile_all_null_bottomup(client):
    """When every row's bottomup_ratio is null, that param's mean/sd are null
    but the endpoint still succeeds."""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _insert_ckg_row(conn, video_id=911, concepts_n=5, depth=3,
                        branch_factor=2.0, convergence_count=0, density=0.1,
                        avg_path_length=1.5, clustering=0.0, bottomup_ratio=None)
        conn.execute("COMMIT")
        conn.close()

        resp = client.get("/api/outline/ckg/profile")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        assert data["lecture_count"] == 1
        assert data["params"]["bottomup_ratio"]["mean"] is None
        assert data["params"]["bottomup_ratio"]["sd"] is None
        # other params still computed
        assert data["params"]["depth"]["mean"] == pytest.approx(3.0)


def test_ckg_profile_empty_table(client):
    """Empty corpus -> lecture_count 0, no crash."""
    with _isolated_ckg_table():
        resp = client.get("/api/outline/ckg/profile")
        assert resp.status_code == 200, resp.text
        data = resp.json()
        assert data["lecture_count"] == 0


def test_ckg_profile_route_not_shadowed_by_video_id(client):
    """'profile' must hit the aggregate endpoint, not be parsed as a video_id
    (which would 422 on int parse or 404)."""
    with _isolated_ckg_table():
        resp = client.get("/api/outline/ckg/profile")
        assert resp.status_code == 200, resp.text
        assert "lecture_count" in resp.json()


# ---- 前端：风格画像卡渲染 ----

def test_profile_frontend_has_view_containers():
    """画像卡三视图容器 id 齐全（小多图 / 单课 / 原型树占位）；雷达已退役。"""
    content = _outline_html()
    for el in ("profile-view-multi", "profile-view-single",
               "profile-view-archetype", "profile-multi-grid",
               "profile-single-chart"):
        assert el in content, f"缺少视图容器 {el}"
    # radar retired for the single-teacher stage
    assert "profile-radar" not in content, "退役的 profile-radar 仍存在"


def test_profile_frontend_fetches_profile_endpoint():
    content = _outline_html()
    assert "/api/outline/ckg/profile" in content


def test_profile_frontend_reuses_decomposition_renderer():
    """画像卡视图复用拆解图渲染器（_buildDecompGraphOption），不再用雷达。"""
    content = _outline_html()
    assert "_buildDecompGraphOption" in content
    assert "renderProfileMulti" in content and "renderProfileSingle" in content
    assert "renderProfile" in content


def test_profile_frontend_has_resize_call():
    """视图容器初始 hidden，渲染后须 setTimeout resize。"""
    content = _outline_html()
    assert "_profileSingleChart" in content or "_profileMultiCharts" in content
    assert content.count(".resize()") >= 1


def test_profile_archetype_synthesizes_from_mean_params():
    """原型示意树由均值参数(depth/branch_factor/relation_density)合成，
    用占位概念名，并显示「示意/非真实课程」角标。"""
    content = _outline_html()
    assert "renderProfileArchetype" in content
    # builds the synthetic tree from the stored means
    assert "_profileMeans.depth" in content
    assert "_profileMeans.branch_factor" in content
    assert "_profileMeans.relation_density" in content
    # renders via the shared decomposition renderer
    assert "_buildDecompGraphOption" in content
    # honest schematic badge (not pretending to be a real lecture)
    assert "profile.archetype_badge" in content
    for lang_val in ("非任何真实课程", "not any real lecture", "実在の講義ではない"):
        assert lang_val in content, f"missing archetype badge text: {lang_val}"


def test_profile_archetype_has_tree_above_graph():
    """原型树视图也是上树下图:示意拆解树 + 示意拆解图,复用树渲染器。"""
    content = _outline_html()
    assert "profile-archetype-tree-chart" in content, "缺少原型树容器"
    assert "_profileArchetypeTreeChart" in content
    # reuses the shared tree option builder
    assert "_buildDecompTreeOption" in content
    # tree built from the synthetic concepts (root excluded) + edges
    assert "_buildDecompTreeOption(concepts.slice(1), edges, rootName)" in content


def test_profile_frontend_switch_page_triggers_load():
    """切到 profile tab 时调用 profile 渲染函数。"""
    content = _outline_html()
    assert ("page === 'profile'" in content)
    assert ("renderProfile" in content or "loadProfile" in content)


def test_profile_i18n_keys_all_three_langs():
    """新增 profile.* i18n 键三语齐全。"""
    content = _outline_html()
    for key in ("profile.title", "profile.based_on", "profile.params_title",
                "profile.rules_title", "profile.empty",
                "profile.view.multi", "profile.view.single",
                "profile.view.archetype", "profile.archetype_badge"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_profile_i18n_radar_dims_all_three_langs():
    """雷达维度名 i18n 键三语齐全。"""
    content = _outline_html()
    for key in ("profile.dim.depth", "profile.dim.branch",
                "profile.dim.bottomup", "profile.dim.convergence",
                "profile.dim.density"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


# ===================================================================
# 教案生成 (Lesson Generator) — POST /api/outline/lesson-gen
# ===================================================================

# 概念 / 边 / 风格化生成 三段 LLM 输出（合法 JSON）。
_LG_CONCEPTS = json.dumps({
    "title": "Decision Tree",
    "concepts": [
        {"name": "Entropy", "definition": "Measure of impurity.",
         "first_para": 0, "parent": "Information Gain"},
        {"name": "Information Gain", "definition": "Reduction in entropy.",
         "first_para": 0, "parent": "Decision Tree"},
        {"name": "Decision Tree", "definition": "A tree-based classifier.",
         "first_para": 0, "parent": "Decision Tree"},
    ]
}, ensure_ascii=False)

# Grounded relations (issue 06): lesson-gen now extracts 讲述关联, not prerequisites.
_LG_RELATIONS = json.dumps({
    "relations": [
        {"from": "Entropy", "to": "Decision Tree", "type": "motivates"},
    ]
}, ensure_ascii=False)

_LG_OUTLINE = json.dumps({
    "outline": [
        {"name": "Decision Tree", "children": [
            {"name": "Information Gain", "children": [
                {"name": "Entropy", "children": []}
            ]}
        ]}
    ],
    "sequence": ["Entropy", "Information Gain", "Decision Tree"],
}, ensure_ascii=False)


def _seed_profile_row(conn):
    """Insert one course_ckg row so the corpus profile is non-empty."""
    _insert_ckg_row(conn, video_id=950, concepts_n=12, depth=5,
                    branch_factor=3.0, convergence_count=2, density=0.3,
                    avg_path_length=2.5, clustering=0.2, bottomup_ratio=0.7)
    conn.execute("COMMIT")


@contextmanager
def _inert_thread():
    """Patch threading.Thread so POST's daemon worker doesn't auto-run.

    The async lesson-gen endpoint starts a daemon thread; in tests we want to
    drive ``_run_lesson_gen`` ourselves deterministically (so mock_ollama
    responses are consumed exactly once). This replaces Thread with a stub
    whose ``start()`` is a no-op.
    """
    class _NoopThread:
        def __init__(self, *a, **k):
            pass

        def start(self):
            pass

    with patch("main_outline.threading.Thread", _NoopThread):
        yield


def test_lesson_gen_material_mode_returns_job_and_done_result(client):
    """material 模式（异步）：POST 返回 202 + job_id；跑完 worker 后
    GET result 返回 done + outline/sequence/concepts/edges/style_label。
    三次 LLM（概念→边→生成）。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        responses = [
            {"response": _LG_CONCEPTS},
            {"response": _LG_RELATIONS},
            {"response": _LG_OUTLINE},
        ]
        calls = []

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            calls.append(prompt)
            return responses.pop(0)

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "some lecture material",
                                     "mode": "material"})
            assert resp.status_code == 202, resp.text
            job_id = resp.json()["job_id"]
            assert job_id

            # Run the worker synchronously (no real thread needed) then poll.
            main_outline._run_lesson_gen(job_id, "some lecture material", "material")

        assert len(calls) == 3  # material → decomposition, relations, generate
        r = client.get(f"/api/outline/lesson-gen/result/{job_id}")
        assert r.status_code == 200, r.text
        data = r.json()
        assert data["state"] == "done"
        result = data["result"]
        assert isinstance(result["outline"], list) and result["outline"]
        assert result["sequence"] == ["Entropy", "Information Gain", "Decision Tree"]
        assert len(result["concepts"]) == 3
        assert result["title"] == "Decision Tree"   # tree root for the output
        # decomposition backbone derived from parents (DT→IG, IG→Entropy)
        assert len(result["decomposition_edges"]) == 2
        assert result["edges"] == result["decomposition_edges"]  # alias
        # grounded relation validated against the concept set
        assert len(result["relations"]) == 1
        assert result["relations"][0]["type"] == "motivates"
        assert result.get("style_label")


def test_lesson_gen_topic_mode_returns_done_result(client):
    """topic 模式（异步）：两次 LLM（主题枚举→生成）→ done result 正常。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        # topic mode now returns a decomposition (concepts with parent) + relations
        enum_resp = json.dumps({
            "title": "决策树",
            "concepts": [
                {"name": "Entropy", "definition": "Impurity.", "parent": "Decision Tree"},
                {"name": "Decision Tree", "definition": "Classifier.", "parent": "决策树"},
            ],
            "relations": [{"from": "Entropy", "to": "Decision Tree", "type": "motivates"}],
        }, ensure_ascii=False)
        gen_resp = json.dumps({
            "outline": [{"name": "Decision Tree", "children": [
                {"name": "Entropy", "children": []}]}],
            "sequence": ["Entropy", "Decision Tree"],
        }, ensure_ascii=False)

        responses = [{"response": enum_resp}, {"response": gen_resp}]
        calls = []

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            calls.append(prompt)
            return responses.pop(0)

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "决策树", "mode": "topic"})
            assert resp.status_code == 202, resp.text
            job_id = resp.json()["job_id"]
            main_outline._run_lesson_gen(job_id, "决策树", "topic")

        assert len(calls) == 2  # topic → enumerate + generate
        r = client.get(f"/api/outline/lesson-gen/result/{job_id}")
        assert r.status_code == 200, r.text
        data = r.json()
        assert data["state"] == "done"
        result = data["result"]
        assert result["sequence"] == ["Entropy", "Decision Tree"]
        assert len(result["concepts"]) == 2
        # Decision Tree → Entropy from parent; relation validated.
        assert len(result["decomposition_edges"]) == 1
        assert len(result["relations"]) == 1
        assert result.get("style_label")


def test_lesson_gen_topic_prompt_is_decomposition():
    """The topic prompt asks for a decomposition (parent) + relations, not
    prerequisite edges."""
    import main_outline
    p = main_outline._LESSON_GEN_TOPIC_PROMPT
    assert "DECOMPOSITION" in p.upper()
    assert '"parent"' in p and '"relations"' in p
    assert "PREREQUISITE" not in p.upper()


def test_lesson_gen_outline_prompt_uses_decomposition_and_relations():
    """The outline prompt feeds decomposition + grounded relations (not prereqs)."""
    import main_outline
    p = main_outline._LESSON_GEN_OUTLINE_PROMPT
    assert "DECOMPOSITION" in p.upper()
    assert "GROUNDED RELATIONS" in p.upper()
    assert "{relation_lines}" in p
    assert "PREREQUISITE" not in p.upper()


def test_lessongen_renders_grounded_relations():
    """The lesson-gen graph renderer threads relations through as a 4th arg and
    draws them as a distinct link set."""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "renderLessonGenGraph(concepts, edges, sequence, relations)" in content
    # the per-plan renderer threads the plan's relations + decomposition edges
    assert "renderLessonGenGraph(plan.concepts, dEdges, plan.sequence, plan.relations)" in content
    assert "plan.decomposition_edges || plan.edges" in content


def test_lessongen_material_extracted_per_page(client):
    """多页资料(以 _PAGE_DELIM 分隔)逐页抽概念:每页一次 LLM 调用,概念合并去重。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        delim = main_outline._PAGE_DELIM
        # 3 pages, each a distinct concept page.
        page1 = '{"title": "T", "concepts": [{"name": "P1A", "definition": "a", "parent": "T"}]}'
        page2 = '{"title": "T", "concepts": [{"name": "P2A", "definition": "b", "parent": "T"}]}'
        page3 = '{"title": "T", "concepts": [{"name": "P3A", "definition": "c", "parent": "T"}]}'
        material = ("PAGE-ONE-TEXT" + delim + "PAGE-TWO-TEXT" + delim + "PAGE-THREE-TEXT")

        prompts = []
        page_responses = [page1, page2, page3]

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            prompts.append(prompt)
            # First 3 calls = per-page concept extraction; then relations; then gen.
            if len(prompts) <= 3:
                return {"response": page_responses[len(prompts) - 1]}
            if len(prompts) == 4:
                return {"response": '{"relations": []}'}
            return {"response": _LG_OUTLINE}

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            job_id = "perpage-job"
            main_outline._lesson_gen_jobs[job_id] = {"state": "running"}
            main_outline._run_lesson_gen(job_id, material, "material")

        # One concept call PER PAGE (3) + relations (1) + generation (1) = 5.
        assert len(prompts) == 5, f"expected 5 calls, got {len(prompts)}"
        # Each per-page prompt carried only its own page text (no whole-doc dump).
        assert "PAGE-ONE-TEXT" in prompts[0] and "PAGE-TWO-TEXT" not in prompts[0]
        assert "PAGE-TWO-TEXT" in prompts[1]
        # Merged concepts cover all three pages.
        r = client.get(f"/api/outline/lesson-gen/result/{job_id}").json()["result"]
        names = {c["name"] for c in r["concepts"]}
        assert {"P1A", "P2A", "P3A"} <= names


def test_lessongen_calls_use_config_num_ctx():
    """material/CKG LLM 调用传入 config 配置的 num_ctx,而非硬编码常量。"""
    import main_outline, inspect
    # the LLM calls live in the extract + build-plan helpers now
    src = (inspect.getsource(main_outline._lessongen_extract_concepts)
           + inspect.getsource(main_outline._lessongen_build_plan))
    assert "num_ctx=num_ctx" in src
    # loader returns (model, temperature, num_ctx)
    assert main_outline._load_ob_llm_config()  # smoke
    assert "_LESSON_GEN_NUM_CTX" not in dir(main_outline), "stale constant remains"


def test_ob_llm_num_ctx_seeded_and_parsed():
    """ob_llm_num_ctx 默认被 seed;_parse_num_ctx 容错并钳制范围。"""
    from main_outline import init_db, _parse_num_ctx
    init_db(str(DB_PATH))
    conn = sqlite3.connect(str(DB_PATH))
    row = conn.execute(
        "SELECT value FROM config WHERE key = 'ob_llm_num_ctx'"
    ).fetchone()
    conn.close()
    assert row is not None, "ob_llm_num_ctx not seeded"
    assert _parse_num_ctx("16384") == 16384
    assert _parse_num_ctx(None) == 8192          # bad → default
    assert _parse_num_ctx("99") == 1024          # clamped to min
    assert _parse_num_ctx("999999999") == 131072  # clamped to max


def test_load_ob_llm_config_returns_num_ctx():
    """_load_ob_llm_config 现在返回 (model, temperature, num_ctx) 三元组。"""
    from main_outline import init_db, _load_ob_llm_config
    init_db(str(DB_PATH))
    model, temp, num_ctx = _load_ob_llm_config()
    assert isinstance(model, str) and isinstance(temp, float)
    assert isinstance(num_ctx, int) and num_ctx >= 1024


def test_config_numctx_frontend_control():
    """设置面板含 num_ctx 输入,并接入 config 读写。"""
    content = _outline_html()
    assert 'id="ob_llm_num_ctx"' in content
    assert "ob_llm_num_ctx: document.getElementById('ob_llm_num_ctx').value" in content
    for key in ("config.numctx_label", "config.numctx_hint"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_call_ollama_generate_accepts_num_ctx():
    """_call_ollama_generate 暴露 num_ctx 参数,只在提供时写入 options。"""
    import main_outline, inspect
    sig = inspect.signature(main_outline._call_ollama_generate)
    assert "num_ctx" in sig.parameters


def test_lessongen_stepper_has_parse_node():
    """流水线含「解析资料」前置节点,且由前端 _lgParseDone 驱动。"""
    content = _outline_html()
    assert "lessongen.step.parse" in content
    assert "_lgParseDone" in content
    # parse node marked done right after resolving the staged file
    assert "_lgParseDone = true" in content


def test_chunk_material_page_vs_segment():
    """逐页 vs 分段:page=按 _PAGE_DELIM 分页;segment=按字数打包。"""
    import main_outline
    delim = main_outline._PAGE_DELIM
    text = "PAGE1" + delim + "PAGE2" + delim + "PAGE3"
    assert main_outline._chunk_material(text, "page") == ["PAGE1", "PAGE2", "PAGE3"]
    # short text → one segment
    assert len(main_outline._chunk_material(text, "segment")) == 1
    # long text → multiple bounded segments
    long_text = "\n".join("line %d content here" % i for i in range(4000))
    segs = main_outline._chunk_material(long_text, "segment")
    assert len(segs) >= 2
    assert all(len(s) <= main_outline._LESSON_GEN_SEGMENT_CHARS + 50 for s in segs)


def test_lesson_gen_rejects_bad_chunk_mode(client):
    """chunk_mode 非法 → 400(在画像检查之前就拦下)。"""
    resp = client.post("/api/outline/lesson-gen",
                       json={"input_text": "x", "mode": "material",
                             "chunk_mode": "bogus"})
    assert resp.status_code == 400, resp.text
    assert "chunk_mode" in resp.json().get("detail", "")


# ---------------------------------------------------------------------------
# Big-lecture → small-units split (two-phase analyze + generate-units)
# ---------------------------------------------------------------------------

def test_partition_concepts_in_order_equal_slices():
    """按 first_para 顺序连续等分;前 n%k 个单元各多 1 个。"""
    from main_outline import _partition_concepts_in_order
    concepts = [{"name": f"C{i}", "first_para": i} for i in range(10)]
    parts = _partition_concepts_in_order(concepts, 3)
    assert [len(p) for p in parts] == [4, 3, 3]   # 10 = 4+3+3
    # contiguous & in order
    assert [c["name"] for c in parts[0]] == ["C0", "C1", "C2", "C3"]
    assert [c["name"] for c in parts[2]] == ["C7", "C8", "C9"]
    # k clamped to concept count
    assert len(_partition_concepts_in_order(concepts, 99)) == 10
    assert len(_partition_concepts_in_order(concepts, 1)) == 1


def test_suggest_unit_count():
    from main_outline import _suggest_unit_count, _CONCEPTS_PER_UNIT
    assert _suggest_unit_count(0) == 1
    assert _suggest_unit_count(8) == 1
    assert _suggest_unit_count(40) == round(40 / _CONCEPTS_PER_UNIT)


def test_lesson_gen_analyze_reports_split(client):
    """analyze 逐页抽概念,返回 concept_count / needs_split / suggested_units。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        delim = main_outline._PAGE_DELIM
        # Build 3 pages whose concepts total > threshold (18) so needs_split=True.
        def page(prefix, n):
            cs = ", ".join(
                '{"name": "%s%d", "definition": "d", "parent": "T"}' % (prefix, i)
                for i in range(n))
            return '{"title": "T", "concepts": [' + cs + ']}'
        material = page("A", 8) + delim + page("B", 8) + delim + page("C", 8)

        responses = [page("A", 8), page("B", 8), page("C", 8)]

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            return {"response": responses.pop(0)}

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen/analyze",
                               json={"input_text": material, "chunk_mode": "page"})
            assert resp.status_code == 202, resp.text
            job_id = resp.json()["job_id"]
            main_outline._run_lesson_gen_analyze(job_id, material, "page")

        res = client.get(f"/api/outline/lesson-gen/result/{job_id}").json()["result"]
        assert res["phase"] == "analyze"
        assert res["concept_count"] == 24
        assert res["needs_split"] is True
        assert res["suggested_units"] == round(24 / main_outline._CONCEPTS_PER_UNIT)
        assert len(res["concepts"]) == 24


def test_lesson_gen_generate_units_makes_k_plans(client):
    """generate-units 把概念切成 K 份,每份产出一份教案。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        concepts = [{"name": f"C{i}", "definition": "d", "parent": "T",
                     "first_para": i} for i in range(12)]

        # Each unit: relations call + generation call.
        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            if "relations" in prompt.lower() or "RELATION" in prompt:
                return {"response": '{"relations": []}'}
            return {"response": _LG_OUTLINE}

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen/generate-units",
                               json={"concepts": concepts, "title": "T", "units": 3})
            assert resp.status_code == 202, resp.text
            job_id = resp.json()["job_id"]
            main_outline._run_lesson_gen_units(job_id, concepts, "T", 3)

        res = client.get(f"/api/outline/lesson-gen/result/{job_id}").json()["result"]
        assert res["phase"] == "units"
        assert res["unit_count"] == 3
        assert len(res["units"]) == 3
        # every unit is a full plan + their concepts cover all 12 (disjoint)
        all_names = []
        for n, u in enumerate(res["units"], start=1):
            assert "outline" in u and "decomposition_edges" in u and "relations" in u
            # No language-specific label baked in — the front-end localizes the
            # tab/heading from unit_index (was "单元 N", now "" + unit_index).
            assert u["title"] == ""
            assert u["unit_index"] == n
            all_names += [c["name"] for c in u["concepts"]]
        assert sorted(all_names) == sorted(c["name"] for c in concepts)


def test_lesson_gen_generate_units_rejects_empty(client):
    resp = client.post("/api/outline/lesson-gen/generate-units",
                       json={"concepts": [], "units": 2})
    assert resp.status_code == 400, resp.text


def test_lesson_plans_table_created():
    """init_db creates the lesson_plans table."""
    from main_outline import init_db
    init_db(str(DB_PATH))
    conn = sqlite3.connect(str(DB_PATH))
    tables = {r[0] for r in conn.execute("SELECT name FROM sqlite_master WHERE type='table'")}
    conn.close()
    assert "lesson_plans" in tables


def test_lesson_plans_crud(client):
    """保存 → 列表 → 取详情 → 重命名 → 删除 全链路。"""
    import main_outline
    from main_outline import init_db
    init_db(str(DB_PATH))

    payload = {"title": "T", "units": [{"title": "单元 1", "outline": [],
               "sequence": ["A"], "concepts": [{"name": "A"}],
               "decomposition_edges": [], "relations": []}]}
    # save
    resp = client.post("/api/outline/lesson-plans",
                       json={"name": "我的教案", "payload": payload})
    assert resp.status_code == 200, resp.text
    pid = resp.json()["id"]

    # list includes it
    lst = client.get("/api/outline/lesson-plans").json()
    assert any(it["id"] == pid and it["name"] == "我的教案" for it in lst)

    # full payload round-trips
    full = client.get(f"/api/outline/lesson-plans/{pid}").json()
    assert full["name"] == "我的教案"
    assert full["payload"]["units"][0]["concepts"][0]["name"] == "A"

    # rename
    r = client.put(f"/api/outline/lesson-plans/{pid}", json={"name": "改名了"})
    assert r.status_code == 200
    assert client.get(f"/api/outline/lesson-plans/{pid}").json()["name"] == "改名了"

    # delete → gone
    client.delete(f"/api/outline/lesson-plans/{pid}")
    assert client.get(f"/api/outline/lesson-plans/{pid}").status_code == 404


def test_lesson_plans_validation(client):
    from main_outline import init_db
    init_db(str(DB_PATH))
    # empty name → 400
    assert client.post("/api/outline/lesson-plans",
                       json={"name": "", "payload": {}}).status_code == 400
    # missing payload → 400
    assert client.post("/api/outline/lesson-plans",
                       json={"name": "x"}).status_code == 400
    # rename/get/missing → 404
    assert client.get("/api/outline/lesson-plans/999999").status_code == 404
    assert client.put("/api/outline/lesson-plans/999999",
                      json={"name": "y"}).status_code == 404


def test_lessongen_save_list_divider_frontend():
    """前端含保存按钮、教案列表、可拖动分隔条。"""
    content = _outline_html()
    assert 'id="lessongen-save-btn"' in content and "function saveLessonPlan" in content
    assert 'id="lessongen-list-body"' in content and "function loadLessonPlanList" in content
    assert "function loadLessonPlan(" in content and "function renameLessonPlan(" in content and "function deleteLessonPlan(" in content
    assert 'id="lessongen-divider"' in content and "_initLessonGenDivider" in content
    assert "/api/outline/lesson-plans" in content
    for key in ("lessongen.save_btn", "lessongen.list_title", "lessongen.delete_confirm"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_console_and_input_are_collapsible():
    """底部监控窗 + 教案输入区都可折叠(默认收起、点击展开)。"""
    content = _outline_html()
    # console: collapsed by default + toggle + latest-message + chevron
    assert "console-collapsed" in content
    assert "function toggleConsole()" in content
    assert 'id="console-latest"' in content and 'id="console-chevron"' in content
    assert 'onclick="toggleConsole()"' in content
    # lesson-gen input: collapsible header + body (hidden by default) + toggle
    assert 'id="lessongen-input-header"' in content
    assert 'id="lessongen-input-body" class="hidden' in content
    assert "function toggleLessonInput()" in content
    assert 'onclick="toggleLessonInput()"' in content


def test_lessongen_unit_split_frontend():
    """前端含两阶段流程:material 走 analyze、拆分询问 + 单元标签页。"""
    content = _outline_html()
    # material → analyze endpoint; split-ask + unit tabs containers
    assert "/api/outline/lesson-gen/analyze" in content
    assert "/api/outline/lesson-gen/generate-units" in content
    assert 'id="lessongen-split-ask"' in content
    assert 'id="lessongen-unit-tabs"' in content
    # phase branching + unit rendering
    assert "R.phase === 'analyze'" in content
    assert "R.phase === 'units'" in content
    assert "_lessonGenRenderUnits" in content
    assert "_lessonGenStartUnits" in content
    for key in ("lessongen.split_msg", "lessongen.split_confirm",
                "lessongen.split_waiting"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_lessongen_chunk_mode_selector_frontend():
    """前端含逐页/分段选择器,并在 POST 里带上 chunk_mode。"""
    content = _outline_html()
    assert 'data-lg-chunk="page"' in content
    assert 'data-lg-chunk="segment"' in content
    assert "_lessonGenChunkMode" in content
    assert "chunk_mode: _lessonGenChunkMode" in content
    for key in ("lessongen.chunk_label", "lessongen.chunk_page",
                "lessongen.chunk_segment"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_graph_symbol_size_is_adaptive():
    """蓝点大小随概念数自适应缩小(避免多节点重叠)。"""
    content = _outline_html()
    # shared decomposition graph builder + the lesson-gen graph both scale size
    assert "14 - n / 12" in content, "decomp graph not adaptive"
    assert "14 - nCount / 12" in content, "lesson-gen graph not adaptive"


def test_lessongen_output_has_tree_above_graph():
    """教案生成输出右半也是上树下图:拆解树 + 拆解+关联图,复用树渲染器。"""
    html_path = Path(r"d:\Project\All for Style\02-outline\outline.html")
    content = html_path.read_text(encoding="utf-8")
    assert "lessongen-tree-chart" in content, "缺少教案树容器"
    assert "function renderLessonGenTree" in content
    assert "_lessonGenTreeChart" in content
    # reuses the shared tree option builder
    assert "_buildDecompTreeOption(concepts || [], edges || [], title)" in content
    # the per-plan renderer draws BOTH tree and graph; the tree root uses the
    # localized unit label (plan.title, else "单元 N / Unit N / ユニット N")
    assert "renderLessonGenTree(plan.concepts, dEdges, _lessonGenPlanLabel(plan, idx))" in content


def test_lesson_gen_no_profile_returns_error_and_starts_no_job(client):
    """全库无已分析课程 → POST 400，不调 LLM、不起任务。"""
    import main_outline
    with _isolated_ckg_table():
        jobs_before = dict(main_outline._lesson_gen_jobs)

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            raise AssertionError("LLM must not be called when no profile exists")

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama):
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "anything", "mode": "material"})
        assert resp.status_code == 400, resp.text
        body = resp.json()
        msg = json.dumps(body, ensure_ascii=False)
        assert "风格" in msg or "profile" in msg.lower() or "分析" in msg
        # 没有新任务被登记。
        assert dict(main_outline._lesson_gen_jobs) == jobs_before


def test_lesson_gen_bad_generation_json_degrades(client):
    """风格化生成输出坏 JSON：worker 不崩，result 仍为 done 且 outline 平铺。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        responses = [
            {"response": _LG_CONCEPTS},
            {"response": _LG_RELATIONS},
            {"response": "this is not json at all <<<"},
        ]

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            return responses.pop(0)

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "material text",
                                     "mode": "material"})
            assert resp.status_code == 202, resp.text
            job_id = resp.json()["job_id"]
            main_outline._run_lesson_gen(job_id, "material text", "material")

        r = client.get(f"/api/outline/lesson-gen/result/{job_id}")
        assert r.status_code == 200, r.text
        data = r.json()
        # 降级而非失败：state done，outline 仍由 concepts/sequence 平铺。
        assert data["state"] == "done"
        result = data["result"]
        assert isinstance(result["outline"], list) and result["outline"]
        assert len(result["sequence"]) == 3


def test_lesson_gen_result_unknown_job_returns_404(client):
    """未知 job_id → 404。"""
    resp = client.get("/api/outline/lesson-gen/result/nonexistent-job")
    assert resp.status_code == 404, resp.text


def test_lesson_gen_job_state_transitions_running_to_done(client):
    """worker 跑完后 _lesson_gen_jobs[job_id] 由 running → done。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        responses = [
            {"response": _LG_CONCEPTS},
            {"response": _LG_RELATIONS},
            {"response": _LG_OUTLINE},
        ]

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            return responses.pop(0)

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "m", "mode": "material"})
            job_id = resp.json()["job_id"]
            # POST 后、worker 跑之前应为 running。
            assert main_outline._lesson_gen_jobs[job_id]["state"] == "running"
            main_outline._run_lesson_gen(job_id, "m", "material")
        assert main_outline._lesson_gen_jobs[job_id]["state"] == "done"


def test_lesson_gen_job_state_error_on_concept_failure(client):
    """概念抽取彻底失败（LLM 抛错）→ worker 置 error，result 含 detail。"""
    import main_outline
    with _isolated_ckg_table():
        conn = main_outline.get_db(str(DB_PATH))
        _seed_profile_row(conn)
        conn.close()

        def mock_ollama(model, prompt, stream=False, temperature=0.0, **kwargs):
            raise RuntimeError("ollama down")

        with patch("main_outline._call_ollama_generate", side_effect=mock_ollama), \
                _inert_thread():
            resp = client.post("/api/outline/lesson-gen",
                               json={"input_text": "m", "mode": "material"})
            job_id = resp.json()["job_id"]
            main_outline._run_lesson_gen(job_id, "m", "material")

        r = client.get(f"/api/outline/lesson-gen/result/{job_id}")
        assert r.status_code == 200, r.text
        data = r.json()
        assert data["state"] == "error"
        assert data.get("detail")


# ---- 前端：教案生成页 ----

def test_lessongen_page_has_mode_switch_and_inputs():
    """#page-lessongen 含模式切换 + 文本框 + 生成按钮。"""
    content = _outline_html()
    # 模式切换控件（data-lg-mode 分段按钮）。
    assert 'data-lg-mode="material"' in content
    assert 'data-lg-mode="topic"' in content
    # 文本输入框 + 生成按钮 id。
    assert 'id="lessongen-input"' in content
    assert 'id="lessongen-generate-btn"' in content


def test_lessongen_calls_endpoint():
    """前端引用 POST /api/outline/lesson-gen。"""
    content = _outline_html()
    assert "/api/outline/lesson-gen" in content


def test_lessongen_renders_outline_and_dag():
    """输出区含层级大纲容器 + ECharts graph DAG + resize 修复。"""
    content = _outline_html()
    assert 'id="lessongen-outline"' in content
    assert 'id="lessongen-dag-chart"' in content
    # ECharts graph 渲染函数。
    assert "renderLessonGenGraph" in content or "lessongen-dag-chart" in content
    assert "type: 'graph'" in content or 'type: "graph"' in content
    # 容器初始 hidden → 渲染后 resize。
    assert "_lessonGenGraphChart" in content


def test_lessongen_switch_page_triggers_init():
    """切到 lessongen tab 时绑定/初始化教案生成逻辑。"""
    content = _outline_html()
    assert "page === 'lessongen'" in content


def test_lessongen_i18n_keys_all_three_langs():
    """新增 lessongen.* i18n 键三语齐全。"""
    content = _outline_html()
    for key in ("lessongen.mode_material", "lessongen.mode_topic",
                "lessongen.generate", "lessongen.placeholder_material",
                "lessongen.placeholder_topic", "lessongen.outline_title",
                "lessongen.dag_title", "lessongen.generating",
                "lessongen.error"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


# ---- 前端：教案生成流水线 stepper + SSE 异步 ----

def test_lessongen_has_stepper_nodes():
    """#page-lessongen 含流水线 stepper 节点容器。"""
    content = _outline_html()
    assert 'id="lessongen-stepper"' in content
    # 六个节点（含「解析资料」前置步骤，用 data-lg-step 索引标识）。
    for i in range(6):
        assert f'data-lg-step="{i}"' in content, f"缺少 stepper 节点 {i}"


def test_lessongen_uses_result_endpoint():
    """前端引用 GET /api/outline/lesson-gen/result/ 取异步结果。"""
    content = _outline_html()
    assert "/api/outline/lesson-gen/result/" in content


def test_lessongen_uses_eventsource_not_polling():
    """异步流程用专用 EventSource（SSE）驱动，不轮询 /result。"""
    content = _outline_html()
    assert "EventSource" in content
    # 不得用 setInterval 轮询 result 端点。
    import re
    # setInterval 块里不应出现 lesson-gen/result（粗粒度断言）。
    assert "setInterval" not in content or \
        not re.search(r"setInterval[\s\S]{0,400}lesson-gen/result", content), \
        "不得用 setInterval 轮询 /result"


def test_lessongen_filters_by_progress_prefix():
    """专用 SSE 仅处理以「教案生成 · 」开头的事件，避免被其他 SSE 干扰。"""
    content = _outline_html()
    assert "教案生成 · " in content


def test_lessongen_stepper_i18n_keys_all_three_langs():
    """stepper 各节点名 + 状态文案 i18n 三语齐全。"""
    content = _outline_html()
    for key in ("lessongen.step.parse", "lessongen.step.profile",
                "lessongen.step.concepts", "lessongen.step.edges",
                "lessongen.step.generate", "lessongen.step.done"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


# ===================================================================
# 教案生成 — 文件上传解析 (extract-file)
# POST /api/outline/lesson-gen/extract-file
# ===================================================================


def test_extract_text_txt_and_md_decode_utf8():
    """.txt / .md：bytes 按 UTF-8 解码为文本。"""
    import main_outline
    data = "你好 hello\n第二行".encode("utf-8")
    assert main_outline._extract_text_from_upload("a.txt", data).strip() == \
        "你好 hello\n第二行"
    assert main_outline._extract_text_from_upload("b.md", data).strip() == \
        "你好 hello\n第二行"


def test_extract_text_txt_bad_bytes_does_not_crash():
    """非 UTF-8 bytes：回退解码，不抛异常。"""
    import main_outline
    data = b"\xff\xfe plain ascii tail"
    out = main_outline._extract_text_from_upload("x.txt", data)
    assert "plain ascii tail" in out


def test_extract_text_docx_reads_paragraphs():
    """.docx：用 python-docx 现造文档，断言抽出已知段落。"""
    import io
    import docx
    import main_outline
    doc = docx.Document()
    doc.add_paragraph("DOCX 标题段落")
    doc.add_paragraph("第二个段落 SECOND")
    buf = io.BytesIO()
    doc.save(buf)
    out = main_outline._extract_text_from_upload("lesson.docx", buf.getvalue())
    assert "DOCX 标题段落" in out
    assert "第二个段落 SECOND" in out


def test_extract_text_pptx_reads_textframes():
    """.pptx：用 python-pptx 现造演示文稿，断言抽出文本框内容。"""
    import io
    from pptx import Presentation
    from pptx.util import Inches
    import main_outline
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    tb.text_frame.text = "PPTX 幻灯片文本 SLIDETEXT"
    buf = io.BytesIO()
    prs.save(buf)
    out = main_outline._extract_text_from_upload("deck.pptx", buf.getvalue())
    assert "PPTX 幻灯片文本 SLIDETEXT" in out


def test_extract_text_pdf_dispatches_to_pdf_branch():
    """.pdf：分派到 pdf 分支（monkeypatch pypdf 验证）。"""
    import main_outline
    fake_pages = [MagicMock(), MagicMock()]
    fake_pages[0].extract_text.return_value = "PDF page one"
    fake_pages[1].extract_text.return_value = "PDF page two"
    fake_reader = MagicMock()
    fake_reader.pages = fake_pages

    import pypdf
    with patch.object(pypdf, "PdfReader", return_value=fake_reader):
        out = main_outline._extract_text_from_upload("doc.pdf", b"%PDF-fake")
    assert "PDF page one" in out
    assert "PDF page two" in out


def test_extract_text_unsupported_extension_raises():
    """不支持的扩展名 → ValueError（端点转 400）。"""
    import main_outline
    with pytest.raises(ValueError):
        main_outline._extract_text_from_upload("image.png", b"\x89PNG")


def test_extract_file_endpoint_txt_returns_payload(client):
    """端点：上传 .txt → 返回 {filename, text, char_count}。"""
    content = "教案资料文本 sample material"
    resp = client.post(
        "/api/outline/lesson-gen/extract-file",
        files={"file": ("note.txt", content.encode("utf-8"), "text/plain")},
    )
    assert resp.status_code == 200, resp.text
    data = resp.json()
    assert data["filename"] == "note.txt"
    assert "教案资料文本" in data["text"]
    assert data["char_count"] == len(data["text"])


def test_extract_file_endpoint_unsupported_returns_400(client):
    """端点：不支持扩展名 → 400，detail 提示支持的类型。"""
    resp = client.post(
        "/api/outline/lesson-gen/extract-file",
        files={"file": ("pic.png", b"\x89PNG\r\n", "image/png")},
    )
    assert resp.status_code == 400, resp.text
    detail = json.dumps(resp.json(), ensure_ascii=False)
    assert "pdf" in detail.lower()


def test_extract_file_endpoint_bad_docx_does_not_500(client):
    """端点：坏文件（声称 .docx 实为乱码）→ 400/422，不 500。"""
    resp = client.post(
        "/api/outline/lesson-gen/extract-file",
        files={"file": ("broken.docx", b"not a real docx zip",
                        "application/octet-stream")},
    )
    assert resp.status_code in (400, 422), resp.text


# ---- 前端：文件上传控件 ----

def test_lessongen_has_file_upload_control():
    """#page-lessongen 含 file input（accept 含各类型）+ 引用 extract-file 端点。"""
    content = _outline_html()
    assert 'type="file"' in content
    # accept 覆盖所有支持类型
    for ext in (".pdf", ".pptx", ".docx", ".txt", ".md"):
        assert ext in content
    assert "/api/outline/lesson-gen/extract-file" in content


def test_lessongen_upload_i18n_keys_all_three_langs():
    """上传相关 i18n 键三语齐全。"""
    content = _outline_html()
    for key in ("lessongen.upload_label", "lessongen.uploading",
                "lessongen.loaded", "lessongen.upload_failed",
                "lessongen.unsupported", "lessongen.staged"):
        assert content.count(f"'{key}'") >= 3, f"{key} 三语不全"


def test_lessongen_upload_is_lazy_parsed_at_generate():
    """上传只暂存文件、不回填粘贴框;在「生成」时才解析。"""
    content = _outline_html()
    # staged-file state + lazy resolver
    assert "_lessonGenFile" in content
    assert "_lessonGenResolveInput" in content
    # upload handler must NOT push extracted text into the textarea
    assert "input.value = data.text" not in content, \
        "upload 仍在回填粘贴框(应改为生成时才解析)"
    # generation resolves (parses) the staged file before POSTing
    assert "await _lessonGenResolveInput()" in content
